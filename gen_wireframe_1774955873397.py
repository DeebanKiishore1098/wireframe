from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import argparse, re

# ----------------------------
# Required utilities
# ----------------------------
def clean_txt(t):
    return re.sub(r'<[^>]*>', '', str(t))[:200]

def _hex_to_rgb(hex_color):
    hc = str(hex_color).strip().lstrip("#")
    if len(hc) == 3:
        hc = "".join([c * 2 for c in hc])
    if len(hc) != 6:
        hc = "000000"
    return RGBColor(int(hc[0:2], 16), int(hc[2:4], 16), int(hc[4:6], 16))

def set_bg(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)

def _within_bounds(x, y, w, h):
    # SAFE AREA: x=[0.3 .. 13.0], y=[0.3 .. 7.2]
    return (x >= 0.3 and y >= 0.3 and (x + w) <= 13.0 and (y + h) <= 7.2)

def add_rect(slide, x, y, w, h, fill_hex, line_hex="#334155", line_w=1.0, radius=0.08):
    if not _within_bounds(x, y, w, h):
        return None
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    shp.line.color.rgb = _hex_to_rgb(line_hex)
    shp.line.width = Pt(line_w)
    try:
        shp.adjustments[0] = float(radius)
    except Exception:
        pass
    return shp

def add_txt(slide, text, x, y, w, h, size_pt, color_hex, bold=False, align=1):
    if not _within_bounds(x, y, w, h):
        return None
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean_txt(text)
    run.font.name = "Calibri"
    run.font.size = Pt(size_pt)
    run.font.bold = bool(bold)
    run.font.color.rgb = _hex_to_rgb(color_hex)
    if align == 2:
        p.alignment = PP_ALIGN.CENTER
    elif align == 3:
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    tf.word_wrap = True
    return tb

def _text_h(lines, size_pt):
    return max(0.2, float(lines) * (float(size_pt) / 72.0) * 1.4)

def _truncate(s, max_chars):
    s = clean_txt(s)
    if len(s) <= max_chars:
        return s
    return s[: max(0, max_chars - 1)] + "…"

def add_input_field(slide, label, x, y, w):
    label = _truncate(label, 40)
    label_h = _text_h(1, 11)
    gap = 0.06
    input_h = 0.42
    total_h = label_h + gap + input_h
    if not _within_bounds(x, y, w, total_h):
        return 0.0
    add_txt(slide, label, x, y, w, label_h, 11, "#94A3B8", False, 1)
    add_rect(slide, x, y + label_h + gap, w, input_h, "#0F172A", "#475569", 1.0, 0.08)
    add_txt(slide, "Enter " + _truncate(label.lower(), 24), x + 0.12, y + label_h + gap + 0.10, w - 0.24, _text_h(1, 11), 11, "#94A3B8", False, 1)
    return total_h

def add_button(slide, text, x, y, w, h, bg_hex="#3B82F6", fg_hex="#FFFFFF"):
    if not _within_bounds(x, y, w, h):
        return None
    add_rect(slide, x, y, w, h, bg_hex, bg_hex, 1.0, 0.08)
    add_txt(slide, _truncate(text, 24), x, y + (h - _text_h(1, 12)) / 2.0, w, _text_h(1, 12), 12, fg_hex, True, 2)
    return True

def add_table_placeholder(slide, headers, x, y, w, h):
    if not _within_bounds(x, y, w, h):
        return None
    add_rect(slide, x, y, w, h, "#1E293B", "#334155", 1.0, 0.08)

    headers = [ _truncate(hd, 20) for hd in (headers or []) ]
    ncol = max(1, min(6, len(headers)))
    col_w = w / float(ncol)

    header_h = 0.45
    row_h = 0.45
    rows = 5  # header + 4 data rows
    needed_h = header_h + 4 * row_h + 0.55  # plus pagination area
    if needed_h > h:
        # shrink rows if needed
        row_h = max(0.32, (h - header_h - 0.55) / 4.0)

    # header background
    add_rect(slide, x + 0.02, y + 0.02, w - 0.04, header_h, "#0F172A", "#334155", 1.0, 0.08)
    for i in range(ncol):
        tx = x + i * col_w
        add_txt(slide, headers[i] if i < len(headers) else "Column", tx + 0.12, y + 0.12, col_w - 0.24, _text_h(1, 11), 11, "#F1F5F9", True, 1)

    # data rows
    sample_rows = [
        ["A-1029", "Active", "High", "2026-02-14", "Owner 1", "Notes…"],
        ["A-1030", "Draft", "Medium", "2026-02-18", "Owner 2", "Notes…"],
        ["A-1031", "Active", "Low", "2026-03-02", "Owner 3", "Notes…"],
        ["A-1032", "Paused", "High", "2026-03-10", "Owner 4", "Notes…"],
    ]
    for r in range(4):
        ry = y + 0.02 + header_h + r * row_h
        tint = "#1E293B" if (r % 2 == 0) else "#172033"
        add_rect(slide, x + 0.02, ry, w - 0.04, row_h, tint, "#334155", 0.8, 0.08)
        for c in range(ncol):
            cell = sample_rows[r][c] if c < len(sample_rows[r]) else "—"
            add_txt(slide, _truncate(cell, 30), x + c * col_w + 0.12, ry + 0.12, col_w - 0.24, _text_h(1, 11), 11, "#94A3B8", False, 1)

    # pagination
    pag_y = y + 0.02 + header_h + 4 * row_h + 0.12
    if _within_bounds(x + 0.02, pag_y, w - 0.04, 0.35):
        add_txt(slide, "Showing 1–4 of 128", x + 0.12, pag_y, 3.0, _text_h(1, 11), 11, "#94A3B8", False, 1)
        # page buttons right
        btn_w = 0.55
        btn_h = 0.32
        right_x = x + w - 0.02 - (btn_w * 4 + 0.12 * 3)
        labels = ["Prev", "1", "2", "Next"]
        for i, lab in enumerate(labels):
            bx = right_x + i * (btn_w + 0.12)
            bg = "#0F172A" if lab != "1" else "#3B82F6"
            fg = "#94A3B8" if lab != "1" else "#FFFFFF"
            add_button(slide, lab, bx, pag_y + 0.02, btn_w, btn_h, bg, fg)

    return True

# ----------------------------
# Layout helpers
# ----------------------------
NAV_H = 0.6
SIDEBAR_W = 2.0
CARD_GAP = 0.2
ROW_GAP = 0.15
SECTION_GAP = 0.3

BG = "#0F172A"
SURFACE = "#1E293B"
BORDER = "#334155"
TEXT_P = "#F1F5F9"
TEXT_S = "#94A3B8"
ACCENT = "#3B82F6"

def add_top_nav(slide, app_name, accent_hex):
    x0, y0, w0, h0 = 0.3, 0.3, 12.7, NAV_H
    add_rect(slide, x0, y0, w0, h0, SURFACE, BORDER, 1.0, 0.08)

    # Logo area
    add_txt(slide, _truncate(app_name, 20), x0 + 0.2, y0 + 0.16, 2.2, _text_h(1, 14), 14, TEXT_P, True, 1)

    # Nav links center
    links = ["Dashboard", "Data", "Reports", "Settings"]
    cx = x0 + 3.0
    link_w = 1.3
    for i, lab in enumerate(links):
        lab = _truncate(lab, 20)
        lx = cx + i * (link_w + 0.15)
        if lx + link_w > (x0 + w0 - 3.0):
            break
        col = TEXT_P if lab == "Dashboard" else TEXT_S
        add_txt(slide, lab, lx, y0 + 0.18, link_w, _text_h(1, 12), 12, col, lab == "Dashboard", 2)
        if lab == "Dashboard":
            add_rect(slide, lx + 0.15, y0 + h0 - 0.10, link_w - 0.30, 0.06, accent_hex, accent_hex, 0.0, 0.08)

    # Right icons
    # Bell
    add_rect(slide, x0 + w0 - 1.55, y0 + 0.14, 0.36, 0.36, "#0F172A", BORDER, 1.0, 0.08)
    add_txt(slide, "🔔", x0 + w0 - 1.55, y0 + 0.16, 0.36, _text_h(1, 12), 12, TEXT_S, False, 2)
    # Avatar
    add_rect(slide, x0 + w0 - 1.10, y0 + 0.12, 0.42, 0.42, "#0F172A", BORDER, 1.0, 0.21)
    add_txt(slide, "U", x0 + w0 - 1.10, y0 + 0.18, 0.42, _text_h(1, 12), 12, TEXT_P, True, 2)
    # User name
    add_txt(slide, "Admin", x0 + w0 - 0.62, y0 + 0.18, 0.55, _text_h(1, 11), 11, TEXT_S, False, 3)

def add_sidebar(slide, accent_hex, active="Dashboard"):
    x, y, w, h = 0.3, 0.3 + NAV_H + 0.15, SIDEBAR_W, 7.2 - (0.3 + NAV_H + 0.15)
    add_rect(slide, x, y, w, h, SURFACE, BORDER, 1.0, 0.08)

    items = [
        ("🏠", "Dashboard"),
        ("📄", "Records"),
        ("📊", "Analytics"),
        ("⚙", "Settings"),
        ("❓", "Help"),
    ]
    y_cursor = y + 0.2
    for icon, label in items:
        row_h = 0.5
        if not _within_bounds(x + 0.12, y_cursor, w - 0.24, row_h):
            break
        is_active = (label == active)
        bg = "#0F172A" if not is_active else accent_hex
        fg = TEXT_S if not is_active else "#FFFFFF"
        add_rect(slide, x + 0.12, y_cursor, w - 0.24, row_h, bg, BORDER, 1.0, 0.08)
        add_txt(slide, icon, x + 0.22, y_cursor + 0.14, 0.35, _text_h(1, 12), 12, fg, False, 1)
        add_txt(slide, _truncate(label, 20), x + 0.55, y_cursor + 0.14, w - 0.85, _text_h(1, 12), 12, fg, is_active, 1)
        y_cursor += row_h + 0.12

def add_page_header(slide, title, breadcrumb, x, y, w, accent_hex, actions=None):
    title = _truncate(title, 60)
    breadcrumb = _truncate(breadcrumb, 80)
    title_h = _text_h(1, 20)
    crumb_h = _text_h(1, 11)
    block_h = title_h + 0.06 + crumb_h + 0.10
    if not _within_bounds(x, y, w, block_h):
        return 0.0

    add_txt(slide, title, x, y, w * 0.65, title_h, 20, TEXT_P, True, 1)
    add_txt(slide, breadcrumb, x, y + title_h + 0.06, w * 0.65, crumb_h, 11, TEXT_S, False, 1)

    # action buttons right
    actions = actions or ["+ New", "Export", "Filter"]
    btn_h = 0.38
    btn_w = 1.05
    gap = 0.12
    total_w = len(actions) * btn_w + (len(actions) - 1) * gap
    bx = x + w - total_w
    by = y + 0.06
    for i, act in enumerate(actions):
        bg = accent_hex if i == 0 else "#0F172A"
        fg = "#FFFFFF" if i == 0 else TEXT_S
        add_button(slide, _truncate(act, 20), bx + i * (btn_w + gap), by, btn_w, btn_h, bg, fg)

    return block_h

def add_kpi_row(slide, kpis, x, y, w, accent_hex):
    n = max(3, min(4, len(kpis)))
    card_h = 1.05
    card_w = (w - (n - 1) * CARD_GAP) / float(n)
    if not _within_bounds(x, y, w, card_h):
        return 0.0
    for i in range(n):
        cx = x + i * (card_w + CARD_GAP)
        add_rect(slide, cx, y, card_w, card_h, SURFACE, BORDER, 1.0, 0.08)
        label, value, trend = kpis[i]
        add_txt(slide, _truncate(label, 24), cx + 0.18, y + 0.18, card_w - 0.36, _text_h(1, 11), 11, TEXT_S, False, 1)
        add_txt(slide, _truncate(value, 18), cx + 0.18, y + 0.44, card_w - 0.36, _text_h(1, 20), 20, TEXT_P, True, 1)
        add_txt(slide, _truncate(trend, 18), cx + 0.18, y + 0.78, card_w - 0.36, _text_h(1, 11), 11, accent_hex, True, 1)
    return card_h

def add_charts_grid(slide, x, y, w, h, labels):
    col_gap = 0.2
    col_w = (w - col_gap) / 2.0
    row_h = (h - ROW_GAP) / 2.0
    # place up to 2 charts (one row) or 4 (2x2) depending on available height
    placed = 0
    for r in range(2):
        for c in range(2):
            if placed >= len(labels):
                return placed
            cx = x + c * (col_w + col_gap)
            cy = y + r * (row_h + ROW_GAP)
            if not _within_bounds(cx, cy, col_w, row_h):
                continue
            add_rect(slide, cx, cy, col_w, row_h, SURFACE, BORDER, 1.0, 0.08)
            add_txt(slide, _truncate(labels[placed], 60), cx + 0.18, cy + 0.18, col_w - 0.36, _text_h(2, 12), 12, TEXT_S, False, 1)
            # inner plot area
            inner_y = cy + 0.55
            inner_h = max(0.4, row_h - 0.75)
            if _within_bounds(cx + 0.18, inner_y, col_w - 0.36, inner_h):
                add_rect(slide, cx + 0.18, inner_y, col_w - 0.36, inner_h, "#0F172A", "#475569", 1.0, 0.08)
                add_txt(slide, "Chart area", cx + 0.18, inner_y + inner_h/2 - _text_h(1, 11)/2, col_w - 0.36, _text_h(1, 11), 11, TEXT_S, False, 2)
            placed += 1
    return placed

def add_filter_bar(slide, x, y, w, accent_hex):
    bar_h = 1.05
    if not _within_bounds(x, y, w, bar_h):
        return 0.0
    add_rect(slide, x, y, w, bar_h, SURFACE, BORDER, 1.0, 0.08)
    # two inputs + search button
    pad = 0.18
    field_w = (w - pad*2 - 1.2 - 0.24*2) / 2.0
    fx1 = x + pad
    fy = y + 0.18
    used1 = add_input_field(slide, "Keyword", fx1, fy, field_w)
    fx2 = fx1 + field_w + 0.24
    used2 = add_input_field(slide, "Status", fx2, fy, field_w)
    bx = x + w - pad - 1.2
    by = y + 0.52
    add_button(slide, "Search", bx, by, 1.2, 0.42, accent_hex, "#FFFFFF")
    return bar_h

# ----------------------------
# Data: first 5 applications
# ----------------------------
BRIEF = {
  "portfolio": {
    "title": "Systech Analytics Application Portfolio",
    "organisation": "Systech Analytics",
    "year": "2026",
    "confidential": True,
    "stats": {
      "totalApps": "15 Applications",
      "domains": "AI, data engineering, automation, and analytics",
      "aiPowered": "7+ AI / LLM-Powered",
      "buildTenure": "Not specified"
    }
  },
  "applications": [
    {
      "id": "01",
      "name": "Maverick",
      "domain": "Casino",
      "category": "AI Generation",
      "tagline": "AI-Powered Promotional Coupon Engine for Casino Players",
      "overview": "Maverick is a personalised incentive generation platform for the gaming and casino industry. It analyses player behaviour, spending history, game preferences, and visit frequency to craft hyper-personalised promotional offers — from free spins and dining credits to VIP room upgrades. Powered by AI, it replaces static promotions with dynamic, data-driven campaigns that adapt in real time to maximise engagement and lifetime value.",
      "features": [
        {"title": "Player segmentation","description": "Segments players by lifetime value, churn risk, and play patterns."},
        {"title": "Automated coupon generation","description": "Generates coupons with configurable reward rules and expiry logic."},
        {"title": "Real-time offer delivery","description": "Delivers offers via SMS, app notifications, and kiosk displays."},
        {"title": "Campaign analytics","description": "Tracks performance with analytics and A/B testing support."},
        {"title": "Systems integration","description": "Integrates with casino POS and loyalty management systems."}
      ],
      "impact": "Increases promotional ROI and player retention by delivering the right offer to the right player at the right moment, replacing guesswork with precision targeting.",
      "metrics": [{"label": "Coverage","value": "Full"},{"label": "Status","value": "Live"},{"label": "Version","value": "v1.0"}],
      "accent": "#F59E0B",
      "icon": "🎰"
    },
    {
      "id": "02",
      "name": "TRBank",
      "domain": "Banking",
      "category": "Data Processing",
      "tagline": "Unstructured-to-Structured Data Transformation for Banking",
      "overview": "TRBank is an intelligent document processing platform for financial institutions. It ingests unstructured banking documents — statements, loan applications, KYC forms, remittance slips, and correspondence — and transforms them into clean, validated, structured datasets ready for analytics, compliance reporting, and core banking system ingestion.",
      "features": [
        {"title": "Multi-format ingestion","description": "Supports PDFs, scanned images, and handwritten forms."},
        {"title": "Entity extraction","description": "Extracts account numbers, dates, transaction amounts, and counterparties."},
        {"title": "Validation rules engine","description": "Flags anomalies, duplicates, and missing fields."},
        {"title": "Structured outputs","description": "Exports JSON, CSV, or writes directly to databases."},
        {"title": "Auditability","description": "Provides audit trail and confidence scoring per extracted field."}
      ],
      "impact": "Eliminates manual data entry, reducing processing time from hours to minutes per batch while improving accuracy and compliance readiness.",
      "metrics": [{"label": "Coverage","value": "Full"},{"label": "Status","value": "Live"},{"label": "Version","value": "v1.0"}],
      "accent": "#2563EB",
      "icon": "🏦"
    },
    {
      "id": "03",
      "name": "Sustainability",
      "domain": "ESG",
      "category": "Emission Analytics",
      "tagline": "Full-Cycle ESG Intelligence and Emission Calculation Platform",
      "overview": "The Sustainability platform is an end-to-end ESG intelligence tool for organisations managing their environmental footprint. It ingests utility bills, fuel logs, travel records, and supply chain data to calculate carbon emissions across Scope 1, 2, and 3. It generates audit-ready ESG reports aligned with GHG Protocol and GRI standards and provides predictive analytics to model emission reduction scenarios for net-zero planning.",
      "features": [
        {"title": "Automated bill parsing","description": "Parses utility bills and maps emission factors intelligently."},
        {"title": "Scope 1–3 accounting","description": "Calculates emissions aligned with the GHG Protocol."},
        {"title": "AI ESG narratives","description": "Generates narrative ESG reports for board and regulatory submission."},
        {"title": "Benchmarking dashboards","description": "Provides peer benchmarking and target-gap analysis."},
        {"title": "Scenario modelling","description": "Models net-zero pathways and reduction roadmaps."}
      ],
      "impact": "Reduces ESG reporting effort from weeks of manual spreadsheets to automated, audit-ready output generated in minutes.",
      "metrics": [{"label": "Coverage","value": "Full"},{"label": "Status","value": "Live"},{"label": "Version","value": "v1.0"}],
      "accent": "#16A34A",
      "icon": "🌿"
    },
    {
      "id": "04",
      "name": "Ecovision",
      "domain": "ESG",
      "category": "Vision AI",
      "tagline": "Video-Powered AI Auditing for Sustainability Compliance",
      "overview": "Ecovision transforms passive CCTV into an active sustainability compliance monitoring system. It analyses live and recorded facility video feeds to detect energy waste, improper waste disposal, safety non-compliance, and resource misuse. Each event is logged with a video clip, timestamp, severity rating, and an AI-generated corrective action recommendation.",
      "features": [
        {"title": "Real-time video analysis","description": "Detects sustainability and compliance events across facilities."},
        {"title": "Incident logging","description": "Extracts clips and classifies severity automatically."},
        {"title": "Corrective recommendations","description": "Generates AI-driven corrective actions per incident."},
        {"title": "Compliance scoring","description": "Provides dashboards with historical trend analytics."},
        {"title": "Platform integration","description": "Integrates natively with the Sustainability platform for unified reporting."}
      ],
      "impact": "Turns periodic manual walkthroughs into continuous, automated, evidence-driven compliance with an always-on audit trail.",
      "metrics": [{"label": "Coverage","value": "Full"},{"label": "Status","value": "Live"},{"label": "Version","value": "v1.0"}],
      "accent": "#16A34A",
      "icon": "🌿"
    },
    {
      "id": "05",
      "name": "Hotel Concierge",
      "domain": "Hospitality",
      "category": "AI Agent",
      "tagline": "AI Avatar Booking Agent for Hotels and Guest Amenities",
      "overview": "Hotel Concierge is a conversational AI agent that handles end-to-end hotel booking and guest services through a lifelike AI avatar interface. Guests can search rooms, check availability, book, and reserve amenities (dining, spa, leisure). The agent maintains multi-turn context and orchestrates reservations via PMS integrations for a seamless, always-available self-service experience.",
      "features": [
        {"title": "Booking lifecycle","description": "Supports booking, modification, and cancellation via natural language."},
        {"title": "Amenity reservations","description": "Books dining, spa, pool, gym, and activities."},
        {"title": "Conversational memory","description": "Maintains multi-turn context and tracks guest preferences."},
        {"title": "AI avatar interface","description": "Provides a lifelike avatar for immersive guest interaction."},
        {"title": "PMS integration","description": "Connects to PMS and channel managers for real-time availability and pricing."}
      ],
      "impact": "Reduces front-desk workload and call volumes while increasing ancillary revenue via 24/7 self-service booking.",
      "metrics": [{"label": "Coverage","value": "Full"},{"label": "Status","value": "Live"},{"label": "Version","value": "v1.0"}],
      "accent": "#0EA5E9",
      "icon": "🏨"
    }
  ]
}

def _domain_item(app):
    name = clean_txt(app.get("name", "App"))
    domain = clean_txt(app.get("domain", "Data"))
    # simple mapping
    if "Casino" in domain:
        return "Offers"
    if "Bank" in domain or "Banking" in domain:
        return "Documents"
    if "ESG" in domain:
        return "Reports"
    if "Hospital" in domain:
        return "Bookings"
    return "Records"

def _table_headers(app):
    domain = clean_txt(app.get("domain", ""))
    if "Casino" in domain:
        return ["Offer ID", "Segment", "Reward", "Channel", "Status", "Expires"]
    if "Bank" in domain:
        return ["Doc ID", "Type", "Source", "Confidence", "Status", "Received"]
    if "ESG" in domain:
        return ["Report ID", "Scope", "Period", "Emissions", "Status", "Owner"]
    if "Hospital" in domain:
        return ["Booking ID", "Guest", "Room", "Dates", "Status", "Channel"]
    return ["ID", "Type", "Owner", "Status", "Updated", "Notes"]

def _kpis(app):
    item = _domain_item(app)
    return [
        (item + " processed", "1,284", "+12% WoW"),
        ("Active items", "312", "+3% WoW"),
        ("Alerts", "18", "-5% WoW"),
        ("SLA met", "96%", "+1.2% WoW"),
    ]

def _chart_labels(app):
    domain = clean_txt(app.get("domain", ""))
    if "Casino" in domain:
        return ["Bar Chart — Redemptions by Day", "Line Chart — ROI by Campaign", "Table — Top Segments", "Area Chart — Visits vs Offers"]
    if "Bank" in domain:
        return ["Bar Chart — Docs by Type", "Line Chart — Extraction Accuracy", "Heatmap — Exceptions by Rule", "Area Chart — Throughput per Hour"]
    if "ESG" in domain:
        return ["Line Chart — Emissions by Month", "Bar Chart — Scope 1/2/3 Breakdown", "Table — Top Sources", "Area Chart — Reduction Scenario"]
    if "Hospital" in domain:
        return ["Line Chart — Bookings by Day", "Bar Chart — Revenue by Room Type", "Table — Top Amenities", "Area Chart — Conversion Funnel"]
    return ["Bar Chart — Volume", "Line Chart — Trend", "Table — Top Items", "Area Chart — Forecast"]

def _form_fields(app):
    item = _domain_item(app)
    domain = clean_txt(app.get("domain", ""))
    if "Casino" in domain:
        return ["Campaign name", "Player segment", "Reward type", "Reward value", "Delivery channel", "Expiry date"]
    if "Bank" in domain:
        return ["Batch name", "Document type", "Ingestion source", "Validation profile", "Output format", "Notification email"]
    if "ESG" in domain:
        return ["Report name", "Reporting period", "Scope selection", "Emission factors set", "Approval workflow", "Owner"]
    if "Hospital" in domain:
        return ["Guest name", "Room type", "Check-in date", "Check-out date", "Amenity package", "Special requests"]
    return ["Name", "Type", "Owner", "Status", "Tags", "Notes"]

# ----------------------------
# Slide builders
# ----------------------------
def add_title_slide(prs, app, org_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    name = clean_txt(app.get("name", "Application"))
    tagline = _truncate(app.get("tagline", ""), 80)
    accent_hex = clean_txt(app.get("accent", ACCENT))

    # Center card
    x, y, w, h = 2.0, 1.6, 9.333, 4.2
    add_rect(slide, x, y, w, h, SURFACE, BORDER, 1.0, 0.08)

    # Icon badge
    add_rect(slide, x + 0.5, y + 0.55, 0.9, 0.9, "#0F172A", BORDER, 1.0, 0.21)
    add_txt(slide, clean_txt(app.get("icon", "⬛")), x + 0.5, y + 0.72, 0.9, _text_h(1, 22), 22, TEXT_P, False, 2)

    add_txt(slide, name, x + 1.55, y + 0.62, w - 2.05, _text_h(1, 36), 36, TEXT_P, True, 1)
    add_txt(slide, tagline, x + 1.55, y + 1.25, w - 2.05, _text_h(2, 18), 18, TEXT_S, False, 1)

    # UI Wireframes label
    add_rect(slide, x + 1.55, y + 2.15, 2.2, 0.45, accent_hex, accent_hex, 1.0, 0.08)
    add_txt(slide, "UI Wireframes", x + 1.55, y + 2.25, 2.2, _text_h(1, 14), 14, "#FFFFFF", True, 2)

    # Footer org
    add_txt(slide, clean_txt(org_name), x, y + h - 0.55, w, _text_h(1, 12), 12, TEXT_S, False, 2)

def add_dashboard_slide(prs, app):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    name = clean_txt(app.get("name", "Application"))
    accent_hex = clean_txt(app.get("accent", ACCENT))

    add_top_nav(slide, name, accent_hex)
    add_sidebar(slide, accent_hex, "Dashboard")

    main_x = 0.3 + SIDEBAR_W + 0.25
    main_w = 13.0 - main_x
    y_cursor = 0.3 + NAV_H + 0.15

    # Page header
    header_h = add_page_header(slide, "Dashboard", "Home / Dashboard", main_x, y_cursor, main_w, accent_hex, ["+ New", "Export", "Filter"])
    y_cursor += header_h + SECTION_GAP

    # KPI row
    kpi_h = add_kpi_row(slide, _kpis(app), main_x, y_cursor, main_w, accent_hex)
    y_cursor += kpi_h + SECTION_GAP

    # Charts grid (2 columns, 2 rows)
    remaining_h = 7.2 - y_cursor
    chart_h = min(remaining_h, 3.6)
    add_charts_grid(slide, main_x, y_cursor, main_w, chart_h, _chart_labels(app))

def add_list_slide(prs, app):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    name = clean_txt(app.get("name", "Application"))
    accent_hex = clean_txt(app.get("accent", ACCENT))

    add_top_nav(slide, name, accent_hex)
    add_sidebar(slide, accent_hex, "Records")

    main_x = 0.3 + SIDEBAR_W + 0.25
    main_w = 13.0 - main_x
    y_cursor = 0.3 + NAV_H + 0.15

    item = _domain_item(app)
    header_h = add_page_header(slide, "Manage " + item, "Home / " + item, main_x, y_cursor, main_w, accent_hex, ["+ New", "Export", "Filter"])
    y_cursor += header_h + SECTION_GAP

    fb_h = add_filter_bar(slide, main_x, y_cursor, main_w, accent_hex)
    y_cursor += fb_h + SECTION_GAP

    # Table
    table_h = 7.2 - y_cursor
    table_h = max(2.2, min(3.9, table_h))
    add_table_placeholder(slide, _table_headers(app), main_x, y_cursor, main_w, table_h)

def add_form_slide(prs, app):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    name = clean_txt(app.get("name", "Application"))
    accent_hex = clean_txt(app.get("accent", ACCENT))

    add_top_nav(slide, name, accent_hex)
    add_sidebar(slide, accent_hex, "Settings")

    main_x = 0.3 + SIDEBAR_W + 0.25
    main_w = 13.0 - main_x
    y_cursor = 0.3 + NAV_H + 0.15

    item = _domain_item(app)
    header_h = add_page_header(slide, "Create New " + item[:-1] if item.endswith("s") else "Create New " + item, "Home / " + item + " / Create", main_x, y_cursor, main_w, accent_hex, ["Save", "Cancel", "Help"])
    y_cursor += header_h + SECTION_GAP

    # Form surface
    form_h = 7.2 - y_cursor - 0.2
    form_h = max(2.8, min(4.6, form_h))
    add_rect(slide, main_x, y_cursor, main_w, form_h, SURFACE, BORDER, 1.0, 0.08)

    pad = 0.25
    col_gap = 0.35
    col_w = (main_w - pad*2 - col_gap) / 2.0
    left_x = main_x + pad
    right_x = left_x + col_w + col_gap
    top_y = y_cursor + pad

    fields = _form_fields(app)
    # 2 columns, track y_cursor per column
    yL = top_y
    yR = top_y
    max_y = y_cursor + form_h - 0.9

    for i, f in enumerate(fields[:6]):
        if i % 2 == 0:
            if yL > max_y:
                continue
            used = add_input_field(slide, f, left_x, yL, col_w)
            if used > 0:
                yL += used + ROW_GAP
        else:
            if yR > max_y:
                continue
            used = add_input_field(slide, f, right_x, yR, col_w)
            if used > 0:
                yR += used + ROW_GAP

    # Buttons bottom-right inside form
    btn_y = y_cursor + form_h - 0.55
    btn_h = 0.42
    btn_w = 1.1
    gap = 0.15
    cancel_x = main_x + main_w - pad - btn_w
    save_x = cancel_x - gap - btn_w
    add_button(slide, "Save", save_x, btn_y, btn_w, btn_h, accent_hex, "#FFFFFF")
    add_button(slide, "Cancel", cancel_x, btn_y, btn_w, btn_h, "#0F172A", TEXT_S)

# ----------------------------
# Main
# ----------------------------
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    org_name = clean_txt(BRIEF.get("portfolio", {}).get("organisation", "Organisation"))

    apps = BRIEF.get("applications", [])[:5]
    for app in apps:
        add_title_slide(prs, app, org_name)
        add_dashboard_slide(prs, app)
        add_list_slide(prs, app)
        feats = app.get("features", []) or []
        if len(feats) >= 3:
            add_form_slide(prs, app)

    prs.save(