import argparse
import json
import os
import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

def rgb(hex_color):
    """Convert hex to RGBColor."""
    hex_color = hex_color.lstrip('#')
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16)
    )

def add_rect(slide, x, y, w, h, fill_hex, line_hex="#334155", line_w=1.0, radius=0.08):
    """Add a rounded rectangle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    # Set fill
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_hex)
    # Set border
    if line_hex:
        shape.line.color.rgb = rgb(line_hex)
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_txt(slide, text, x, y, w, h, size=12, color_hex="#F1F5F9", bold=False, align=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.text = str(text)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.font.size = Pt(size)
    p.font.color.rgb = rgb(color_hex)
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_kpi_tile(slide, x, y, w, h, kpi_data, colors):
    """Add a KPI card tile."""
    # Background
    bg_color = "#1E293B" # Surface color
    accent = kpi_data.get("accentColor", colors.get("primary", "#3B82F6"))
    
    add_rect(slide, x, y, w, h, bg_color, line_hex=accent, line_w=1.5)
    
    # Label
    add_txt(slide, kpi_data.get("title", "KPI"), x + 0.1, y + 0.1, w - 0.2, 0.3, size=10, color_hex="#94A3B8")
    
    # Value
    val = kpi_data.get("sampleValue", "0")
    add_txt(slide, val, x + 0.1, y + 0.35, w - 0.2, 0.5, size=24, color_hex="#F1F5F9", bold=True)
    
    # Change
    change = kpi_data.get("sampleChange", "")
    change_color = colors.get("success", "#107C10") if "+" in change else colors.get("danger", "#D13438")
    add_txt(slide, change, x + 0.1, y + 0.85, w - 0.2, 0.2, size=9, color_hex=change_color)

def add_chart(slide, x, y, w, h, chart_data, colors):
    """Add a chart object."""
    chart_type_map = {
        "line": XL_CHART_TYPE.LINE,
        "clustered_bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "horizontal_bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "stacked_bar": XL_CHART_TYPE.COLUMN_STACKED,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "pie": XL_CHART_TYPE.PIE,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "area": XL_CHART_TYPE.AREA,
        "radar": XL_CHART_TYPE.RADAR,
        "funnel": XL_CHART_TYPE.COLUMN_CLUSTERED # pptx doesn't have native funnel easily
    }
    
    c_type = chart_type_map.get(chart_data.get("type"), XL_CHART_TYPE.COLUMN_CLUSTERED)
    
    # Title rectangle (background for chart)
    add_rect(slide, x, y, w, h, "#1E293B", line_hex="#334155")
    add_txt(slide, chart_data.get("title", ""), x + 0.2, y + 0.1, w - 0.4, 0.4, size=14, color_hex="#F1F5F9", bold=True)
    
    # Chart Data
    cfg = chart_data.get("cfg", {})
    categories = cfg.get("cat", ["A", "B", "C", "D"])
    values = cfg.get("data", [10, 20, 30, 40])
    
    if not values or len(values) == 0:
        values = [10, 20, 30, 40] # Fallback
        
    try:
        # Chart is placed slightly inside the card
        chart_x, chart_y = x + 0.2, y + 0.6
        chart_w, chart_h = w - 0.4, h - 0.8
        
        if c_type == XL_CHART_TYPE.XY_SCATTER:
            from pptx.chart.data import XyChartData
            chart_payload = XyChartData()
            series = chart_payload.add_series(chart_data.get("title", ""))
            # Scatter expects X, Y pairs. If we only have category/val, we use category index as X.
            for i, val in enumerate(values):
                series.add_data_point(i, val)
        else:
            chart_payload = CategoryChartData()
            chart_payload.categories = categories
            series_label = cfg.get("label", chart_data.get("title", "Series 1"))
            chart_payload.add_series(series_label, values)

        chart_graphic = slide.shapes.add_chart(
            c_type, Inches(chart_x), Inches(chart_y), Inches(chart_w), Inches(chart_h), chart_payload
        )
        chart = chart_graphic.chart
        
        # Basic styling
        chart.has_legend = True
        chart.legend.font.size = Pt(8)
        chart.legend.font.color.rgb = rgb("#94A3B8")
        
        # Plot area styling
        try:
            chart.plot_area.format.fill.solid()
            chart.plot_area.format.fill.fore_color.rgb = rgb("#1E293B")
        except:
            pass
            
        # Axis styling (not all charts have axes)
        try:
            category_axis = chart.category_axis
            category_axis.tick_labels.font.size = Pt(8)
            category_axis.tick_labels.font.color.rgb = rgb("#94A3B8")
            
            value_axis = chart.value_axis
            value_axis.tick_labels.font.size = Pt(8)
            value_axis.tick_labels.font.color.rgb = rgb("#94A3B8")
        except:
            pass
            
    except Exception as e:
        print(f"Error adding chart: {e}")
        # Add a placeholder if chart fails
        add_rect(slide, x+0.2, y+0.6, w-0.4, h-0.8, "#0F172A", line_hex="#334155")
        add_txt(slide, f"Chart Placeholder: {chart_data.get('type')}", x+0.4, y+h/2, w-0.8, 0.5, align=PP_ALIGN.CENTER)

def main():
    parser = argparse.ArgumentParser(description="Deterministic Dashboard Layout Engine")
    parser.add_argument("--payload", help="JSON payload file")
    parser.add_argument("--output", required=True, help="Output PPTX filename")
    args = parser.parse_args()

    # Load payload
    if args.payload:
        with open(args.payload, 'r') as f:
            data = json.load(f)
    else:
        # If no payload file, try reading from stdin
        try:
            data = json.loads(sys.stdin.read())
        except:
            print("Error: No valid payload provided.")
            return

    prs = Presentation()
    
    # Set slide dimensions for 16:9 (13.333 x 7.5 inches)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    colors = data.get("colors", {
        "primary": "#3B82F6",
        "secondary": "#1E40AF",
        "success": "#107C10",
        "warning": "#FFB900",
        "danger": "#D13438"
    })
    
    client_name = data.get("clientName", "IntelliFrame")
    industry = data.get("industry", "Business")
    page_count = data.get("pageCount", 1)
    
    # Map visualizations by page
    pages = {}
    for viz in data.get("visualizations", []):
        p_num = viz.get("page", 1)
        if p_num not in pages:
            pages[p_num] = []
        pages[p_num].append(viz)
    
    # Generate slides
    for p_num in range(1, page_count + 1):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = rgb("#0F172A") # Dark Navy
        
        # Header
        add_rect(slide, 0, 0, 13.333, 0.6, "#1E293B", line_hex=None)
        add_txt(slide, f"{client_name} | {industry} Analytics", 0.3, 0.1, 8, 0.4, size=18, bold=True)
        add_txt(slide, f"Page {p_num} of {page_count}", 11.5, 0.1, 1.5, 0.4, size=10, align=PP_ALIGN.RIGHT)
        
        # Filter Bar Placeholder
        add_rect(slide, 0.3, 0.7, 12.7, 0.4, "#0F172A", line_hex="#334155")
        add_txt(slide, " Filters: Date Range [2026-01-01 - 2026-04-02] | Region: All | Segment: All ", 0.4, 0.75, 12, 0.3, size=9)
        
        # Layout components for this page
        page_viz = pages.get(p_num, [])
        
        # Grid constants
        MARGIN_X = 0.3
        MARGIN_Y = 1.3 # Below header and filter bar
        PAGE_W = 12.7
        PAGE_H = 5.9
        GAP = 0.2
        
        # Calculate row heights dynamically or use a base
        # In this simplistic deterministic version, we'll follow the row/col logic from the prompt
        # 4 columns grid
        COL_W = (PAGE_W - (3 * GAP)) / 4
        
        for viz in page_viz:
            row_idx = viz.get("row", 1)
            width_str = viz.get("width", "25%")
            
            # Simple heuristic for row Y position
            # Cards are usually row 1 (0.3 y), Charts start row 2 (1.5 y)
            # This logic should ideally be more robust based on previous rows
            if row_idx == 1:
                row_y = MARGIN_Y
                row_h = 1.2
            else:
                row_y = MARGIN_Y + 1.2 + GAP + ((row_idx - 2) * (2.2 + GAP))
                row_h = 2.2
            
            # X position based on "index" in row or sequential placement
            # The layout-engine is supposed to provide row/width
            # We'll assume the visualizations are ordered correctly by row then position
            # This deterministic engine will place them side-by-side in each row
            
            # Track current X for each row
            if not hasattr(main, 'row_x_map'):
                main.row_x_map = {}
            row_key = f"{p_num}_{row_idx}"
            cur_x = main.row_x_map.get(row_key, MARGIN_X)
            
            if width_str == "25%":
                w = COL_W
            elif width_str == "50%":
                w = (COL_W * 2) + GAP
            elif width_str == "100%":
                w = PAGE_W
            else:
                w = COL_W # Default
                
            if viz.get("type") == "kpi-card":
                add_kpi_tile(slide, cur_x, row_y, w, row_h, viz, colors)
            else:
                add_chart(slide, cur_x, row_y, w, row_h, viz, colors)
                
            main.row_x_map[row_key] = cur_x + w + GAP

    # Save presentation
    prs.save(args.output)
    print(f"Finished generating {args.output}")

if __name__ == "__main__":
    main()
