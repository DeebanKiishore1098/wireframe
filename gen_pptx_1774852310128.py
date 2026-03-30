#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ──────────────────────────────────────────────────────────────────────────────
# Embedded input JSON (STRICT) as BRIEF
# ──────────────────────────────────────────────────────────────────────────────
BRIEF = {
  "portfolio": {
    "title": "Timex Ecom BI Solution (Microsoft Fabric + Power BI) — Application Portfolio Brief",
    "organisation": "SYSTECH",
    "year": "2025",
    "confidential": True,
    "stats": {
      "totalApps": "1",
      "domains": "E-commerce Analytics",
      "aiPowered": "No",
      "buildTenure": "Not specified"
    }
  },
  "applications": [
    {
      "id": "01",
      "name": "Timex Ecom BI & Reporting Platform",
      "domain": "E-commerce Analytics",
      "category": "Business Intelligence / Reporting",
      "tagline": "Microsoft Fabric-based data integration and Power BI dashboards for Timex India e-commerce sales performance",
      "overview": "Implements a comprehensive BI solution for Timex Group India Limited using Microsoft Fabric to automate data integration and ETL, consolidating multi-partner e-commerce sales data and delivering Power BI dashboards for primary/secondary sales and brand/product performance. The solution supports tracking sales patterns, inventory turnover, and customer/location dimensions across 1P, 3P, and Dotcom channels in India.",
      "features": [
        {
          "title": "Multi-source data ingestion",
          "description": "Consolidates data from Oracle, Bizom, Olabi, and CSV files from partner portals plus emailed datasets into a unified analytics layer."
        },
        {
          "title": "Fabric ETL automation to OneLake",
          "description": "Automates integration and transformation workflows in Microsoft Fabric with curated outputs landed in OneLake."
        },
        {
          "title": "Power BI dashboards for sales visibility",
          "description": "Interactive reporting for primary sales (region/product/customer), secondary sales (channel/product), and brand/product performance."
        },
        {
          "title": "Standardized KPI model",
          "description": "Defines core KPI calculations including Net Price (SP*Qty), NSV (SP-10%), NRV (SP-GST), and GM%."
        },
        {
          "title": "Channel and partner segmentation",
          "description": "Supports analysis across partner types (1P/3P/Dotcom) and portals (e.g., Flipkart, Myntra, JW.com, Shantam, Timehut, NDSL, VRPT)."
        },
        {
          "title": "Row-level security framework",
          "description": "RLS roles defined for All Portals, Regional (1P), State (3P), and Dotcom to restrict visibility by seller/channel grouping."
        }
      ],
      "impact": "Improves decision-making and operational efficiency by centralizing e-commerce sales data, standardizing KPIs, and enabling secure self-service insights across channels and brand groups.",
      "metrics": [
        {
          "label": "Coverage",
          "value": "India (1P, 3P, Dotcom)"
        },
        {
          "label": "Refresh Frequency",
          "value": "Quarterly (CSV sources)"
        },
        {
          "label": "Version",
          "value": "v1.0"
        },
        {
          "label": "Platform",
          "value": "Microsoft Fabric + Power BI"
        }
      ],
      "accent": "#2563EB",
      "icon": "📊"
    }
  ]
}

# ──────────────────────────────────────────────────────────────────────────────
# Design System (Dark Mode)
# ──────────────────────────────────────────────────────────────────────────────
BG = RGBColor(0x08, 0x0C, 0x14)          # #080C14
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xC7, 0xD2, 0xFE)       # light indigo-ish
LIGHT_GRAY = RGBColor(0xD1, 0xD5, 0xDB)  # gray-300
CARD = RGBColor(0x0F, 0x17, 0x2A)        # slate-900-ish
CARD2 = RGBColor(0x0B, 0x12, 0x24)
ACCENT_BLUE = RGBColor(0x25, 0x63, 0xEB) # #2563EB
ACCENT_INDIGO = RGBColor(0x4F, 0x46, 0xE5) # #4F46E5

FONT_HEAD = "Calibri"   # Inter-like fallback
FONT_BODY = "Calibri"

def hex_to_rgb(hex_str: str) -> RGBColor:
    s = (hex_str or "").strip()
    if s.startswith("#"):
        s = s[1:]
    if len(s) != 6:
        return ACCENT_BLUE
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))

def set_slide_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def style_shape_fill(shape, fill_rgb: RGBColor, line_rgb=None, line_width_pt=1.0):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_rgb
    line = shape.line
    if line_rgb is None:
        line.color.rgb = fill_rgb
        line.width = Pt(0)
    else:
        line.color.rgb = line_rgb
        line.width = Pt(line_width_pt)

def add_textbox(slide, x, y, w, h, text, font_size=18, bold=False, color=WHITE,
                font_name=FONT_BODY, align=None, line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if align is not None:
        p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    return tb

def add_round_card(slide, x, y, w, h, fill_rgb=CARD, line_rgb=RGBColor(0x1F, 0x29, 0x3D), radius_shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(radius_shape, x, y, w, h)
    style_shape_fill(shp, fill_rgb, line_rgb=line_rgb, line_width_pt=1.0)
    return shp

def set_shape_text(shape, lines, font_sizes, colors, bolds=None, font_name=FONT_BODY, left=Inches(0.25), top=Inches(0.15)):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = left
    tf.margin_right = left
    tf.margin_top = top
    tf.margin_bottom = top
    if bolds is None:
        bolds = [False] * len(lines)

    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = txt
        run.font.name = font_name
        run.font.size = Pt(font_sizes[i] if i < len(font_sizes) else font_sizes[-1])
        run.font.bold = bool(bolds[i] if i < len(bolds) else bolds[-1])
        run.font.color.rgb = colors[i] if i < len(colors) else colors[-1]

def add_badge(slide, x, y, w, h, text, fill_rgb, text_rgb=WHITE, font_size=12):
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    style_shape_fill(badge, fill_rgb, line_rgb=None)
    tf = badge.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_BODY
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = text_rgb
    return badge

def add_bullets(shape, title, items, title_color=WHITE, bullet_color=LIGHT_GRAY, accent=ACCENT_BLUE):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.2)
    tf.margin_bottom = Inches(0.15)

    # Title
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    r0.font.name = FONT_HEAD
    r0.font.size = Pt(16)
    r0.font.bold = True
    r0.font.color.rgb = title_color
    p0.space_after = Pt(10)

    # Bullets
    for it in items:
        p = tf.add_paragraph()
        p.level = 0
        p.text = ""  # use runs to control colors
        p.space_after = Pt(6)

        r = p.add_run()
        r.text = "• "
        r.font.name = FONT_BODY
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = accent

        r2 = p.add_run()
        r2.text = it
        r2.font.name = FONT_BODY
        r2.font.size = Pt(12)
        r2.font.color.rgb = bullet_color

def build_presentation(output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 widescreen
    prs.slide_height = Inches(7.5)

    # Use blank layout
    blank = prs.slide_layouts[6]

    portfolio = BRIEF["portfolio"]
    apps = BRIEF["applications"]

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 1: HERO
    # ──────────────────────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1, BG)

    # Top accent bar
    bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.18))
    style_shape_fill(bar, ACCENT_INDIGO, line_rgb=None)

    # Title block
    add_textbox(
        s1, Inches(0.8), Inches(0.75), Inches(12.0), Inches(1.3),
        portfolio["title"],
        font_size=34, bold=True, color=WHITE, font_name=FONT_HEAD
    )

    sub = f'{portfolio["organisation"]}  •  {portfolio["year"]}'
    add_textbox(
        s1, Inches(0.82), Inches(1.75), Inches(9.0), Inches(0.5),
        sub,
        font_size=16, bold=False, color=LIGHT_GRAY, font_name=FONT_BODY
    )

    if portfolio.get("confidential", False):
        add_badge(s1, Inches(10.9), Inches(1.6), Inches(1.9), Inches(0.42), "CONFIDENTIAL", fill_rgb=RGBColor(0xB9, 0x1C, 0x1C), text_rgb=WHITE, font_size=12)

    # Stats dashboard row
    stats = portfolio.get("stats", {})
    stat_items = [
        ("Total Apps", stats.get("totalApps", "")),
        ("Domains", stats.get("domains", "")),
        ("AI Powered", stats.get("aiPowered", "")),
        ("Build Tenure", stats.get("buildTenure", "")),
    ]

    row_y = Inches(2.55)
    card_w = Inches(3.05)
    card_h = Inches(1.1)
    gap = Inches(0.35)
    x0 = Inches(0.8)

    for i, (k, v) in enumerate(stat_items):
        x = x0 + i * (card_w + gap)
        card = add_round_card(s1, x, row_y, card_w, card_h, fill_rgb=CARD, line_rgb=RGBColor(0x1E, 0x29, 0x3B))
        set_shape_text(
            card,
            [k.upper(), v],
            [11, 22],
            [MUTED, WHITE],
            bolds=[True, True],
            font_name=FONT_BODY,
            left=Inches(0.25),
            top=Inches(0.12),
        )

    # Footer hint
    add_textbox(
        s1, Inches(0.8), Inches(6.9), Inches(12.0), Inches(0.4),
        "Application Portfolio • Dark Mode Brief",
        font_size=11, bold=False, color=RGBColor(0x94, 0xA3, 0xB8), font_name=FONT_BODY
    )

    # ──────────────────────────────────────────────────────────────────────────
    # Slide 2: CATALOG (grid of cards)
    # ──────────────────────────────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2, BG)

    add_textbox(s2, Inches(0.8), Inches(0.55), Inches(12.0), Inches(0.7),
                "Application Catalog", font_size=28, bold=True, color=WHITE, font_name=FONT_HEAD)
    add_textbox(s2, Inches(0.82), Inches(1.15), Inches(12.0), Inches(0.4),
                "Overview of applications in this portfolio", font_size=13, bold=False, color=LIGHT_GRAY, font_name=FONT_BODY)

    # Grid settings
    left = Inches(0.8)
    top = Inches(1.75)
    cols = 2
    card_w = Inches(6.05)
    card_h = Inches(1.65)
    hgap = Inches(0.45)
    vgap = Inches(0.4)

    for idx, app in enumerate(apps):
        r = idx // cols
        c = idx % cols
        x = left + c * (card_w + hgap)
        y = top + r * (card_h + vgap)

        accent = hex_to_rgb(app.get("accent"))
        card = add_round_card(s2, x, y, card_w, card_h, fill_rgb=CARD2, line_rgb=RGBColor(0x1E, 0x29, 0x3B))

        # Accent stripe
        stripe = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.12), card_h)
        style_shape_fill(stripe, accent, line_rgb=None)

        # Icon "bubble"
        icon_box = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.35), y + Inches(0.35), Inches(0.95), Inches(0.95))
        style_shape_fill(icon_box, RGBColor(0x11, 0x22, 0x44), line_rgb=accent, line_width_pt=1.5)
        tf = icon_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = app.get("icon", "⬤")
        run.font.size = Pt(28)
        run.font.name = FONT_HEAD
        run.font.bold = True
        run.font.color.rgb = WHITE

        # Text on card
        name = f'{app.get("id", "").strip()}  {app.get("name", "")}'.strip()
        domain = app.get("domain", "")
        category = app.get("category", "")

        title_tb = s2.shapes.add_textbox(x + Inches(1.45), y + Inches(0.28), card_w - Inches(1.7), Inches(0.6))
        ttf = title_tb.text_frame
        ttf.clear()
        p0 = ttf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = name
        r0.font.name = FONT_HEAD
        r0.font.size = Pt(16)
        r0.font.bold = True
        r0.font.color.rgb = WHITE

        sub_tb = s2.shapes.add_textbox(x + Inches(1.45), y + Inches(0.88), card_w - Inches(1.7), Inches(0.6))
        stf = sub_tb.text_frame
        stf.clear()
        p1 = stf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = f"{domain}  •  {category}"
        r1.font.name = FONT_BODY
        r1.font.size = Pt(12)
        r1.font.color.rgb = LIGHT_GRAY

        # Accent pill (domain)
        add_badge(s2, x + card_w - Inches(2.05), y + Inches(0.28), Inches(1.75), Inches(0.38),
                  "PORTFOLIO APP", fill_rgb=accent, text_rgb=WHITE, font_size=11)

    # ──────────────────────────────────────────────────────────────────────────
    # Detail slides (one per application)
    # ──────────────────────────────────────────────────────────────────────────
    for app in apps:
        s = prs.slides.add_slide(blank)
        set_slide_bg(s, BG)
        accent = hex_to_rgb(app.get("accent"))

        # Header bar + title
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.18))
        style_shape_fill(bar, accent, line_rgb=None)

        add_textbox(s, Inches(0.8), Inches(0.55), Inches(11.7), Inches(0.7),
                    app.get("name", ""), font_size=28, bold=True, color=WHITE, font_name=FONT_HEAD)

        add_textbox(s, Inches(0.82), Inches(1.15), Inches(11.0), Inches(0.5),
                    app.get("tagline", ""), font_size=13, bold=False, color=LIGHT_GRAY, font_name=FONT_BODY)

        # Right top badge
        add_badge(s, Inches(10.95), Inches(0.58), Inches(2.0), Inches(0.42),
                  app.get("domain", "Domain"), fill_rgb=RGBColor(0x11, 0x22, 0x44), text_rgb=WHITE, font_size=11)

        # Layout: Left (Overview & Impact), Right (Features & Metrics)
        margin_x = Inches(0.8)
        content_top = Inches(1.75)
        content_h = Inches(5.45)
        left_w = Inches(6.2)
        right_w = Inches(5.55)
        gap = Inches(0.45)

        left_x = margin_x
        right_x = margin_x + left_w + gap

        # Left card
        left_card = add_round_card(s, left_x, content_top, left_w, content_h, fill_rgb=CARD, line_rgb=RGBColor(0x1E, 0x29, 0x3B))

        # Icon bubble on left
        icon = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_x + Inches(0.35), content_top + Inches(0.35), Inches(1.05), Inches(1.05))
        style_shape_fill(icon, RGBColor(0x0B, 0x1B, 0x3A), line_rgb=accent, line_width_pt=1.5)
        tf = icon.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = app.get("icon", "⬤")
        r.font.name = FONT_HEAD
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = WHITE

        # Overview title
        ov_title = s.shapes.add_textbox(left_x + Inches(1.55), content_top + Inches(0.35), left_w - Inches(1.9), Inches(0.45))
        ovtf = ov_title.text_frame
        ovtf.clear()
        p0 = ovtf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = "Overview"
        r0.font.name = FONT_HEAD
        r0.font.size = Pt(16)
        r0.font.bold = True
        r0.font.color.rgb = WHITE

        # Overview body
        ov_body = s.shapes.add_textbox(left_x + Inches(0.35), content_top + Inches(1.55), left_w - Inches(0.7), Inches(2.55))
        ovbf = ov_body.text_frame
        ovbf.clear()
        ovbf.word_wrap = True
        p1 = ovbf.paragraphs[0]
        p1.text = app.get("overview", "")
        p1.font.name = FONT_BODY
        p1.font.size = Pt(12.5)
        p1.font.color.rgb = LIGHT_GRAY
        p1.space_after = Pt(8)

        # Impact section
        impact_box = add_round_card(s, left_x + Inches(0.35), content_top + Inches(4.25), left_w - Inches(0.7), Inches(1.05),
                                    fill_rgb=RGBColor(0x0B, 0x12, 0x24), line_rgb=accent, line_width_pt=1.2)
        itf = impact_box.text_frame
        itf.clear()
        itf.word_wrap = True
        itf.margin_left = Inches(0.25)
        itf.margin_right = Inches(0.25)
        itf.margin_top = Inches(0.12)
        itf.margin_bottom = Inches(0.12)

        pI0 = itf.paragraphs[0]
        rI0 = pI0.add_run()
        rI0.text = "Impact"
        rI0.font.name = FONT_HEAD
        rI0.font.size = Pt(14)
        rI0.font.bold = True
        rI0.font.color.rgb = WHITE
        pI0.space_after = Pt(4)

        pI1 = itf.add_paragraph()
        rI1 = pI1.add_run()
        rI1.text = app.get("impact", "")
        rI1.font.name = FONT_BODY
        rI1.font.size = Pt(12)
        rI1.font.color.rgb = LIGHT_GRAY

        # Right column: features card
        feat_card = add_round_card(s, right_x, content_top, right_w, Inches(3.85), fill_rgb=CARD, line_rgb=RGBColor(0x1E, 0x29, 0x3B))
        feature_lines = []
        for f in app.get("features", []):
            t = (f.get("title") or "").strip()
            d = (f.get("description") or "").strip()
            feature_lines.append(f"{t}: {d}" if d else t)

        add_bullets(feat_card, "Key Features", feature_lines, title_color=WHITE, bullet_color=LIGHT_GRAY, accent=accent)

        # Right column: metrics card
        metrics_card = add_round_card(s, right_x, content_top + Inches(4.05), right_w, Inches(1.4),
                                      fill_rgb=RGBColor(0x0B, 0x12, 0x24), line_rgb=RGBColor(0x1E, 0x29, 0x3B))

        mtf = metrics_card.text_frame
        mtf.clear()
        mtf.word_wrap = True
        mtf.margin_left = Inches(0.25)
        mtf.margin_right = Inches(0.25)
        mtf.margin_top = Inches(0.15)
        mtf.margin_bottom = Inches(0.12)

        pM0 = mtf.paragraphs[0]
        rM0 = pM0.add_run()
        rM0.text = "Metrics"
        rM0.font.name = FONT_HEAD
        rM0.font.size = Pt(14)
        rM0.font.bold = True
        rM0.font.color.rgb = WHITE
        pM0.space_after = Pt(6)

        # Render metrics in two columns within the metrics card using separate textboxes
        metrics = app.get("metrics", [])
        # Fallback to simple list if not enough space
        col1 = metrics[0::2]
        col2 = metrics[1::2]

        m_x = right_x + Inches(0.25)
        m_y = content_top + Inches(4.45)
        m_w = (right_w - Inches(0.6)) / 2
        m_h = Inches(0.9)

        def metric_block(slide, x, y, w, h, pairs):
            tb = slide.shapes.add_textbox(x, y, w, h)
            tf = tb.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            for i, m in enumerate(pairs):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_after = Pt(6)
                label = m.get("label", "")
                value = m.get("value", "")
                r1 = p.add_run()
                r1.text = f"{label}: "
                r1.font.name = FONT_BODY
                r1.font.size = Pt(11.5)
                r1.font.bold = True
                r1.font.color.rgb = MUTED

                r2 = p.add_run()
                r2.text = value
                r2.font.name = FONT_BODY
                r2.font.size = Pt(11.5)
                r2.font.bold = False
                r2.font.color.rgb = WHITE

        metric_block(s, m_x, m_y, m_w, m_h, col1)
        metric_block(s, m_x + m_w + Inches(0.1), m_y, m_w, m_h, col2)

        # Footer
        add_textbox(s, Inches(0.8), Inches(7.05), Inches(12.0), Inches(0.35),
                    f'{portfolio["organisation"]} • {portfolio["year"]} • Application Detail {app.get("id","")}',
                    font_size=10.5, bold=False, color=RGBColor(0x94, 0xA3, 0xB8), font_name=FONT_BODY)

    prs.save(output_path)

def main():
    parser = argparse.ArgumentParser(description="Generate Application Portfolio PPTX (Dark Mode) using python-pptx.")
    parser.add_argument("--output", required=True, help="Output PPTX file path")
    args = parser.parse_args()
    build_presentation(args.output)

if __name__ == "__main__":
    main()