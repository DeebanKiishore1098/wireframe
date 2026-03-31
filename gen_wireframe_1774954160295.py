import argparse
import re
from math import ceil

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# =========================
# Renderer Utilities (MUST DEFINE)
# =========================
def clean_txt(t):
    t = "" if t is None else str(t)
    return re.sub(r"<[^>]*>", "", t)


def _hex_to_rgb(hex_color):
    hx = clean_txt(hex_color).strip().lstrip("#")
    if len(hx) != 6:
        hx = "FFFFFF"
    return RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))


def set_bg(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def add_txt(slide, text, x, y, w, h, size, color_hex, bold=False, align=1):
    # positional-args-only contract enforced by caller usage (no kwargs)
    tx = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean_txt(text)

    font = run.font
    font.name = "Helvetica"
    font.size = Pt(float(size))
    font.bold = bool(bold)
    font.color.rgb = _hex_to_rgb(color_hex)

    if align == 2:
        p.alignment = PP_ALIGN.CENTER
    elif align == 3:
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    return tx


def add_rect(slide, x, y, w, h, fill_hex, line_hex, line_w=1.0, r=0.1):
    # Rounded rectangle surface card
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex_to_rgb(fill_hex)

    shp.line.color.rgb = _hex_to_rgb(line_hex)
    shp.line.width = Pt(float(line_w))

    try:
        # python-pptx uses 0..1 adjustment for roundness on rounded rects
        shp.adjustments[0] = max(0.0, min(1.0, float(r)))
    except Exception:
        pass

    return shp


def add_pill(slide, x, y, text, bg_hex, fg_hex, size=10):
    # Simple rounded badge
    t = clean_txt(text)
    if not t:
        t = "—"
    # Approx width calc on 8pt grid; keep luxury spacing.
    w = max(0.9, min(3.2, 0.18 * len(t) + 0.7))
    h = 0.34
    pill = add_rect(slide, x, y, w, h, bg_hex, bg_hex, 0.0, 0.8)
    add_txt(slide, t, x + 0.18, y + 0.06, w - 0.36, h - 0.08, size, fg_hex, True, 1)
    return pill


# =========================
# Design Tokens (Premium Design System v2.0)
# =========================
BASE_BG = "#080C14"     # 60%
SURFACE = "#0F172A"     # 30%
ACCENT = "#2563EB"      # 10% (Electric Blue - no purple)

TEXT_PRIMARY = "#E6EDF7"
TEXT_MUTED = "#9AA6BC"
LINE_SOFT = "#1D2A44"
BASE_DARK = "#05080F"

GOLD = 1.618

# Slide geometry (16:9)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Layout columns (Golden Ratio)
CONTENT_W = SLIDE_W_IN * 0.62
SIDEBAR_W = SLIDE_W_IN * 0.38

# Margins (8-pt grid inspired; in inches ~ 0.0833 per 6pt; here use clean 0.5/0.75)
MARGIN_X = 0.75
MARGIN_Y = 0.60
GUTTER = 0.40

CONTENT_X = MARGIN_X
SIDEBAR_X = CONTENT_X + CONTENT_W + GUTTER

TOPBAR_H = 0.55  # subtle header area height


def _safe_get(d, path, default=""):
    cur = d
    try:
        for key in path:
            cur = cur[key]
        return cur
    except Exception:
        return default


def _chunk_text(text, max_chars=170):
    t = clean_txt(text).strip()
    if not t:
        return []
    words = t.split()
    out = []
    buf = []
    count = 0
    for w in words:
        if count + len(w) + (1 if buf else 0) > max_chars:
            out.append(" ".join(buf))
            buf = [w]
            count = len(w)
        else:
            buf.append(w)
            count += len(w) + (1 if buf else 0)
    if buf:
        out.append(" ".join(buf))
    return out


def _fit_font(base, min_sz=10, max_sz=54):
    return max(min_sz, min(max_sz, float(base)))


# =========================
# Slide Templates
# =========================
def slide_hero(prs, brief):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, BASE_DARK)

    # Top subtle rule
    add_rect(slide, MARGIN_X, MARGIN_Y - 0.12, SLIDE_W_IN - 2 * MARGIN_X, 0.02, LINE_SOFT, LINE_SOFT, 0.0, 0.0)

    title = clean_txt(_safe_get(brief, ["portfolio", "title"], "Application Portfolio"))
    org = clean_txt(_safe_get(brief, ["portfolio", "organisation"], ""))
    year = clean_txt(_safe_get(brief, ["portfolio", "year"], ""))
    confidential = _safe_get(brief, ["portfolio", "confidential"], False)

    # Left content block
    left_x = CONTENT_X
    left_w = CONTENT_W

    # Title sizing: Heading = Body * 1.618
    body_sz = 22
    heading_sz = _fit_font(body_sz * GOLD, 44, 56)  # per spec "Pt 44+"

    add_txt(slide, title, left_x, 1.25, left_w, 1.35, heading_sz, TEXT_PRIMARY, True, 1)

    subline_parts = [p for p in [org, year] if p]
    subline = " · ".join(subline_parts) if subline_parts else " "
    add_txt(slide, subline, left_x, 2.55, left_w, 0.50, body_sz, TEXT_MUTED, False, 1)

    if confidential:
        add_pill(slide, left_x, 3.05, "CONFIDENTIAL", ACCENT, "#FFFFFF", 10)
        add_txt(slide, "Internal distribution only.", left_x + 1.45, 3.05, left_w - 1.6, 0.40, 12, TEXT_MUTED, False, 1)

    # Right stats panel (4 cards)
    stats = _safe_get(brief, ["portfolio", "stats"], {}) or {}
    stat_items = [
        ("Total Apps", _safe_get(stats, ["totalApps"], "—")),
        ("Domains", _safe_get(stats, ["domains"], "—")),
        ("AI / LLM", _safe_get(stats, ["aiPowered"], "—")),
        ("Build Tenure", _safe_get(stats, ["buildTenure"], "—")),
    ]
    card_w = SIDEBAR_W - (MARGIN_X)  # within right region; keep internal padding
    card_x = SIDEBAR_X
    top_y = 1.10
    gap_y = 0.22
    card_h = 1.18

    # Sidebar title
    add_txt(slide, "Portfolio Highlights", card_x, top_y - 0.55, card_w, 0.35, 14, TEXT_MUTED, True, 1)

    for i, (lbl, val) in enumerate(stat_items):
        y = top_y + i * (card_h + gap_y)
        add_rect(slide, card_x, y, card_w, card_h, SURFACE, LINE_SOFT, 1.0, 0.18)
        add_txt(slide, clean_txt(lbl).upper(), card_x + 0.30, y + 0.22, card_w - 0.60, 0.28, 10, TEXT_MUTED, True, 1)
        add_txt(slide, clean_txt(val), card_x + 0.30, y + 0.52, card_w - 0.60, 0.55, 22, TEXT_PRIMARY, True, 1)

    # Footer
    add_txt(slide, "Systech Analytics · Application Portfolio", MARGIN_X, SLIDE_H_IN - 0.55, SLIDE_W_IN - 2 * MARGIN_X, 0.35, 10, TEXT_MUTED, False, 1)
    return slide


def slide_catalog(prs, brief):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BASE_DARK)

    apps = _safe_get(brief, ["applications"], []) or []
    title = clean_txt(_safe_get(brief, ["portfolio", "title"], "Application Catalog"))

    # Header
    add_txt(slide, "CATALOG", MARGIN_X, MARGIN_Y, SLIDE_W_IN - 2 * MARGIN_X, 0.35, 10, TEXT_MUTED, True, 1)
    add_txt(slide, title, MARGIN_X, MARGIN_Y + 0.25, SLIDE_W_IN - 2 * MARGIN_X, 0.60, 26, TEXT_PRIMARY, True, 1)

    # Asymmetric layout: left "filter/legend" panel + right grid of cards
    panel_x = MARGIN_X
    panel_y = 1.40
    panel_w = 3.45
    panel_h = SLIDE_H_IN - panel_y - MARGIN_Y

    grid_x = panel_x + panel_w + GUTTER
    grid_y = panel_y
    grid_w = SLIDE_W_IN - MARGIN_X - grid_x
    grid_h = panel_h

    add_rect(slide, panel_x, panel_y, panel_w, panel_h, SURFACE, LINE_SOFT, 1.0, 0.16)
    add_txt(slide, "Overview", panel_x + 0.30, panel_y + 0.26, panel_w - 0.60, 0.35, 14, TEXT_PRIMARY, True, 1)
    add_txt(
        slide,
        clean_txt("All applications in this portfolio, grouped as cards for rapid scanning. Detail slides follow for each app."),
        panel_x + 0.30,
        panel_y + 0.70,
        panel_w - 0.60,
        0.95,
        11,
        TEXT_MUTED,
        False,
        1,
    )

    # Domain quick chips (top 8 unique)
    domains = []
    for a in apps:
        d = clean_txt(_safe_get(a, ["domain"], ""))
        if d and d not in domains:
            domains.append(d)
    domains = domains[:8]

    add_txt(slide, "Domains", panel_x + 0.30, panel_y + 1.70, panel_w - 0.60, 0.30, 11, TEXT_MUTED, True, 1)
    px, py = panel_x + 0.30, panel_y + 2.05
    for idx, d in enumerate(domains):
        add_pill(slide, px, py, d, "#0B1222", TEXT_PRIMARY, 9)
        py += 0.44
        if py > panel_y + panel_h - 0.60:
            break

    # Grid cards (no overlaps; 8-pt grid feel via consistent gaps)
    card_gap = 0.22
    cols = 3
    card_w = (grid_w - (cols - 1) * card_gap) / cols
    card_h = 1.55

    rows = max(1, int((grid_h + card_gap) // (card_h + card_gap)))
    capacity = rows * cols
    visible_apps = apps[:capacity]

    for i, app in enumerate(visible_apps):
        r = i // cols
        c = i % cols
        x = grid_x + c * (card_w + card_gap)
        y = grid_y + r * (card_h + card_gap)

        accent = clean_txt(_safe_get(app, ["accent"], ACCENT)) or ACCENT
        name = clean_txt(_safe_get(app, ["name"], "Untitled"))
        aid = clean_txt(_safe_get(app, ["id"], ""))
        domain = clean_txt(_safe_get(app, ["domain"], ""))
        cat = clean_txt(_safe_get(app, ["category"], ""))
        tagline = clean_txt(_safe_get(app, ["tagline"], ""))

        add_rect(slide, x, y, card_w, card_h, SURFACE, LINE_SOFT, 1.0, 0.16)

        # Accent bar
        add_rect(slide, x, y, card_w, 0.06, accent, accent, 0.0, 0.0)

        # Top row: ID + pills
        add_txt(slide, f"{aid}".strip(), x + 0.20, y + 0.16, 0.60, 0.30, 10, TEXT_MUTED, True, 1)
        # Category pill (right aligned visually)
        pill_text = cat if cat else "Category"
        add_pill(slide, x + card_w - 1.90, y + 0.14, pill_text, "#0B1222", TEXT_PRIMARY, 9)

        add_txt(slide, name, x + 0.20, y + 0.50, card_w - 0.40, 0.35, 14, TEXT_PRIMARY, True, 1)
        add_txt(slide, domain, x + 0.20, y + 0.86, card_w - 0.40, 0.28, 10, TEXT_MUTED, True, 1)

        # Tagline (2 lines max)
        tag = tagline
        if len(tag) > 95:
            tag = tag[:92].rstrip() + "…"
        add_txt(slide, tag, x + 0.20, y + 1.10, card_w - 0.40, 0.40, 10.5, TEXT_MUTED, False, 1)

    # Footer
    add_txt(slide, f"Total applications: {clean_txt(str(len(apps)))}", MARGIN_X, SLIDE_H_IN - 0.55, SLIDE_W_IN - 2 * MARGIN_X, 0.35, 10, TEXT_MUTED, False, 1)
    return slide


def slide_detail(prs, brief, app):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BASE_DARK)

    # Safe reads
    aid = clean_txt(_safe_get(app, ["id"], ""))
    name = clean_txt(_safe_get(app, ["name"], "Application"))
    domain = clean_txt(_safe_get(app, ["domain"], ""))
    cat = clean_txt(_safe_get(app, ["category"], ""))
    icon = clean_txt(_safe_get(app, ["icon"], ""))
    accent = clean_txt(_safe_get(app, ["accent"], ACCENT)) or ACCENT
    tagline = clean_txt(_safe_get(app, ["tagline"], ""))
    overview = clean_txt(_safe_get(app, ["overview"], ""))
    impact = clean_txt(_safe_get(app, ["impact"], ""))

    features = _safe_get(app, ["features"], []) or []
    metrics = _safe_get(app, ["metrics"], []) or []

    # Header strip
    add_txt(slide, "APPLICATION DETAIL", MARGIN_X, MARGIN_Y, SLIDE_W_IN - 2 * MARGIN_X, 0.30, 10, TEXT_MUTED, True, 1)
    add_rect(slide, MARGIN_X, MARGIN_Y + 0.32, SLIDE_W_IN - 2 * MARGIN_X, 0.02, LINE_SOFT, LINE_SOFT, 0.0, 0.0)

    # Left content (62%)
    left_x = CONTENT_X
    left_y = 1.15
    left_w = CONTENT_W
    left_h = SLIDE_H_IN - left_y - MARGIN_Y

    # App identity block
    add_pill(slide, left_x, left_y, f"ID {aid}" if aid else "ID —", "#0B1222", TEXT_PRIMARY, 9)
    add_pill(slide, left_x + 1.05, left_y, domain if domain else "Domain —", "#0B1222", TEXT_PRIMARY, 9)
    add_pill(slide, left_x + 2.55, left_y, cat if cat else "Category —", "#0B1222", TEXT_PRIMARY, 9)

    # Accent marker + icon
    add_rect(slide, left_x, left_y + 0.52, 0.12, 0.70, accent, accent, 0.0, 0.0)
    add_txt(slide, icon, left_x + 0.18, left_y + 0.50, 0.60, 0.60, 22, TEXT_PRIMARY, False, 1)

    body_sz = 12
    heading_sz = _fit_font(body_sz * GOLD, 18, 28)

    add_txt(slide, name, left_x + 0.82, left_y + 0.52, left_w - 0.82, 0.55, heading_sz, TEXT_PRIMARY, True, 1)

    if tagline:
        add_txt(slide, tagline, left_x + 0.82, left_y + 1.05, left_w - 0.82, 0.45, 12, TEXT_MUTED, False, 1)

    # Overview (chunked - Miller's Law)
    section_y = left_y + 1.60
    add_txt(slide, "Overview", left_x, section_y, left_w, 0.35, 12, TEXT_PRIMARY, True, 1)
    section_y += 0.38

    chunks = _chunk_text(overview, 175)[:4]  # keep to 4 chunks for breathing room
    if not chunks:
        chunks = ["—"]
    for ch in chunks:
        add_txt(slide, ch, left_x, section_y, left_w, 0.42, 11, TEXT_MUTED, False, 1)
        section_y += 0.42

    # Impact callout
    if impact:
        callout_h = 0.72
        add_rect(slide, left_x, section_y + 0.12, left_w, callout_h, "#0B1222", LINE_SOFT, 1.0, 0.16)
        add_rect(slide, left_x, section_y + 0.12, 0.10, callout_h, accent, accent, 0.0, 0.0)
        add_txt(slide, "Impact", left_x + 0.22, section_y + 0.18, left_w - 0.30, 0.26, 10, TEXT_PRIMARY, True, 1)
        imp = impact if len(impact) <= 170 else impact[:167].rstrip() + "…"
        add_txt(slide, imp, left_x + 0.22, section_y + 0.40, left_w - 0.30, 0.40, 10.5, TEXT_MUTED, False, 1)
        section_y += callout_h + 0.28
    else:
        section_y += 0.20

    # Features (5 items)
    add_txt(slide, "Key Features", left_x, section_y, left_w, 0.35, 12, TEXT_PRIMARY, True, 1)
    section_y += 0.40

    feat_card_gap = 0.16
    feat_card_h = 0.62
    max_feats = 5
    for i in range(min(max_feats, len(features))):
        f = features[i] or {}
        ft = clean_txt(_safe_get(f, ["title"], "Feature"))
        fd = clean_txt(_safe_get(f, ["description"], ""))
        y = section_y + i * (feat_card_h + feat_card_gap)

        add_rect(slide, left_x, y, left_w, feat_card_h, SURFACE, LINE_SOFT, 1.0, 0.14)
        add_rect(slide, left_x, y, 0.08, feat_card_h, accent, accent, 0.0, 0.0)
        add_txt(slide, ft, left_x + 0.18, y + 0.12, left_w - 0.30, 0.25, 11, TEXT_PRIMARY, True, 1)
        if fd:
            desc = fd if len(fd) <= 150 else fd[:147].rstrip() + "…"
            add_txt(slide, desc, left_x + 0.18, y + 0.33, left_w - 0.30, 0.26, 10, TEXT_MUTED, False, 1)

    # Right sidebar (38%)
    sb_x = SIDEBAR_X
    sb_y = 1.15
    sb_w = SIDEBAR_W - MARGIN_X
    sb_h = SLIDE_H_IN - sb_y - MARGIN_Y

    # Technical specs container
    add_rect(slide, sb_x, sb_y, sb_w, sb_h, SURFACE, LINE_SOFT, 1.0, 0.16)
    add_txt(slide, "Specs & Metrics", sb_x + 0.30, sb_y + 0.22, sb_w - 0.60, 0.35, 14, TEXT_PRIMARY, True, 1)

    # Domain pills section
    add_txt(slide, "Classification", sb_x + 0.30, sb_y + 0.65, sb_w - 0.60, 0.30, 11, TEXT_MUTED, True, 1)
    py = sb_y + 0.95
    add_pill(slide, sb_x + 0.30, py, domain if domain else "Domain —", "#0B1222", TEXT_PRIMARY, 9)
    py += 0.44
    add_pill(slide, sb_x + 0.30, py, cat if cat else "Category —", "#0B1222", TEXT_PRIMARY, 9)
    py += 0.54

    # Metrics list
    add_txt(slide, "Metrics", sb_x + 0.30, py, sb_w - 0.60, 0.30, 11, TEXT_MUTED, True, 1)
    py += 0.36

    if not metrics:
        metrics = [{"label": "Coverage", "value": "—"}, {"label": "Status", "value": "—"}, {"label": "Version", "value": "—"}]

    row_h = 0.54
    max_rows = int((sb_y + sb_h - py - 0.40) // (row_h + 0.12))
    show = metrics[:max_rows]

    for m in show:
        lbl = clean_txt(_safe_get(m, ["label"], "—"))
        val = clean_txt(_safe_get(m, ["value"], "—"))
        add_rect(slide, sb_x + 0.30, py, sb_w - 0.60, row_h, "#0B1222", LINE_SOFT, 1.0, 0.14)
        add_txt(slide, lbl, sb_x + 0.46, py + 0.14, (sb_w - 0.60) * 0.62, 0.26, 10, TEXT_MUTED, True, 1)
        add_txt(slide, val, sb_x + 0.46 + (sb_w - 0.60) * 0.62, py + 0.11, (sb_w - 0.60) * 0.38 - 0.16, 0.30, 11, TEXT_PRIMARY, True, 3)
        py += row_h + 0.12

    # Footer
    port = clean_txt(_safe_get(brief, ["portfolio", "organisation"], "Systech Analytics"))
    add_txt(slide, f"{port} · {name}", MARGIN_X, SLIDE_H_IN - 0.55, SLIDE_W_IN - 2 * MARGIN_X, 0.35, 10, TEXT_MUTED, False, 1)

    return slide


# =========================
# Main
# =========================
def main():
    parser = argparse.ArgumentParser(description="IntelliFrame Wireframe Deck Generator (v2.0 - PPTX Master)")
    parser.add_argument("--output", required=True, help="Output PPTX file path")
    args = parser.parse_args()

    # Brief (STRICT JSON) embedded as Python dict
    brief = {
        "portfolio": {
            "title": "Systech Analytics Application Portfolio",
            "organisation": "Systech Analytics",
            "year": "2026",
            "confidential": True,
            "stats": {
                "totalApps": "15 Applications",
                "domains": "5+ Industries",
                "aiPowered": "7+ AI / LLM-Powered",
                "buildTenure": "Not specified",
            },
        },
        "applications": [
            {
                "id": "01",
                "name": "Maverick",
                "domain": "Casino",
                "category": "AI Generation",
                "tagline": "AI-Powered Promotional Coupon Engine for Casino Players",
                "overview": "Maverick is a personalised incentive generation platform built for the gaming and casino industry. It analyses player behaviour, spending history, game preferences, and visit frequency to craft hyper-personalised promotional offers — from free spins and dining credits to VIP room upgrades. Powered by AI, Maverick replaces static promotions with dynamic, data-driven campaigns that adapt in real time to maximise engagement and lifetime value.",
                "features": [
                    {"title": "Player segmentation", "description": "Segments players by lifetime value, churn risk, and play patterns."},
                    {"title": "Automated coupon generation", "description": "Generates offers with configurable reward rules and expiry logic."},
                    {"title": "Real-time offer delivery", "description": "Delivers offers via SMS, app notifications, and kiosk displays."},
                    {"title": "Campaign analytics", "description": "Tracks performance with A/B testing support."},
                    {"title": "Systems integration", "description": "Integrates with casino POS and loyalty management systems."},
                ],
                "impact": "Increases promotional ROI and player retention by delivering the right offer to the right player at the right moment, replacing guesswork with precision targeting.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#F59E0B",
                "icon": "🎰",
            },
            {
                "id": "02",
                "name": "TRBank",
                "domain": "Banking",
                "category": "Data Processing",
                "tagline": "Unstructured-to-Structured Data Transformation for Banking",
                "overview": "TRBank is an intelligent document processing platform designed for financial institutions. It ingests unstructured banking documents — statements, loan applications, KYC forms, remittance slips, and correspondence — and transforms them into clean, validated, structured datasets ready for analytics, compliance reporting, and core banking ingestion.",
                "features": [
                    {"title": "Multi-format ingestion", "description": "Supports PDFs, scanned images, and handwritten forms."},
                    {"title": "Entity extraction", "description": "Extracts account numbers, dates, amounts, and counterparties."},
                    {"title": "Validation rules engine", "description": "Flags anomalies, duplicates, and missing fields."},
                    {"title": "Structured outputs", "description": "Exports to JSON, CSV, and direct database write formats."},
                    {"title": "Auditability", "description": "Provides audit trail and confidence scoring per extracted field."},
                ],
                "impact": "Eliminates manual data entry, reducing processing time from hours to minutes per batch while improving accuracy and compliance readiness.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#2563EB",
                "icon": "🏦",
            },
            {
                "id": "03",
                "name": "Sustainability",
                "domain": "ESG",
                "category": "Emission Analytics",
                "tagline": "Full-Cycle ESG Intelligence and Emission Calculation Platform",
                "overview": "The Sustainability platform is an end-to-end ESG intelligence tool for organisations managing environmental footprint. It ingests utility bills, fuel logs, travel records, and supply chain data to calculate carbon emissions across Scope 1, 2, and 3. It generates audit-ready ESG reports aligned with GHG Protocol and GRI standards and provides predictive analytics to model emission reduction scenarios for net-zero planning.",
                "features": [
                    {"title": "Utility bill parsing", "description": "Automates bill parsing with intelligent emission factor mapping."},
                    {"title": "Scope 1–3 carbon accounting", "description": "Calculates emissions aligned with the GHG Protocol."},
                    {"title": "AI ESG narratives", "description": "Generates board- and regulator-ready narrative reports."},
                    {"title": "Benchmarking dashboards", "description": "Provides peer benchmarking and target-gap analysis."},
                    {"title": "Scenario modelling", "description": "Models net-zero pathways and reduction roadmap planning."},
                ],
                "impact": "Reduces ESG reporting effort from weeks of manual spreadsheets to automated, audit-ready output generated in minutes.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#16A34A",
                "icon": "🌿",
            },
            {
                "id": "04",
                "name": "Ecovision",
                "domain": "ESG",
                "category": "Vision AI",
                "tagline": "Video-Powered AI Auditing for Sustainability Compliance",
                "overview": "Ecovision brings computer vision into sustainability compliance by transforming CCTV into an active monitoring system. It analyses live and recorded facility video to detect energy waste, improper waste disposal, safety non-compliance, and resource misuse. Each event is logged with a video clip, timestamp, severity rating, and an AI-generated corrective action recommendation.",
                "features": [
                    {"title": "Real-time video analysis", "description": "Detects sustainability and compliance events across facilities."},
                    {"title": "Incident logging", "description": "Extracts clips and classifies severity for each detected event."},
                    {"title": "Corrective actions", "description": "Generates AI recommendations to address detected incidents."},
                    {"title": "Compliance scoring", "description": "Tracks scores with historical trend analytics."},
                    {"title": "ESG integration", "description": "Integrates natively with the Sustainability platform for unified reporting."},
                ],
                "impact": "Turns periodic walkthrough audits into continuous, automated, evidence-driven compliance with an always-on audit trail.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#16A34A",
                "icon": "🌿",
            },
            {
                "id": "05",
                "name": "Hotel Concierge",
                "domain": "Hospitality",
                "category": "AI Agent",
                "tagline": "AI Avatar Booking Agent for Hotels and Guest Amenities",
                "overview": "Hotel Concierge is a conversational AI agent that handles end-to-end hotel booking and guest services via a lifelike AI avatar. Guests can search rooms, check availability, book, and reserve amenities like dining and spa. The agent maintains multi-turn context and orchestrates reservations through hotel PMS integrations for always-available self-service.",
                "features": [
                    {"title": "End-to-end booking", "description": "Supports room search, booking, modification, and cancellation."},
                    {"title": "Amenity reservations", "description": "Books dining, spa, pool, gym, and activities."},
                    {"title": "Conversational memory", "description": "Maintains multi-turn context with guest preference tracking."},
                    {"title": "AI avatar interface", "description": "Provides an immersive, personalised avatar-based experience."},
                    {"title": "PMS integration", "description": "Connects to PMS and channel managers for live availability and pricing."},
                ],
                "impact": "Reduces front-desk workload and call volumes while increasing ancillary revenue via 24/7 self-service booking.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#EC4899",
                "icon": "🏨",
            },
            {
                "id": "06",
                "name": "VetAI",
                "domain": "HR Tech",
                "category": "AI Screening",
                "tagline": "End-to-End AI-Powered Candidate Screening and Interview Platform",
                "overview": "VetAI is a full-stack AI recruitment platform that automates screening from job posting to hiring decision. It conducts live AI video interviews, evaluates technical, behavioural, and communication competencies in real time, and produces structured hiring reports. Proctoring and silence detection support integrity, while role templates enable rapid setup across job types.",
                "features": [
                    {"title": "AI-led video interviews", "description": "Runs live interviews with dynamic, contextual follow-up questions."},
                    {"title": "Proctoring & integrity", "description": "Detects anomalies and generates interview integrity scoring."},
                    {"title": "Multi-dimensional evaluation", "description": "Assesses technical, behavioural, and communication performance."},
                    {"title": "Structured hiring reports", "description": "Generates role-fit scores, rankings, and decision-ready summaries."},
                    {"title": "Pipeline management", "description": "Provides candidate tracking dashboards with analytics."},
                ],
                "impact": "Cuts time-to-screen by up to 70%, letting recruiters focus on final-stage evaluation while AI scales initial assessment consistently.",
                "metrics": [
                    {"label": "Coverage", "value": "Full"},
                    {"label": "Status", "value": "Live"},
                    {"label": "Version", "value": "v1.0"},
                    {"label": "Time-to-screen reduction", "value": "Up to 70%"},
                ],
                "accent": "#8B5CF6",
                "icon": "🧑‍💼",
            },
            {
                "id": "07",
                "name": "Resonance",
                "domain": "L&D",
                "category": "Communication AI",
                "tagline": "AI Communication Coaching for Professional Upskilling",
                "overview": "Resonance is an AI communication coaching platform that helps professionals improve confidence and effectiveness. Using speech analysis and LLM-based evaluation, it measures tone, clarity, pacing, filler words, and persuasive impact across practice sessions and real conversations. It delivers personalised feedback and exercises with progress tracked over time.",
                "features": [
                    {"title": "Speech analytics", "description": "Analyses tone, pace, clarity, and filler word frequency."},
                    {"title": "Actionable feedback", "description": "Produces AI feedback reports with specific improvement suggestions."},
                    {"title": "Scenario practice", "description": "Supports presentations, negotiations, and interview simulations."},
                    {"title": "Progress tracking", "description": "Visualises skill growth across sessions over time."},
                    {"title": "Team analytics", "description": "Gives L&D managers cohort-level monitoring and insights."},
                ],
                "impact": "Improves organisational communication quality, reducing meeting inefficiencies and strengthening stakeholder engagement.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#06B6D4",
                "icon": "🎓",
            },
            {
                "id": "08",
                "name": "Orbit",
                "domain": "Internal Tool",
                "category": "Project Management",
                "tagline": "Systech's Internal Project and Ticket Tracking System",
                "overview": "Orbit is Systech's internal project management and ticket tracking platform. It centralises task creation, assignment, status tracking, and milestone management across projects and client engagements. Built to replace spreadsheets and email threads, it provides leadership a real-time view of delivery health, resourcing, and sprint progress across the organisation.",
                "features": [
                    {"title": "Projects & sprints", "description": "Creates projects/sprints with milestones, deadlines, and priorities."},
                    {"title": "Ticket lifecycle", "description": "Manages tickets from backlog to completion sign-off."},
                    {"title": "Workload visibility", "description": "Assigns owners and shows individual workload across projects."},
                    {"title": "Search & filters", "description": "Filters by priority/status with global search across work items."},
                    {"title": "Collaboration threads", "description": "Per-ticket activity feeds and comment threads for async collaboration."},
                ],
                "impact": "Creates a single source of truth for delivery, replacing fragmented communication with a structured operational hub.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#64748B",
                "icon": "🛠",
            },
            {
                "id": "09",
                "name": "SysRank",
                "domain": "Internal Tool",
                "category": "Technical Assessment",
                "tagline": "Systech's Proprietary Technical Assessment and Benchmarking Platform",
                "overview": "SysRank is Systech's internal equivalent of HackerRank for evaluating technical talent via coding challenges, SQL assessments, data engineering problems, and timed sets. Used for both hiring screening and internal benchmarking, it includes managed question banks by domain/difficulty/track and dashboards for results and skill heatmaps.",
                "features": [
                    {"title": "Timed challenges", "description": "Runs assessments across Python, SQL, and data engineering tracks."},
                    {"title": "Auto-grading", "description": "Auto-graded test cases with partial scoring and detailed breakdowns."},
                    {"title": "Candidate portal", "description": "Candidate-facing portal with proctoring and integrity controls."},
                    {"title": "Question bank management", "description": "Organises questions by domain, difficulty, and technology track."},
                    {"title": "Results analytics", "description": "Leaderboards, score history, and skill heatmaps for benchmarking."},
                ],
                "impact": "Standardises objective talent evaluation across hiring and internal capability benchmarking.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#64748B",
                "icon": "🛠",
            },
            {
                "id": "10",
                "name": "AeroIntel",
                "domain": "Aviation",
                "category": "RAG · Conversational AI",
                "tagline": "RAG-Powered AI Assistant for Airport CCR Technicians",
                "overview": "AeroIntel is a retrieval-augmented generation (RAG) assistant for airport Common Communication Room (CCR) technicians. It enables natural language queries over a hybrid knowledge base of structured operational databases and unstructured technical documents to surface procedures, manuals, fault histories, and protocols quickly during live incidents.",
                "features": [
                    {"title": "Hybrid RAG search", "description": "Combines structured database querying with document retrieval."},
                    {"title": "Conversational interface", "description": "Context-aware, multi-turn natural language interactions."},
                    {"title": "Guided procedures", "description": "Retrieves step-by-step maintenance procedures with resolution support."},
                    {"title": "Fault history matching", "description": "Finds similar incidents via fault history and pattern matching."},
                    {"title": "Agentic orchestration", "description": "Synthesises answers across multiple sources for complex queries."},
                ],
                "impact": "Reduces time-to-resolution by giving CCR engineers instant access to the right information during incidents.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#0EA5E9",
                "icon": "✈",
            },
            {
                "id": "11",
                "name": "SysMart",
                "domain": "Retail",
                "category": "AI Chatbot",
                "tagline": "Databricks-Powered AI Chatbot for Retail Operations",
                "overview": "SysMart is a conversational AI platform for retail operations built on Databricks AI serving endpoints. It combines live SQL database querying with RAG over product catalogues, policies, and FAQs to answer customer and internal questions. Freshchat integration enables deployment in existing support channels with escalation to human agents when needed.",
                "features": [
                    {"title": "Live database querying", "description": "Answers questions via natural language queries against inventory and order systems."},
                    {"title": "RAG over knowledge bases", "description": "Retrieves from catalogues, return policies, and operational FAQs."},
                    {"title": "Freshchat deployment", "description": "Integrates into Freshchat for customer-facing self-service."},
                    {"title": "Human handoff", "description": "Escalates complex queries to human agents with context."},
                    {"title": "Usage analytics", "description": "Tracks usage and categorises queries for continuous improvement."},
                ],
                "impact": "Reduces support ticket volume by automating routine queries, freeing agents for higher-value interactions.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#F97316",
                "icon": "🛒",
            },
            {
                "id": "12",
                "name": "Intelliframe",
                "domain": "Design",
                "category": "Generative AI",
                "tagline": "AI Wireframe Generator for Dashboards and Data Products",
                "overview": "Intelliframe generates annotated dashboard and application wireframes from natural language briefs or data schemas. Product teams describe metrics, user journeys, and layout preferences; Intelliframe produces structured mockups with component annotations, layout logic, and data-binding suggestions for developer handoff, with iterative refinement via follow-up prompts.",
                "features": [
                    {"title": "NL-to-wireframe generation", "description": "Creates wireframes for dashboards, apps, and portals from briefs."},
                    {"title": "Schema-aware layouts", "description": "Suggests layouts based on the underlying data model."},
                    {"title": "Annotated components", "description": "Adds UX rationale and interaction behaviour notes."},
                    {"title": "Export formats", "description": "Exports to image, PDF, and developer-ready specification documents."},
                    {"title": "Iterative refinement", "description": "Improves results through conversational follow-up prompts."},
                ],
                "impact": "Shrinks wireframing and design iteration from weeks to hours, reducing rework and misalignment across teams.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#A855F7",
                "icon": "🎨",
            },
            {
                "id": "13",
                "name": "CCTV (AiCCTV)",
                "domain": "Security",
                "category": "Vision Analytics",
                "tagline": "AI-Powered Surveillance and Vision Analytics Platform",
                "overview": "AiCCTV transforms passive camera infrastructure into an active security intelligence layer. It processes live and recorded feeds to detect anomalies, unauthorised access, crowd density violations, safety hazards, and behavioural patterns of concern. It generates real-time alerts with video clip evidence and incident metadata for proactive response and auditability.",
                "features": [
                    {"title": "Anomaly & intrusion detection", "description": "Detects threats across multiple simultaneous camera feeds."},
                    {"title": "Crowd & zone monitoring", "description": "Monitors density and restricted access by zones."},
                    {"title": "Safety hazard detection", "description": "Detects falls, unattended objects, and fire indicators."},
                    {"title": "Evidence-backed alerts", "description": "Dispatches alerts with video clips and rich incident metadata."},
                    {"title": "Incident analytics", "description": "Heatmaps, historical analytics, and compliance reporting."},
                ],
                "impact": "Enables proactive security operations at scale, reducing response time and providing evidence-backed records for investigations and audits.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#EF4444",
                "icon": "📷",
            },
            {
                "id": "14",
                "name": "DataOne",
                "domain": "Data Engineering",
                "category": "AI Tooling",
                "tagline": "Unified AI Data Engineering Toolkit Across Fabric, Snowflake and Databricks",
                "overview": "DataOne is a unified AI data engineering platform spanning Microsoft Fabric, Snowflake, and Databricks through a single interface. Using purpose-built MCP (Model Context Protocol) servers, it exposes platform capabilities to AI agents and automation workflows for natural language pipeline creation, schema exploration, query execution, and data asset management without switching consoles or writing boilerplate integration code.",
                "features": [
                    {"title": "MCP servers", "description": "Purpose-built MCP servers for Fabric, Snowflake, and Databricks."},
                    {"title": "Natural language pipelines", "description": "Creates pipelines and orchestrates transformations via language."},
                    {"title": "Cross-platform discovery", "description": "Explores schemas and queries data lineage across platforms."},
                    {"title": "AI agent operations", "description": "Automates data quality checks, monitoring, and alerting with agents."},
                    {"title": "Workflow automation", "description": "Integrates with no-code orchestration for pipeline workflows."},
                ],
                "impact": "Reduces data engineering toil by unifying three major platforms under a language-driven interface for agents and analysts.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#1D4ED8",
                "icon": "🧱",
            },
            {
                "id": "15",
                "name": "Chef",
                "domain": "Food and Beverage",
                "category": "AI Avatar",
                "tagline": "AI-Powered Video Avatar for Food and Beverage Experiences",
                "overview": "Chef is an AI video avatar application for food and beverage experiences. A lifelike AI presenter delivers personalised menu recommendations, ingredient explanations, allergen guidance, and step-by-step cooking walkthroughs via natural language conversation. It blends engaging content with commerce, enabling brands to offer an on-demand AI host with human-like warmth at scale.",
                "features": [
                    {"title": "Lifelike video avatar", "description": "Natural, context-aware conversational interaction through a video presenter."},
                    {"title": "Personalised recommendations", "description": "Suggests dishes based on preferences and dietary restrictions."},
                    {"title": "Guided cooking", "description": "Step-by-step walkthroughs with real-time Q&A."},
                    {"title": "Allergen & nutrition guidance", "description": "Provides allergen info, nutrition breakdowns, and substitutions."},
                    {"title": "Commerce integration", "description": "Integrates with POS and e-commerce for in-conversation ordering."},
                ],
                "impact": "Creates a differentiated always-on engagement channel for F&B brands, combining human-like warmth with scalable AI personalisation.",
                "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
                "accent": "#F43F5E",
                "icon": "🍽",
            },
        ],
    }

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    # Slides
    slide_hero(prs, brief)
    slide_catalog(prs, brief)

    apps = _safe_get(brief, ["applications"], []) or []
    for app in apps:
        slide_detail(prs, brief, app or {})

    prs.save(args.output)


if __name__ == "__main__":
    main()