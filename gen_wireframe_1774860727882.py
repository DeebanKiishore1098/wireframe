import argparse
import re
from math import ceil

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ──────────────────────────────────────────────────────────────────────────────
# Renderer Utilities (MUST DEFINE)
# ──────────────────────────────────────────────────────────────────────────────

def clean_txt(t):
    if t is None:
        return ""
    t = str(t)
    return re.sub(r'<[^>]*>', '', t).strip()


def _hex_to_rgb(hex_color):
    h = clean_txt(hex_color).lstrip("#")
    if len(h) != 6:
        h = "FFFFFF"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, hex):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex)


def add_txt(slide, text, x, y, w, h, size, color_hex, bold=False, align=0):
    text = clean_txt(text)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)
    tf.vertical_anchor = MSO_ANCHOR.TOP

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align

    r = p.runs[0]
    r.font.name = "Helvetica"
    r.font.size = Pt(size)
    r.font.bold = bool(bold)
    r.font.color.rgb = _hex_to_rgb(color_hex)
    return tb


def add_rect(slide, x, y, w, h, fill_hex, line_hex, line_w=1.0, r=0.1):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.adjustments[0] = float(r)
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    shp.line.color.rgb = _hex_to_rgb(line_hex)
    shp.line.width = Pt(max(0.0, float(line_w)))
    return shp


def add_pill(slide, x, y, text, bg_hex, fg_hex, size=10):
    text = clean_txt(text)
    pad_x = 0.16
    base_w = 0.35 + (len(text) * 0.085)
    base_h = 0.32
    shp = add_rect(slide, x, y, base_w + pad_x, base_h, bg_hex, bg_hex, 0.0, r=0.5)
    tb = slide.shapes.add_textbox(Inches(x + 0.10), Inches(y + 0.02), Inches(base_w + pad_x - 0.20), Inches(base_h - 0.04))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Helvetica"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = _hex_to_rgb(fg_hex)
    return shp


# ──────────────────────────────────────────────────────────────────────────────
# Design System Constants
# ──────────────────────────────────────────────────────────────────────────────

BASE_BG = "#05080F"
BASE = "#080C14"       # 60%
SURFACE = "#0F172A"    # 30%
ACCENT = "#2563EB"     # 10% (no purple)

TEXT_PRIMARY = "#E6EEFf"
TEXT_MUTED = "#9AA6BF"
STROKE = "#1E293B"

GOLDEN = 1.618

# Slide size (PowerPoint default widescreen 13.333 x 7.5 in)
SLIDE_W = 13.333
SLIDE_H = 7.5

# Strict 8-pt grid mapped to inches: 8pt = 0.111..."
GRID = 0.125  # close to 8–9pt; we keep consistent spacing in 0.125 increments
GUTTER = 0.50  # generous whitespace ~32pt
MARGIN = 0.75  # outer margin ~48-54pt

CONTENT_RATIO = 0.62
SIDEBAR_RATIO = 0.38

CONTENT_W = (SLIDE_W - 2 * MARGIN) * CONTENT_RATIO
SIDEBAR_W = (SLIDE_W - 2 * MARGIN) * SIDEBAR_RATIO
CONTENT_X = MARGIN
SIDEBAR_X = MARGIN + CONTENT_W

TOP_Y = MARGIN
BOTTOM_Y = SLIDE_H - MARGIN

LINE_H_MULT = 1.5


def _safe_get(d, path, default=""):
    cur = d
    try:
        for k in path:
            if isinstance(cur, dict):
                cur = cur.get(k, default)
            else:
                return default
        return cur if cur is not None else default
    except Exception:
        return default


def _coerce_list(x):
    if isinstance(x, list):
        return x
    if x is None:
        return []
    return [x]


def _clamp(n, lo, hi):
    return max(lo, min(hi, n))


def _accent_sanitized(hex_color):
    # Enforce "NO PURPLE" rule: if a provided accent is purple-ish, fall back to ACCENT.
    h = clean_txt(hex_color).lstrip("#").upper()
    if len(h) != 6:
        return ACCENT
    r = int(h[0:2], 16)
    g = int(h[2:4], 16)
    b = int(h[4:6], 16)
    # crude heuristic: purple-ish if R and B high and G low
    if (r > 120 and b > 120 and g < 110) or (b > r and r > 90 and g < 90):
        return ACCENT
    return "#" + h


def _chunk_text(text, max_chars=180):
    # Miller's Law-ish chunking: split into ~max_chars segments by sentences.
    text = clean_txt(text)
    if not text:
        return []
    parts = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    cur = ""
    for p in parts:
        p = p.strip()
        if not p:
            continue
        if len(cur) + len(p) + 1 <= max_chars:
            cur = (cur + " " + p).strip()
        else:
            if cur:
                chunks.append(cur)
            cur = p
    if cur:
        chunks.append(cur)
    return chunks


# ──────────────────────────────────────────────────────────────────────────────
# Slide Builders
# ──────────────────────────────────────────────────────────────────────────────

def build_hero(prs, brief):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BASE_BG)

    title = clean_txt(_safe_get(brief, ["portfolio", "title"], "Application Portfolio"))
    org = clean_txt(_safe_get(brief, ["portfolio", "organisation"], ""))
    year = clean_txt(_safe_get(brief, ["portfolio", "year"], ""))
    confidential = bool(_safe_get(brief, ["portfolio", "confidential"], False))

    # Header band
    add_rect(slide, MARGIN, TOP_Y, SLIDE_W - 2 * MARGIN, 0.72, BASE, STROKE, 1.0, r=0.18)
    add_txt(slide, f"{org}".strip() if org else " ", MARGIN + 0.30, TOP_Y + 0.10, CONTENT_W - 0.60, 0.30, 14, TEXT_MUTED, False, 0)
    right_hdr = f"{year}".strip()
    if confidential:
        right_hdr = (right_hdr + "  •  CONFIDENTIAL").strip("  •")
    add_txt(slide, right_hdr if right_hdr else " ", SIDEBAR_X + 0.20, TOP_Y + 0.10, SIDEBAR_W - 0.40, 0.30, 12, TEXT_MUTED, False, PP_ALIGN.RIGHT)

    # Main columns surfaces
    main_y = TOP_Y + 0.90
    main_h = SLIDE_H - MARGIN - main_y
    add_rect(slide, CONTENT_X, main_y, CONTENT_W, main_h, SURFACE, STROKE, 1.0, r=0.16)
    add_rect(slide, SIDEBAR_X + 0.10, main_y, SIDEBAR_W - 0.10, main_h, SURFACE, STROKE, 1.0, r=0.16)

    # Left: Title
    title_size = 48
    add_txt(slide, title, CONTENT_X + 0.50, main_y + 0.55, CONTENT_W - 1.00, 1.40, title_size, TEXT_PRIMARY, True, 0)

    subtitle = "Ultra-dark portfolio wireframes • Multi-slide overview"
    add_txt(slide, subtitle, CONTENT_X + 0.52, main_y + 2.05, CONTENT_W - 1.04, 0.45, 14, TEXT_MUTED, False, 0)

    # Portfolio stats (4 cards)
    stats = _safe_get(brief, ["portfolio", "stats"], {}) if isinstance(_safe_get(brief, ["portfolio", "stats"], {}), dict) else {}
    total_apps = clean_txt(stats.get("totalApps", str(len(_coerce_list(_safe_get(brief, ["applications"], [])))) or "—"))
    domains = clean_txt(stats.get("domains", "—"))
    ai_powered = clean_txt(stats.get("aiPowered", "—"))
    build_tenure = clean_txt(stats.get("buildTenure", "—"))

    cards = [
        ("Total Apps", total_apps),
        ("Domains", domains),
        ("AI-Powered", ai_powered),
        ("Build Tenure", build_tenure),
    ]

    card_x = SIDEBAR_X + 0.30
    card_w = SIDEBAR_W - 0.50
    card_h = (main_h - 0.90) / 4.0
    card_y = main_y + 0.35
    for i, (k, v) in enumerate(cards):
        y = card_y + i * card_h
        add_rect(slide, card_x, y, card_w, card_h - 0.18, BASE, STROKE, 1.0, r=0.18)
        add_txt(slide, k, card_x + 0.28, y + 0.18, card_w - 0.56, 0.28, 11, TEXT_MUTED, True, 0)
        add_txt(slide, v if v else "—", card_x + 0.28, y + 0.52, card_w - 0.56, card_h - 0.80, 18, TEXT_PRIMARY, True, 0)

        # Accent bar
        add_rect(slide, card_x + 0.22, y + (card_h - 0.28) - 0.18, 0.75, 0.08, ACCENT, ACCENT, 0.0, r=0.4)

    return slide


def build_catalog(prs, brief):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BASE_BG)

    apps = _coerce_list(_safe_get(brief, ["applications"], []))

    # Title row
    add_rect(slide, MARGIN, TOP_Y, SLIDE_W - 2 * MARGIN, 0.72, BASE, STROKE, 1.0, r=0.18)
    add_txt(slide, "Application Catalog", MARGIN + 0.30, TOP_Y + 0.16, 7.8, 0.40, 22, TEXT_PRIMARY, True, 0)
    add_txt(slide, f"{len(apps)} apps", SLIDE_W - MARGIN - 2.0, TOP_Y + 0.18, 2.0, 0.36, 12, TEXT_MUTED, False, PP_ALIGN.RIGHT)

    # Asymmetric layout surfaces
    main_y = TOP_Y + 0.90
    main_h = SLIDE_H - MARGIN - main_y

    add_rect(slide, CONTENT_X, main_y, CONTENT_W, main_h, SURFACE, STROKE, 1.0, r=0.16)
    add_rect(slide, SIDEBAR_X + 0.10, main_y, SIDEBAR_W - 0.10, main_h, SURFACE, STROKE, 1.0, r=0.16)

    # Right sidebar: legend / filters
    add_txt(slide, "Legend", SIDEBAR_X + 0.38, main_y + 0.40, SIDEBAR_W - 0.70, 0.35, 14, TEXT_PRIMARY, True, 0)
    add_txt(slide, "Cards show ID, name, domain and category. Accent uses corporate blue.",
            SIDEBAR_X + 0.38, main_y + 0.80, SIDEBAR_W - 0.70, 0.90, 11, TEXT_MUTED, False, 0)

    add_txt(slide, "Domains (sample)", SIDEBAR_X + 0.38, main_y + 1.85, SIDEBAR_W - 0.70, 0.30, 12, TEXT_PRIMARY, True, 0)

    # domain pills (unique up to 10)
    domains = []
    for a in apps:
        d = clean_txt(_safe_get(a, ["domain"], ""))
        if d and d not in domains:
            domains.append(d)
    domains = domains[:10]

    px = SIDEBAR_X + 0.38
    py = main_y + 2.25
    row_w = SIDEBAR_W - 0.85
    cur_x = px
    cur_y = py
    for d in domains:
        pill_w = 0.35 + (len(d) * 0.085) + 0.16
        if cur_x + pill_w > px + row_w:
            cur_x = px
            cur_y += 0.42
        add_pill(slide, cur_x, cur_y, d, BASE, TEXT_PRIMARY, 9)
        cur_x += pill_w + 0.14

    # Left: grid of app cards (no overlaps, 8pt-ish spacing)
    inner_x = CONTENT_X + 0.35
    inner_y = main_y + 0.35
    inner_w = CONTENT_W - 0.70
    inner_h = main_h - 0.70

    cols = 2
    gap = 0.25
    card_w = (inner_w - gap) / cols
    card_h = 1.28  # fixed height for consistent grid
    rows = max(1, int(inner_h // (card_h + gap)))
    capacity = rows * cols

    # If too many, increase rows by reducing card height slightly
    if len(apps) > capacity:
        # Compute required rows; then adjust card_h to fit
        req_rows = ceil(len(apps) / cols)
        card_h = (inner_h - (req_rows - 1) * gap) / req_rows
        card_h = _clamp(card_h, 1.00, 1.28)

    for idx, app in enumerate(apps):
        r = idx // cols
        c = idx % cols
        x = inner_x + c * (card_w + gap)
        y = inner_y + r * (card_h + gap)
        if y + card_h > inner_y + inner_h + 0.01:
            break

        app_id = clean_txt(_safe_get(app, ["id"], f"{idx+1:02d}"))
        name = clean_txt(_safe_get(app, ["name"], "Untitled"))
        domain = clean_txt(_safe_get(app, ["domain"], "—"))
        category = clean_txt(_safe_get(app, ["category"], "—"))
        tagline = clean_txt(_safe_get(app, ["tagline"], ""))

        add_rect(slide, x, y, card_w, card_h, BASE, STROKE, 1.0, r=0.16)

        # Accent strip (uniform corporate accent)
        add_rect(slide, x + 0.18, y + 0.18, 0.10, card_h - 0.36, ACCENT, ACCENT, 0.0, r=0.4)

        add_txt(slide, f"{app_id}  {name}", x + 0.36, y + 0.18, card_w - 0.54, 0.32, 13, TEXT_PRIMARY, True, 0)
        add_txt(slide, tagline if tagline else " ", x + 0.36, y + 0.52, card_w - 0.54, 0.45, 10.5, TEXT_MUTED, False, 0)

        # Pills: domain + category
        add_pill(slide, x + 0.36, y + card_h - 0.44, domain, SURFACE, TEXT_PRIMARY, 9)
        add_pill(slide, x + 0.36 + (0.35 + len(domain) * 0.085 + 0.16) + 0.14, y + card_h - 0.44, category, SURFACE, TEXT_PRIMARY, 9)

    return slide


def build_detail(prs, brief, app):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BASE_BG)

    # Top title row
    add_rect(slide, MARGIN, TOP_Y, SLIDE_W - 2 * MARGIN, 0.72, BASE, STROKE, 1.0, r=0.18)

    app_id = clean_txt(_safe_get(app, ["id"], "—"))
    name = clean_txt(_safe_get(app, ["name"], "Untitled"))
    domain = clean_txt(_safe_get(app, ["domain"], "—"))
    category = clean_txt(_safe_get(app, ["category"], "—"))
    tagline = clean_txt(_safe_get(app, ["tagline"], ""))
    icon = clean_txt(_safe_get(app, ["icon"], ""))
    accent = _accent_sanitized(_safe_get(app, ["accent"], ACCENT))

    add_txt(slide, f"{app_id}  {name}", MARGIN + 0.30, TOP_Y + 0.14, 8.5, 0.45, 22, TEXT_PRIMARY, True, 0)
    add_txt(slide, f"{domain} • {category}", SLIDE_W - MARGIN - 4.0, TOP_Y + 0.18, 4.0, 0.36, 12, TEXT_MUTED, False, PP_ALIGN.RIGHT)

    # Main panels
    main_y = TOP_Y + 0.90
    main_h = SLIDE_H - MARGIN - main_y
    add_rect(slide, CONTENT_X, main_y, CONTENT_W, main_h, SURFACE, STROKE, 1.0, r=0.16)
    add_rect(slide, SIDEBAR_X + 0.10, main_y, SIDEBAR_W - 0.10, main_h, SURFACE, STROKE, 1.0, r=0.16)

    # Left: identity block
    add_rect(slide, CONTENT_X + 0.35, main_y + 0.35, CONTENT_W - 0.70, 1.55, BASE, STROKE, 1.0, r=0.18)

    # Accent + icon
    add_rect(slide, CONTENT_X + 0.55, main_y + 0.55, 0.12, 1.15, accent, accent, 0.0, r=0.4)
    add_txt(slide, icon if icon else " ", CONTENT_X + 0.78, main_y + 0.50, 0.60, 0.60, 26, TEXT_PRIMARY, True, PP_ALIGN.CENTER)

    add_txt(slide, name, CONTENT_X + 1.45, main_y + 0.48, CONTENT_W - 2.00, 0.45, 20, TEXT_PRIMARY, True, 0)
    add_txt(slide, tagline if tagline else " ", CONTENT_X + 1.45, main_y + 0.95, CONTENT_W - 2.00, 0.60, 11.5, TEXT_MUTED, False, 0)

    # Pills row
    px = CONTENT_X + 1.45
    py = main_y + 1.35
    add_pill(slide, px, py, domain, SURFACE, TEXT_PRIMARY, 9)
    add_pill(slide, px + (0.35 + len(domain) * 0.085 + 0.16) + 0.14, py, category, SURFACE, TEXT_PRIMARY, 9)

    # Overview (chunked)
    overview = clean_txt(_safe_get(app, ["overview"], ""))
    chunks = _chunk_text(overview, max_chars=190)

    ov_y = main_y + 2.05
    ov_h = 2.05
    add_rect(slide, CONTENT_X + 0.35, ov_y, CONTENT_W - 0.70, ov_h, BASE, STROKE, 1.0, r=0.18)
    add_txt(slide, "Overview", CONTENT_X + 0.60, ov_y + 0.18, CONTENT_W - 1.20, 0.30, 13, TEXT_PRIMARY, True, 0)

    text_y = ov_y + 0.55
    line_h = 0.42
    for i in range(min(4, len(chunks))):
        add_txt(slide, f"• {chunks[i]}", CONTENT_X + 0.60, text_y + i * line_h, CONTENT_W - 1.20, 0.40, 11, TEXT_MUTED, False, 0)

    # Features
    feats = _coerce_list(_safe_get(app, ["features"], []))
    ft_y = ov_y + ov_h + 0.30
    ft_h = main_y + main_h - 0.35 - ft_y
    add_rect(slide, CONTENT_X + 0.35, ft_y, CONTENT_W - 0.70, ft_h, BASE, STROKE, 1.0, r=0.18)
    add_txt(slide, "Key Features", CONTENT_X + 0.60, ft_y + 0.18, CONTENT_W - 1.20, 0.30, 13, TEXT_PRIMARY, True, 0)

    # Feature list (up to 6)
    fy = ft_y + 0.58
    max_items = 6
    item_h = 0.62
    for i, f in enumerate(feats[:max_items]):
        title = clean_txt(_safe_get(f, ["title"], "Feature"))
        desc = clean_txt(_safe_get(f, ["description"], ""))
        box_y = fy + i * item_h
        if box_y + item_h > ft_y + ft_h - 0.10:
            break
        add_rect(slide, CONTENT_X + 0.60, box_y, CONTENT_W - 1.20, 0.54, SURFACE, STROKE, 1.0, r=0.14)
        add_txt(slide, title, CONTENT_X + 0.78, box_y + 0.08, (CONTENT_W - 1.20) * 0.42, 0.34, 11.5, TEXT_PRIMARY, True, 0)
        add_txt(slide, desc, CONTENT_X + 0.78 + (CONTENT_W - 1.20) * 0.42, box_y + 0.08,
                (CONTENT_W - 1.20) * 0.58 - 0.18, 0.34, 10.5, TEXT_MUTED, False, 0)

    # Right sidebar: specs + metrics + impact
    sb_x = SIDEBAR_X + 0.30
    sb_w = SIDEBAR_W - 0.50

    # Technical specs card (structured)
    add_rect(slide, sb_x, main_y + 0.35, sb_w, 1.45, BASE, STROKE, 1.0, r=0.18)
    add_txt(slide, "Technical Specs", sb_x + 0.26, main_y + 0.52, sb_w - 0.52, 0.30, 13, TEXT_PRIMARY, True, 0)

    add_txt(slide, "Surface", sb_x + 0.26, main_y + 0.86, sb_w * 0.40, 0.26, 10.5, TEXT_MUTED, False, 0)
    add_txt(slide, "Web App / Platform", sb_x + sb_w * 0.46, main_y + 0.86, sb_w * 0.50, 0.26, 10.5, TEXT_PRIMARY, False, 0)

    add_txt(slide, "Pattern", sb_x + 0.26, main_y + 1.12, sb_w * 0.40, 0.26, 10.5, TEXT_MUTED, False, 0)
    add_txt(slide, "Service + UI", sb_x + sb_w * 0.46, main_y + 1.12, sb_w * 0.50, 0.26, 10.5, TEXT_PRIMARY, False, 0)

    add_txt(slide, "Accent", sb_x + 0.26, main_y + 1.38, sb_w * 0.40, 0.26, 10.5, TEXT_MUTED, False, 0)
    # show sanitized accent chip (still respects "NO PURPLE" via sanitizer)
    add_rect(slide, sb_x + sb_w * 0.46, main_y + 1.40, 0.40, 0.18, accent, accent, 0.0, r=0.4)
    add_txt(slide, accent, sb_x + sb_w * 0.46 + 0.48, main_y + 1.34, sb_w * 0.50 - 0.48, 0.30, 10.5, TEXT_PRIMARY, False, 0)

    # Metrics
    metrics = _coerce_list(_safe_get(app, ["metrics"], []))
    met_y = main_y + 1.95
    met_h = 2.10
    add_rect(slide, sb_x, met_y, sb_w, met_h, BASE, STROKE, 1.0, r=0.18)
    add_txt(slide, "Metrics", sb_x + 0.26, met_y + 0.18, sb_w - 0.52, 0.30, 13, TEXT_PRIMARY, True, 0)

    my = met_y + 0.58
    row_h = 0.40
    for i, m in enumerate(metrics[:6]):
        lbl = clean_txt(_safe_get(m, ["label"], ""))
        val = clean_txt(_safe_get(m, ["value"], ""))
        y = my + i * row_h
        if y + row_h > met_y + met_h - 0.12:
            break
        add_txt(slide, lbl if lbl else "—", sb_x + 0.26, y, sb_w * 0.45, 0.30, 10.5, TEXT_MUTED, False, 0)
        add_txt(slide, val if val else "—", sb_x + sb_w * 0.52, y, sb_w * 0.46, 0.30, 10.5, TEXT_PRIMARY, True, 0)

    # Impact
    impact = clean_txt(_safe_get(app, ["impact"], ""))
    imp_y = met_y + met_h + 0.30
    imp_h = main_y + main_h - 0.35 - imp_y
    add_rect(slide, sb_x, imp_y, sb_w, imp_h, BASE, STROKE, 1.0, r=0.18)
    add_txt(slide, "Impact", sb_x + 0.26, imp_y + 0.18, sb_w - 0.52, 0.30, 13, TEXT_PRIMARY, True, 0)
    add_txt(slide, impact if impact else "—", sb_x + 0.26, imp_y + 0.56, sb_w - 0.52, imp_h - 0.70, 11, TEXT_MUTED, False, 0)

    # Footer page marker
    add_txt(slide, "Systech Analytics • Portfolio Detail", MARGIN, SLIDE_H - MARGIN + 0.10, 6.5, 0.30, 9.5, TEXT_MUTED, False, 0)

    return slide


# ──────────────────────────────────────────────────────────────────────────────
# Main
# ──────────────────────────────────────────────────────────────────────────────

BRIEF = {
  "portfolio": {
    "title": "Systech Analytics Application Portfolio",
    "organisation": "Systech Analytics",
    "year": "2026",
    "confidential": True,
    "stats": {
      "totalApps": "15",
      "domains": "AI, Data Engineering, Automation, Analytics",
      "aiPowered": "7+",
      "buildTenure": "Not specified"
    }
  },
  "applications": [
    {
      "id": "01",
      "name": "Maverick",
      "domain": "Casino",
      "category": "AI Generation",
      "tagline": "AI-Powered Promotional Coupon Engine for Casino Players",
      "overview": "Maverick is a personalised incentive generation platform built for the gaming and casino industry. It analyses player behaviour, spending history, game preferences, and visit frequency to craft hyper-personalised promotional offers — from free spins and dining credits to VIP room upgrades. Powered by AI, Maverick replaces static promotions with dynamic, data-driven campaigns that adapt in real time to maximise engagement and lifetime value.",
      "features": [
        {"title": "Player segmentation", "description": "Segments players by lifetime value, churn risk, and play patterns."},
        {"title": "Automated coupon generation", "description": "Generates coupons with configurable reward rules and expiry logic."},
        {"title": "Real-time offer delivery", "description": "Delivers offers via SMS, app notifications, and kiosk displays."},
        {"title": "Campaign analytics", "description": "Provides performance analytics with A/B testing support."},
        {"title": "Systems integration", "description": "Integrates with casino POS and loyalty management systems."}
      ],
      "impact": "Increases promotional ROI and player retention by delivering the right offer to the right player at the right moment, replacing guesswork with precision targeting.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#8E24AA",
      "icon": "🎰"
    },
    {
      "id": "02",
      "name": "TRBank",
      "domain": "Banking",
      "category": "Data Processing",
      "tagline": "Unstructured-to-Structured Data Transformation for Banking",
      "overview": "TRBank is an intelligent document processing platform designed for financial institutions. It ingests unstructured banking documents — account statements, loan applications, KYC forms, remittance slips, and free-text correspondence — and transforms them into clean, validated, structured datasets ready for downstream analytics, compliance reporting, and core banking system ingestion.",
      "features": [
        {"title": "Multi-format ingestion", "description": "Ingests PDFs, scanned images, and handwritten forms."},
        {"title": "Entity extraction", "description": "Extracts account numbers, dates, amounts, and counterparties."},
        {"title": "Validation rules engine", "description": "Flags anomalies, duplicates, and missing fields."},
        {"title": "Structured outputs", "description": "Outputs JSON, CSV, and supports direct database writes."},
        {"title": "Auditability", "description": "Provides audit trails and confidence scoring per extracted field."}
      ],
      "impact": "Eliminates manual data entry, reducing processing time from hours to minutes per batch while improving accuracy and compliance readiness.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#1565C0",
      "icon": "🏦"
    },
    {
      "id": "03",
      "name": "Sustainability",
      "domain": "ESG",
      "category": "Emission Analytics",
      "tagline": "Full-Cycle ESG Intelligence and Emission Calculation Platform",
      "overview": "The Sustainability platform is an end-to-end ESG intelligence tool for organisations managing their environmental footprint. It ingests utility bills, fuel logs, travel records, and supply chain data to calculate carbon emissions across Scope 1, 2, and 3. It generates audit-ready ESG reports aligned with GHG Protocol and GRI standards, and provides predictive analytics to model emission reduction scenarios for net-zero planning.",
      "features": [
        {"title": "Automated bill parsing", "description": "Parses utility bills and maps emissions factors intelligently."},
        {"title": "Carbon accounting", "description": "Calculates Scope 1, 2, and 3 aligned with the GHG Protocol."},
        {"title": "Narrative ESG reporting", "description": "Generates AI-written ESG narratives for board/regulatory submission."},
        {"title": "Benchmarking dashboards", "description": "Provides peer benchmarking and target-gap analysis."},
        {"title": "Scenario modelling", "description": "Models net-zero scenarios and reduction roadmaps."}
      ],
      "impact": "Reduces ESG reporting effort from weeks of spreadsheets to audit-ready output in minutes, enabling confident sustainability reporting.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#2E7D32",
      "icon": "🌿"
    },
    {
      "id": "04",
      "name": "Ecovision",
      "domain": "ESG",
      "category": "Vision AI",
      "tagline": "Video-Powered AI Auditing for Sustainability Compliance",
      "overview": "Ecovision transforms passive CCTV into an active sustainability compliance monitoring system. It analyses live and recorded facility video feeds to detect energy waste, improper waste disposal, safety non-compliance, and resource misuse. Each event is logged with an extracted clip, timestamp, severity rating, and an AI-generated corrective action recommendation.",
      "features": [
        {"title": "Real-time video analysis", "description": "Detects sustainability and compliance events across facilities."},
        {"title": "Incident logging", "description": "Logs incidents with clip extraction and severity classification."},
        {"title": "Corrective actions", "description": "Generates AI recommendations for corrective action per incident."},
        {"title": "Compliance scoring", "description": "Provides scoring dashboards with historical trend analytics."},
        {"title": "Platform integration", "description": "Integrates natively with the Sustainability platform for unified reporting."}
      ],
      "impact": "Turns periodic walkthroughs into continuous, automated, evidence-driven auditing with an always-on trail for regulators and internal teams.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#2E7D32",
      "icon": "🌿"
    },
    {
      "id": "05",
      "name": "Hotel Concierge",
      "domain": "Hospitality",
      "category": "AI Agent",
      "tagline": "AI Avatar Booking Agent for Hotels and Guest Amenities",
      "overview": "Hotel Concierge is a conversational AI agent that handles end-to-end hotel booking and guest services through a lifelike AI avatar. Guests can search rooms, check availability, book, modify/cancel stays, and reserve amenities like dining and spa. The agent maintains multi-turn context and orchestrates reservations via PMS integrations for real-time availability and pricing.",
      "features": [
        {"title": "End-to-end booking", "description": "Supports room search, booking, modification, and cancellation in natural language."},
        {"title": "Amenity reservations", "description": "Books dining, spa, pool, gym, and activities."},
        {"title": "Conversation memory", "description": "Maintains multi-turn context with guest preference tracking."},
        {"title": "AI avatar interface", "description": "Provides an immersive, personalised lifelike avatar experience."},
        {"title": "PMS integration", "description": "Integrates with PMS/channel managers for live availability and pricing."}
      ],
      "impact": "Reduces front-desk workload and call volumes while increasing ancillary revenue through 24/7 self-service booking via a human-like AI interface.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#00838F",
      "icon": "🏨"
    },
    {
      "id": "06",
      "name": "VetAI",
      "domain": "HR Tech",
      "category": "AI Screening",
      "tagline": "End-to-End AI-Powered Candidate Screening and Interview Platform",
      "overview": "VetAI automates candidate screening from job posting to hiring decision. It runs AI-powered live video interviews, evaluates responses across communication, technical accuracy, and behavioural competencies, and produces structured hiring reports. Video proctoring and silence detection support integrity, while role templates enable rapid setup across job types.",
      "features": [
        {"title": "AI-led interviews", "description": "Conducts live video interviews with contextual follow-up questioning."},
        {"title": "Proctoring & integrity", "description": "Detects anomalies and generates interview integrity scores."},
        {"title": "Multi-dimensional evaluation", "description": "Scores technical, behavioural, and communication competencies."},
        {"title": "Hiring reports", "description": "Produces structured role-fit scores and candidate rankings."},
        {"title": "Pipeline management", "description": "Provides a candidate dashboard with tracking and analytics."}
      ],
      "impact": "Cuts time-to-screen by up to 70%, letting recruiters focus on final evaluation while AI scales consistent initial assessment.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}, {"label": "Time-to-screen reduction", "value": "Up to 70%"}],
      "accent": "#6A1B9A",
      "icon": "🧑‍💼"
    },
    {
      "id": "07",
      "name": "Resonance",
      "domain": "L&D",
      "category": "Communication AI",
      "tagline": "AI Communication Coaching for Professional Upskilling",
      "overview": "Resonance is an AI communication coaching platform that uses speech analysis and LLM-based evaluation to assess tone, clarity, pacing, filler words, and persuasive impact. It delivers personalised feedback and structured exercises, tracking progress across sessions for individuals and cohorts with team-level analytics for L&D managers.",
      "features": [
        {"title": "Speech analytics", "description": "Analyses tone, pace, clarity, and filler word frequency in real time."},
        {"title": "Actionable feedback", "description": "Generates specific, actionable improvement reports."},
        {"title": "Scenario practice", "description": "Supports practice for presentations, negotiations, and interviews."},
        {"title": "Progress tracking", "description": "Tracks skill growth over time with visualisations."},
        {"title": "Cohort insights", "description": "Provides team-level analytics for L&D managers."}
      ],
      "impact": "Improves communication effectiveness across organisations, reducing meeting inefficiencies and strengthening stakeholder engagement.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#3949AB",
      "icon": "🗣"
    },
    {
      "id": "08",
      "name": "Orbit",
      "domain": "Internal Tool",
      "category": "Project Management",
      "tagline": "Systech's Internal Project and Ticket Tracking System",
      "overview": "Orbit is Systech's internal project management and ticket tracking platform. It centralises task creation, assignment, status tracking, and milestone management across projects and client engagements. Built to replace spreadsheets and email threads, it provides leadership real-time visibility into delivery health, resource allocation, and sprint progress.",
      "features": [
        {"title": "Project & sprint setup", "description": "Creates projects/sprints with milestones, deadlines, and priorities."},
        {"title": "Ticket lifecycle", "description": "Manages tickets from backlog to completion sign-off."},
        {"title": "Workload visibility", "description": "Shows assignments and individual workload across projects."},
        {"title": "Search & filters", "description": "Provides status/priority filtering and global search."},
        {"title": "Collaboration feed", "description": "Per-ticket activity feeds and comment threads for async collaboration."}
      ],
      "impact": "Provides a single source of truth for delivery, replacing fragmented communications with a structured, searchable operational hub.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#546E7A",
      "icon": "🛠"
    },
    {
      "id": "09",
      "name": "SysRank",
      "domain": "Internal Tool",
      "category": "Technical Assessment",
      "tagline": "Systech's Proprietary Technical Assessment and Benchmarking Platform",
      "overview": "SysRank is Systech's internal equivalent of HackerRank for evaluating technical talent through coding challenges, SQL assessments, and data engineering problem sets. It supports objective screening in hiring and internal benchmarking. Question banks are organised by domain, difficulty, and technology track, with dashboards for score history and skill heatmaps.",
      "features": [
        {"title": "Timed assessments", "description": "Runs timed coding challenges across Python, SQL, and data engineering."},
        {"title": "Auto-grading", "description": "Auto-grades with partial scoring and detailed breakdowns."},
        {"title": "Secure portal", "description": "Candidate-facing assessment portal with proctoring/integrity controls."},
        {"title": "Question bank management", "description": "Organises questions by domain, difficulty, and technology track."},
        {"title": "Performance insights", "description": "Leaderboards, score history, and skill heatmaps."}
      ],
      "impact": "Standardises objective talent evaluation for hiring and internal capability benchmarking across the organisation.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#546E7A",
      "icon": "🧪"
    },
    {
      "id": "10",
      "name": "AeroIntel",
      "domain": "Aviation",
      "category": "RAG · Conversational AI",
      "tagline": "RAG-Powered AI Assistant for Airport CCR Technicians",
      "overview": "AeroIntel is a retrieval-augmented generation (RAG) assistant for airport CCR technicians. It supports natural language queries over a hybrid knowledge base of structured operational databases and unstructured technical documents, surfacing procedures, manuals, fault histories, and protocols with context-aware multi-turn conversation and agentic orchestration.",
      "features": [
        {"title": "Hybrid RAG search", "description": "Combines structured database querying with document retrieval."},
        {"title": "Conversational interface", "description": "Context-aware, multi-turn natural language interaction."},
        {"title": "Guided maintenance procedures", "description": "Retrieves step-by-step procedures with resolution support."},
        {"title": "Fault history matching", "description": "Looks up fault histories and matches similar incidents."},
        {"title": "Agentic orchestration", "description": "Synthesises answers across multiple sources for complex queries."}
      ],
      "impact": "Reduces time-to-resolution during live incidents by giving technicians instant access to the right information when every minute counts.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#1E88E5",
      "icon": "✈"
    },
    {
      "id": "11",
      "name": "SysMart",
      "domain": "Retail",
      "category": "AI Chatbot",
      "tagline": "Databricks-Powered AI Chatbot for Retail Operations",
      "overview": "SysMart is a conversational AI platform for retail operations built on Databricks AI serving endpoints. It combines live SQL database querying with RAG over product catalogues, return policies, and FAQs to support customers and internal teams. Freshchat integration enables deployment into existing support channels with escalation to human agents.",
      "features": [
        {"title": "Live database Q&A", "description": "Answers natural language questions against inventory and order databases."},
        {"title": "RAG knowledge base", "description": "Retrieves from catalogues, return policies, and operational FAQs."},
        {"title": "Freshchat integration", "description": "Deploys directly into customer-facing support channels."},
        {"title": "Human handoff", "description": "Escalates complex queries to human agents seamlessly."},
        {"title": "Analytics dashboard", "description": "Usage analytics and query categorisation for continuous improvement."}
      ],
      "impact": "Reduces support ticket volume by automating routine queries, freeing agents for higher-value customer interactions.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#EF6C00",
      "icon": "🛒"
    },
    {
      "id": "12",
      "name": "Intelliframe",
      "domain": "Design",
      "category": "Generative AI",
      "tagline": "AI Wireframe Generator for Dashboards and Data Products",
      "overview": "Intelliframe generates annotated dashboard and application wireframes from natural language briefs or data schemas. Product managers and analysts describe required metrics, user journeys, and layout preferences, and the tool produces structured mockups with component annotations, layout logic, and data-binding suggestions, exportable for developer handoff and iterative refinement.",
      "features": [
        {"title": "NL-to-wireframe", "description": "Generates wireframes for dashboards, apps, and portals from prompts."},
        {"title": "Schema-aware layouts", "description": "Suggests layouts based on the underlying data model."},
        {"title": "Annotated components", "description": "Adds UX rationale and interaction behaviour notes."},
        {"title": "Export formats", "description": "Exports to image, PDF, and developer-ready specification documents."},
        {"title": "Iterative refinement", "description": "Refines wireframes through conversational follow-ups."}
      ],
      "impact": "Compresses weeks of wireframing into hours, aligning product/design/engineering faster and reducing rework and misalignment.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#D81B60",
      "icon": "🧩"
    },
    {
      "id": "13",
      "name": "CCTV (AiCCTV)",
      "domain": "Security",
      "category": "Vision Analytics",
      "tagline": "AI-Powered Surveillance and Vision Analytics Platform",
      "overview": "AiCCTV turns passive camera infrastructure into an active security intelligence layer. It analyses live and recorded feeds to detect anomalies, unauthorised access, crowd density violations, safety hazards, and behavioural patterns of concern. It generates real-time alerts with video clip evidence and incident metadata, plus historical analytics and compliance reporting.",
      "features": [
        {"title": "Anomaly & intrusion detection", "description": "Detects unauthorised access and anomalies across multiple camera feeds."},
        {"title": "Crowd & zone monitoring", "description": "Monitors crowd density and restricted zone violations."},
        {"title": "Safety hazard detection", "description": "Detects falls, unattended objects, and fire indicators."},
        {"title": "Evidence-backed alerts", "description": "Dispatches alerts with clip evidence and incident metadata."},
        {"title": "Historical analytics", "description": "Heatmaps, incident trends, and compliance reporting."}
      ],
      "impact": "Enables proactive security operations at scale, reducing response time and providing evidence-backed records for audits and investigations.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#C62828",
      "icon": "📷"
    },
    {
      "id": "14",
      "name": "DataOne",
      "domain": "Data Engineering",
      "category": "AI Tooling",
      "tagline": "Unified AI Data Engineering Toolkit Across Fabric, Snowflake and Databricks",
      "overview": "DataOne is a unified AI data engineering platform bridging Microsoft Fabric, Snowflake, and Databricks through purpose-built MCP (Model Context Protocol) servers. It enables natural language pipeline creation, cross-platform schema exploration, query execution, and data asset management, and connects AI agents and automation workflows without switching consoles or writing boilerplate integration code.",
      "features": [
        {"title": "MCP servers", "description": "Purpose-built MCP servers for Fabric, Snowflake, and Databricks."},
        {"title": "NL pipeline orchestration", "description": "Creates and orchestrates transformations via natural language."},
        {"title": "Cross-platform exploration", "description": "Explores schemas and queries data lineage across platforms."},
        {"title": "AI agent operations", "description": "Automates data quality checks, monitoring, and alerting with agents."},
        {"title": "Workflow automation", "description": "Integrates no-code automation for pipeline orchestration."}
      ],
      "impact": "Reduces data engineering toil by providing a single language-driven interface across three major data platforms, eliminating context switching and manual API work.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#0D47A1",
      "icon": "🧱"
    },
    {
      "id": "15",
      "name": "Chef",
      "domain": "Food and Beverage",
      "category": "AI Avatar",
      "tagline": "AI-Powered Video Avatar for Food and Beverage Experiences",
      "overview": "Chef is an AI video avatar for food and beverage experiences. A lifelike presenter delivers personalised menu recommendations, ingredient explanations, allergen guidance, and step-by-step cooking walkthroughs via natural language conversation. It blends engaging content with commerce for restaurants, food brands, and recipe platforms, with integrations for in-conversation ordering.",
      "features": [
        {"title": "Video avatar host", "description": "Lifelike AI presenter with natural, context-aware conversation."},
        {"title": "Personalised recommendations", "description": "Suggests dishes based on preferences and dietary restrictions."},
        {"title": "Guided cooking", "description": "Provides step-by-step walkthroughs with real-time Q&A."},
        {"title": "Allergen & nutrition guidance", "description": "Explains allergens, nutrition, and ingredient substitutions."},
        {"title": "Commerce integration", "description": "Integrates with POS/e-commerce for in-conversation ordering."}
      ],
      "impact": "Creates an always-on engagement channel for F&B brands, combining human-like warmth with scalable, consistent AI personalisation.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#6D4C41",
      "icon": "🍽"
    }
  ]
}


def main():
    parser = argparse.ArgumentParser(description="Generate Systech Analytics Application Portfolio wireframe deck.")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    brief = BRIEF if isinstance(BRIEF, dict) else {}

    # Slides: HERO, CATALOG, then one DETAIL per app
    build_hero(prs, brief)
    build_catalog(prs, brief)

    apps = _coerce_list(_safe_get(brief, ["applications"], []))
    for app in apps:
        if not isinstance(app, dict):
            continue
        build_detail(prs, brief, app)

    prs.save(args.output)


if __name__ == "__main__":
    main()