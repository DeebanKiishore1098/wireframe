import argparse
import re
from math import ceil

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# -----------------------------
# Renderer Utilities (REQUIRED)
# -----------------------------
def clean_txt(t):
    if t is None:
        return ""
    t = str(t)
    return re.sub(r"<[^>]*>", "", t).strip()


def _hex_to_rgb(hex_color):
    h = clean_txt(hex_color).lstrip("#")
    if len(h) != 6:
        h = "FFFFFF"
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def add_txt(slide, text, x, y, w, h, size, color_hex, bold=False, align=0):
    # POSitional args only as required by caller; keep signature as given.
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = clean_txt(text)
    p.alignment = align

    run = p.runs[0]
    run.font.name = "Helvetica"
    run.font.size = Pt(size)
    run.font.bold = bool(bold)
    run.font.color.rgb = _hex_to_rgb(color_hex)
    return tb


def add_rect(slide, x, y, w, h, fill_hex, line_hex, line_w=1.0, r=0.1):
    # Rounded surface cards.
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.adjustments[0] = float(r)

    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(fill_hex)

    shape.line.color.rgb = _hex_to_rgb(line_hex)
    shape.line.width = Pt(float(line_w))
    return shape


def add_pill(slide, x, y, text, bg_hex, fg_hex, size=10):
    t = clean_txt(text)
    # Simple width heuristic: characters * 0.11" + padding
    w = max(0.9, min(3.2, 0.55 + 0.11 * len(t)))
    h = 0.34
    pill = add_rect(slide, x, y, w, h, bg_hex, bg_hex, 1.0, r=0.9)
    add_txt(slide, t, x + 0.14, y + 0.06, w - 0.28, h - 0.12, size, fg_hex, True, 0)
    return pill


# -----------------------------
# Data (STRICT JSON as provided)
# -----------------------------
DATA = {
  "portfolio": {
    "title": "Systech Analytics · Application Portfolio",
    "organisation": "Systech Analytics",
    "year": "2026",
    "confidential": True,
    "stats": {
      "totalApps": "15 Applications",
      "domains": "5+ Industries",
      "aiPowered": "7+ AI / LLM-Powered",
      "buildTenure": "Not specified"
    }
  },
  "applications": [
    {
      "id": "01",
      "name": "Maverick",
      "domain": "Casino",
      "category": "AI Generation",
      "tagline": "AI-Powered Promotional Coupon Engine for Casino Players",
      "overview": "Maverick is a personalised incentive generation platform built for the gaming and casino industry. It analyses player behaviour, spending history, game preferences, and visit frequency to craft hyper-personalised promotional offers — from free spins and dining credits to VIP room upgrades. Powered by AI, Maverick replaces static, one-size-fits-all promotions with dynamic, data-driven campaigns that adapt in real time to maximise player engagement and lifetime value.",
      "features": [
        {"title": "Player Segmentation", "description": "Segments players by lifetime value, churn risk, and play patterns."},
        {"title": "Automated Coupon Generation", "description": "Generates coupons with configurable reward rules and expiry logic."},
        {"title": "Real-Time Offer Delivery", "description": "Delivers offers via SMS, app notifications, and kiosk displays."},
        {"title": "Campaign Analytics & A/B Testing", "description": "Tracks performance and supports A/B testing for optimisation."},
        {"title": "POS & Loyalty Integrations", "description": "Integrates with casino POS and loyalty management systems."}
      ],
      "impact": "Increases promotional ROI and player retention by delivering the right offer to the right player at the right moment, replacing guesswork with precision targeting.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#7C3AED",
      "icon": "🎰"
    },
    {
      "id": "02",
      "name": "TRBank",
      "domain": "Banking",
      "category": "Data Processing",
      "tagline": "Unstructured-to-Structured Data Transformation for Banking",
      "overview": "TRBank is an intelligent document processing platform designed for financial institutions. It ingests unstructured banking documents — account statements, loan applications, KYC forms, remittance slips, and free-text correspondence — and transforms them into clean, validated, structured datasets ready for downstream analytics, compliance reporting, and core banking system ingestion.",
      "features": [
        {"title": "Multi-Format Ingestion", "description": "Ingests PDFs, scanned images, and handwritten forms."},
        {"title": "Entity Extraction", "description": "Extracts account numbers, dates, transaction amounts, and counterparties."},
        {"title": "Validation Rules Engine", "description": "Flags anomalies, duplicates, and missing fields."},
        {"title": "Structured Outputs", "description": "Exports JSON/CSV and supports direct database writes."},
        {"title": "Auditability", "description": "Provides audit trail and confidence scoring per extracted field."}
      ],
      "impact": "Eliminates manual data entry for banking operations teams, reducing document processing time from hours to minutes per batch while improving accuracy and compliance readiness.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#2563EB",
      "icon": "🏦"
    },
    {
      "id": "03",
      "name": "Sustainability",
      "domain": "ESG",
      "category": "Emission Analytics",
      "tagline": "Full-Cycle ESG Intelligence and Emission Calculation Platform",
      "overview": "The Sustainability platform is an end-to-end ESG intelligence tool for organisations managing their environmental footprint. It ingests utility bills, fuel logs, travel records, and supply chain data to calculate carbon emissions across Scope 1, 2, and 3 categories. The platform generates audit-ready ESG reports aligned with GHG Protocol and GRI standards, and provides predictive analytics to model emission reduction scenarios for net-zero planning.",
      "features": [
        {"title": "Automated Bill Parsing", "description": "Parses utility bills and maps emission factors intelligently."},
        {"title": "Scope 1–3 Carbon Accounting", "description": "Calculates emissions aligned with the GHG Protocol."},
        {"title": "AI-Generated ESG Narratives", "description": "Generates board- and regulator-ready ESG narrative reports."},
        {"title": "Benchmarking & Dashboards", "description": "Provides peer benchmarking and target-gap analysis dashboards."},
        {"title": "Scenario Modelling", "description": "Models reduction roadmaps for net-zero planning."}
      ],
      "impact": "Reduces ESG reporting effort from weeks of manual spreadsheet work to automated, audit-ready output generated in minutes — enabling organisations to meet sustainability reporting obligations with confidence.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#16A34A",
      "icon": "🌿"
    },
    {
      "id": "04",
      "name": "Ecovision",
      "domain": "ESG",
      "category": "Vision AI",
      "tagline": "Video-Powered AI Auditing for Sustainability Compliance",
      "overview": "Ecovision brings computer vision into the sustainability space, transforming passive CCTV infrastructure into an active compliance monitoring system. It continuously analyses live and recorded facility video feeds to detect energy waste, improper waste disposal, safety non-compliance, and resource misuse. Each detected event is logged with an extracted video clip, timestamp, severity rating, and an AI-generated corrective action recommendation.",
      "features": [
        {"title": "Continuous Video Auditing", "description": "Analyses live and recorded facility feeds for sustainability/compliance events."},
        {"title": "Incident Logging", "description": "Logs events with clip extraction, timestamps, and severity classification."},
        {"title": "Corrective Recommendations", "description": "Generates AI-driven corrective action recommendations per incident."},
        {"title": "Compliance Scoring", "description": "Provides dashboards with historical trends and compliance scores."},
        {"title": "Unified ESG Reporting", "description": "Integrates natively with the Sustainability platform for reporting."}
      ],
      "impact": "Transforms sustainability auditing from periodic manual walkthroughs into a continuous, automated, evidence-driven compliance function — providing an always-on audit trail for regulators and internal teams.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#16A34A",
      "icon": "👁"
    },
    {
      "id": "05",
      "name": "Hotel Concierge",
      "domain": "Hospitality",
      "category": "AI Agent",
      "tagline": "AI Avatar Booking Agent for Hotels and Guest Amenities",
      "overview": "Hotel Concierge is a conversational AI agent that handles end-to-end hotel booking and guest services through a lifelike AI avatar interface. Guests interact in natural language to search rooms, check availability, make bookings, and reserve amenities including dining, spa, and leisure activities. The agent maintains conversational context across multi-turn interactions and orchestrates reservations through backend hotel PMS integrations, delivering a seamless, always-available self-service experience.",
      "features": [
        {"title": "End-to-End Booking", "description": "Supports room search, booking, modification, and cancellation via natural language."},
        {"title": "Amenity Reservations", "description": "Books dining, spa, pool, gym, and activities from the same interface."},
        {"title": "Multi-Turn Memory", "description": "Maintains conversation context and tracks guest preferences."},
        {"title": "AI Avatar Interface", "description": "Provides an immersive, personalised avatar-driven guest experience."},
        {"title": "PMS & Channel Integration", "description": "Connects to PMS/channel managers for real-time availability and pricing."}
      ],
      "impact": "Reduces front-desk workload and call centre volumes while increasing ancillary revenue by enabling 24/7 self-service booking through an engaging, human-like AI interface.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#F59E0B",
      "icon": "🏨"
    },
    {
      "id": "06",
      "name": "VetAI",
      "domain": "HR Tech",
      "category": "AI Screening",
      "tagline": "End-to-End AI-Powered Candidate Screening and Interview Platform",
      "overview": "VetAI is a full-stack AI recruitment platform that automates the candidate screening process from job posting through to hiring decision. It conducts live AI-powered video interviews, evaluates responses across communication quality, technical accuracy, and behavioural competencies in real time, and delivers structured hiring reports to recruiters. Video proctoring and silence detection ensure interview integrity, while job-title-based configuration templates allow rapid setup for diverse roles across any industry.",
      "features": [
        {"title": "AI Video Interviews", "description": "Conducts live interviews with dynamic, contextual follow-up questions."},
        {"title": "Integrity Controls", "description": "Video proctoring, anomaly detection, and integrity scoring."},
        {"title": "Multi-Dimensional Evaluation", "description": "Scores technical, behavioural, and communication competencies."},
        {"title": "Structured Hiring Reports", "description": "Generates role-fit scores, rankings, and recruiter-ready summaries."},
        {"title": "Pipeline Management", "description": "Provides a dashboard for candidate tracking and recruiting analytics."}
      ],
      "impact": "Cuts time-to-screen by up to 70%, enabling recruiters to focus on final-stage evaluation while AI handles initial assessment at scale — consistently and without bias.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}, {"label": "Time-to-screen reduction", "value": "Up to 70%"}],
      "accent": "#8B5CF6",
      "icon": "🧑"
    },
    {
      "id": "07",
      "name": "Resonance",
      "domain": "L&D",
      "category": "Communication AI",
      "tagline": "AI Communication Coaching for Professional Upskilling",
      "overview": "Resonance is an AI-powered communication coaching platform designed to help professionals develop confident, effective communication skills. Using speech analysis and LLM-based evaluation, it assesses tone, clarity, pacing, filler word usage, and persuasive impact during practice sessions and real conversations. Personalised feedback loops and structured exercises help users build lasting communication habits, with progress tracked and visualised across sessions over time.",
      "features": [
        {"title": "Real-Time Speech Analysis", "description": "Analyses tone, pace, clarity, and filler word frequency."},
        {"title": "Actionable Feedback Reports", "description": "Generates specific improvement suggestions and coaching notes."},
        {"title": "Scenario Practice", "description": "Supports presentations, negotiations, and interview practice sessions."},
        {"title": "Progress Tracking", "description": "Visualises skill growth across sessions and time."},
        {"title": "Team Analytics", "description": "Provides cohort-level insights for L&D managers."}
      ],
      "impact": "Helps organisations build a culture of clear, confident communication — reducing meeting inefficiencies and improving stakeholder engagement across teams at every level.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#06B6D4",
      "icon": "🗣"
    },
    {
      "id": "08",
      "name": "Orbit",
      "domain": "Internal Tool",
      "category": "Project Management",
      "tagline": "Systech's Internal Project and Ticket Tracking System",
      "overview": "Orbit is Systech's purpose-built internal project management and ticket tracking platform. It centralises task creation, assignment, status tracking, and milestone management across all active projects and client engagements. Built to replace scattered spreadsheets and email threads, Orbit gives team leads and senior leadership a real-time, always-accurate operational view of delivery health, resource allocation, and sprint progress across the entire organisation.",
      "features": [
        {"title": "Projects & Sprints", "description": "Creates projects/sprints with milestones, deadlines, and priorities."},
        {"title": "Ticket Lifecycle", "description": "Manages tickets from backlog through completion sign-off."},
        {"title": "Workload Visibility", "description": "Assigns work and visualises individual workload across projects."},
        {"title": "Search & Filtering", "description": "Global search with priority/status filters across active projects."},
        {"title": "Collaboration Threads", "description": "Per-ticket activity feeds and async comment threads."}
      ],
      "impact": "Gives Systech leadership a single source of truth for all project delivery, replacing fragmented communication channels with a structured, searchable operational hub.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#64748B",
      "icon": "🗂"
    },
    {
      "id": "09",
      "name": "SysRank",
      "domain": "Internal Tool",
      "category": "Technical Assessment",
      "tagline": "Systech's Proprietary Technical Assessment and Benchmarking Platform",
      "overview": "SysRank is Systech's internal equivalent of HackerRank — a proprietary platform for evaluating technical talent through structured coding challenges, SQL assessments, data engineering problems, and timed problem sets. It is used both in the hiring pipeline to screen candidates objectively and internally to benchmark developer and analyst skill levels across the team. Question banks are managed by domain, difficulty level, and technology track.",
      "features": [
        {"title": "Timed Assessments", "description": "Runs coding and SQL challenges across data engineering tracks."},
        {"title": "Auto-Grading", "description": "Auto-graded test cases with partial scoring and breakdowns."},
        {"title": "Candidate Portal", "description": "Candidate-facing portal with proctoring and integrity controls."},
        {"title": "Question Bank Management", "description": "Organises questions by domain, difficulty, and technology track."},
        {"title": "Results & Benchmarking", "description": "Leaderboards, score history, and skill heatmaps for benchmarking."}
      ],
      "impact": "Standardises talent evaluation at Systech, ensuring consistent, objective assessment for both external hiring and internal capability development and benchmarking.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#64748B",
      "icon": "🧪"
    },
    {
      "id": "10",
      "name": "AeroIntel",
      "domain": "Aviation",
      "category": "RAG · Conversational AI",
      "tagline": "RAG-Powered AI Assistant for Airport CCR Technicians",
      "overview": "AeroIntel is a retrieval-augmented generation (RAG) AI assistant built specifically for airport Common Communication Room (CCR) technicians. It enables natural language queries over a hybrid knowledge base combining structured operational databases and unstructured technical documents. Technicians can instantly surface maintenance procedures, equipment manuals, fault histories, and operational protocols — without navigating complex document repositories or waiting for expert availability during live incidents.",
      "features": [
        {"title": "Hybrid RAG Search", "description": "Combines structured database querying with document retrieval."},
        {"title": "Conversational Interface", "description": "Context-aware, multi-turn natural language queries."},
        {"title": "Guided Procedures", "description": "Retrieves step-by-step maintenance procedures for resolution support."},
        {"title": "Fault History Matching", "description": "Finds similar incidents and patterns from historical faults."},
        {"title": "Agentic Orchestration", "description": "Synthesises answers across multiple sources for complex queries."}
      ],
      "impact": "Dramatically reduces time-to-resolution for technical faults by giving CCR engineers instant access to the right information during live airport incidents — when every minute counts.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#0EA5E9",
      "icon": "✈"
    },
    {
      "id": "11",
      "name": "SysMart",
      "domain": "Retail",
      "category": "AI Chatbot",
      "tagline": "Databricks-Powered AI Chatbot for Retail Operations",
      "overview": "SysMart is a conversational AI platform for retail operations, built on Databricks AI serving endpoints. It combines live SQL-based database querying with retrieval-augmented generation over product catalogues, return policies, and FAQ knowledge bases to answer questions from both customers and internal operations teams. The Freshchat integration deploys the chatbot directly into the existing customer support channel, enabling seamless handoff between AI and human agents when escalation is needed.",
      "features": [
        {"title": "Live DB Q&A", "description": "Answers questions via natural language queries against inventory and order databases."},
        {"title": "RAG Knowledge Layer", "description": "Retrieves and synthesises from catalogues, policies, and FAQs."},
        {"title": "Freshchat Integration", "description": "Deploys directly into Freshchat for customer support workflows."},
        {"title": "Smart Escalation", "description": "Escalates complex cases with human handoff when needed."},
        {"title": "Analytics Dashboard", "description": "Tracks usage, categorises queries, and supports continuous improvement."}
      ],
      "impact": "Reduces support ticket volume and handles routine product, order, and policy queries automatically — freeing support agents to focus on higher-value, complex customer interactions.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#F97316",
      "icon": "🛒"
    },
    {
      "id": "12",
      "name": "Intelliframe",
      "domain": "Design",
      "category": "Generative AI",
      "tagline": "AI Wireframe Generator for Dashboards and Data Products",
      "overview": "Intelliframe accelerates the design-to-development pipeline by generating annotated dashboard and application wireframes from natural language briefs or data schemas. Product managers and business analysts describe what they need — the metrics to display, the intended user journey, layout preferences — and Intelliframe produces structured wireframe mockups complete with component annotations, layout logic, and data binding suggestions ready for developer handoff.",
      "features": [
        {"title": "NL-to-Wireframe Generation", "description": "Generates wireframes for dashboards, apps, and portals from briefs."},
        {"title": "Schema-Aware Layout", "description": "Suggests layouts based on underlying data models and schemas."},
        {"title": "Annotated Components", "description": "Adds UX rationale and interaction behaviour notes per component."},
        {"title": "Export Options", "description": "Exports to image, PDF, and developer-ready specification formats."},
        {"title": "Iterative Refinement", "description": "Refines outputs via follow-up conversational prompts."}
      ],
      "impact": "Compresses weeks of wireframing and design iteration into hours, aligning product, design, and engineering teams around a shared visual brief much faster — reducing rework and misalignment.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#EC4899",
      "icon": "🧩"
    },
    {
      "id": "13",
      "name": "CCTV (AiCCTV)",
      "domain": "Security",
      "category": "Vision Analytics",
      "tagline": "AI-Powered Surveillance and Vision Analytics Platform",
      "overview": "AiCCTV is an intelligent video surveillance solution that transforms passive camera infrastructure into an active security intelligence layer. It processes live and recorded feeds to detect anomalies, unauthorised access, crowd density violations, safety hazards, and behavioural patterns of concern. Alerts are generated in real time with accompanying video clip evidence and incident metadata, enabling security operations centres to act proactively rather than reactively on emerging threats.",
      "features": [
        {"title": "Multi-Feed Detection", "description": "Real-time anomaly and intrusion detection across multiple camera feeds."},
        {"title": "Crowd & Zone Monitoring", "description": "Crowd density monitoring and restricted-zone access alerts."},
        {"title": "Safety Hazard Detection", "description": "Detects falls, unattended objects, and fire indicators."},
        {"title": "Evidence-Backed Alerts", "description": "Dispatches alerts with video clips and incident metadata."},
        {"title": "Incident Analytics", "description": "Historical analytics, heatmaps, and compliance reporting."}
      ],
      "impact": "Enables proactive security operations at scale across large facilities — reducing incident response time and providing evidence-backed records for compliance investigations and regulatory audits.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#DC2626",
      "icon": "📷"
    },
    {
      "id": "14",
      "name": "DataOne",
      "domain": "Data Engineering",
      "category": "AI Tooling",
      "tagline": "Unified AI Data Engineering Toolkit Across Fabric, Snowflake and Databricks",
      "overview": "DataOne is Systech's unified AI data engineering platform that bridges Microsoft Fabric, Snowflake, and Databricks under a single intelligent interface. Through purpose-built MCP (Model Context Protocol) servers, DataOne exposes the full capability of each platform to AI agents and automation workflows — enabling natural language pipeline creation, cross-platform schema exploration, query execution, and data asset management without switching between separate consoles or writing boilerplate integration code.",
      "features": [
        {"title": "MCP Server Layer", "description": "MCP servers purpose-built for Fabric, Snowflake, and Databricks."},
        {"title": "NL Pipeline Orchestration", "description": "Creates and orchestrates pipelines via natural language."},
        {"title": "Cross-Platform Discovery", "description": "Explores schemas and queries data lineage across platforms."},
        {"title": "AI Agent Automations", "description": "Automates data quality checks, monitoring, and alerting."},
        {"title": "No-Code Integrations", "description": "Connects to no-code workflow automation for orchestration."}
      ],
      "impact": "Reduces data engineering toil by giving AI agents and analysts a unified, language-driven interface across three major data platforms simultaneously — eliminating context-switching and manual API wrangling.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#22C55E",
      "icon": "🧱"
    },
    {
      "id": "15",
      "name": "Chef",
      "domain": "AI Avatar",
      "category": "Food and Beverage",
      "tagline": "AI-Powered Video Avatar for Food and Beverage Experiences",
      "overview": "Chef is an AI video avatar application tailored for the food and beverage industry. A lifelike AI presenter delivers personalised menu recommendations, ingredient explanations, allergen guidance, and step-by-step cooking walkthroughs — all through natural language conversation. The platform blends engaging content with commerce, enabling restaurants, food brands, and recipe platforms to offer an interactive, on-demand AI host available around the clock with the warmth of a human presenter.",
      "features": [
        {"title": "Conversational Video Avatar", "description": "Lifelike presenter with context-aware conversational interaction."},
        {"title": "Personalised Recommendations", "description": "Suggests dishes based on preferences and dietary restrictions."},
        {"title": "Guided Cooking Walkthroughs", "description": "Step-by-step cooking guidance with real-time Q&A."},
        {"title": "Allergen & Nutrition Guidance", "description": "Provides allergens, nutrition breakdowns, and substitutions."},
        {"title": "Commerce Integrations", "description": "Integrates with POS/e-commerce for in-conversation ordering."}
      ],
      "impact": "Creates a differentiated, always-on digital engagement channel for F&B brands — combining the warmth and personalisation of a human host with the scalability and consistency of AI.",
      "metrics": [{"label": "Coverage", "value": "Full"}, {"label": "Status", "value": "Live"}, {"label": "Version", "value": "v1.0"}],
      "accent": "#A16207",
      "icon": "👨"
    }
  ]
}


# -----------------------------
# Design System Constants
# -----------------------------
BASE_BG = "#05080F"        # Ultra Dark base
BASE_60 = "#080C14"        # Deep Midnight
SURFACE_30 = "#0F172A"     # Slate
ACCENT_10 = "#2563EB"      # Electric Blue (global)
TXT_PRIMARY = "#E5E7EB"
TXT_MUTED = "#94A3B8"
LINE_SUBTLE = "#1F2937"

PHI = 1.618  # Golden Ratio

# Slide size (16:9 default): 13.333 x 7.5 in
SLIDE_W = 13.333
SLIDE_H = 7.5

# Layout ratios
LEFT_RATIO = 0.62
RIGHT_RATIO = 0.38

# Margins / grid (8-pt grid translated to inches: 8pt = 0.111in)
M = 0.72  # generous margin (approx 52pt), aligned with 8pt grid spirit
G = 0.24  # gutter (approx 17pt)


def _safe_get(d, path, default=""):
    cur = d
    for k in path:
        if isinstance(cur, dict) and k in cur:
            cur = cur[k]
        else:
            return default
    return cur


def _chunk_text(s, max_chars=220):
    # Miller's Law chunking: short paragraph blocks.
    t = clean_txt(s)
    if not t:
        return []
    # Prefer sentence split; fallback to char chunk.
    parts = re.split(r"(?<=[.!?])\s+", t)
    chunks, buf = [], ""
    for p in parts:
        if not p:
            continue
        if len(buf) + len(p) + 1 <= max_chars:
            buf = (buf + " " + p).strip()
        else:
            if buf:
                chunks.append(buf)
            buf = p.strip()
    if buf:
        chunks.append(buf)
    return chunks[:3]  # keep tight


def _accent_normalize(hex_color):
    # Enforce "NO PURPLE": if accent missing or looks purple, use global blue.
    h = clean_txt(hex_color).upper()
    if not h.startswith("#") or len(h) != 7:
        return ACCENT_10
    # crude purple detection: high R and B, lower G
    r = int(h[1:3], 16)
    g = int(h[3:5], 16)
    b = int(h[5:7], 16)
    if (r > 120 and b > 120 and g < 120):  # likely purple/magenta
        return ACCENT_10
    return h


def _add_footer(slide, left_text, right_text):
    add_txt(slide, clean_txt(left_text), M, SLIDE_H - 0.45, 7.5, 0.3, 10, TXT_MUTED, False, 0)
    add_txt(slide, clean_txt(right_text), SLIDE_W - M - 3.5, SLIDE_H - 0.45, 3.5, 0.3, 10, TXT_MUTED, False, PP_ALIGN.RIGHT)


def build_hero(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BASE_BG)

    # Regions
    content_w = SLIDE_W - 2 * M
    left_w = content_w * LEFT_RATIO - G / 2
    right_w = content_w * RIGHT_RATIO - G / 2
    left_x = M
    right_x = M + left_w + G
    top_y = M
    body_h = SLIDE_H - 2 * M - 0.55  # reserve footer

    # Background panels (subtle)
    add_rect(slide, left_x, top_y, left_w, body_h, BASE_60, LINE_SUBTLE, 1.0, r=0.12)
    add_rect(slide, right_x, top_y, right_w, body_h, SURFACE_30, LINE_SUBTLE, 1.0, r=0.12)

    title = _safe_get(data, ["portfolio", "title"], "Application Portfolio")
    org = _safe_get(data, ["portfolio", "organisation"], "")
    year = _safe_get(data, ["portfolio", "year"], "")
    confidential = _safe_get(data, ["portfolio", "confidential"], False)

    # Typography scale
    body_pt = 18
    h1_pt = max(44, int(body_pt * PHI * 1.5))  # 44+
    h2_pt = int(body_pt * PHI)

    # Left content
    add_txt(slide, clean_txt(title), left_x + 0.6, top_y + 0.8, left_w - 1.2, 1.6, h1_pt, TXT_PRIMARY, True, 0)
    subtitle = " ".join([t for t in [org, year] if clean_txt(t)])
    if confidential:
        subtitle = (subtitle + " · " if subtitle else "") + "CONFIDENTIAL"
    add_txt(slide, clean_txt(subtitle), left_x + 0.6, top_y + 2.55, left_w - 1.2, 0.6, h2_pt, TXT_MUTED, True, 0)

    add_txt(
        slide,
        clean_txt("Luxury in layout. Clarity in systems. A concise catalog of Systech applications across domains."),
        left_x + 0.6, top_y + 3.35, left_w - 1.2, 1.2,
        body_pt, TXT_PRIMARY, False, 0
    )

    # Right: 4 stat cards
    stats = _safe_get(data, ["portfolio", "stats"], {}) or {}
    stat_items = [
        ("Total", _safe_get(stats, ["totalApps"], "—")),
        ("Domains", _safe_get(stats, ["domains"], "—")),
        ("AI / LLM", _safe_get(stats, ["aiPowered"], "—")),
        ("Tenure", _safe_get(stats, ["buildTenure"], "—")),
    ]
    card_gap = 0.28
    card_h = (body_h - 0.9 - card_gap * 3) / 4
    cy = top_y + 0.45
    for (label, value) in stat_items:
        add_rect(slide, right_x + 0.45, cy, right_w - 0.9, card_h, BASE_60, LINE_SUBTLE, 1.0, r=0.16)
        add_txt(slide, clean_txt(label).upper(), right_x + 0.75, cy + 0.25, right_w - 1.5, 0.35, 11, TXT_MUTED, True, 0)
        add_txt(slide, clean_txt(value), right_x + 0.75, cy + 0.65, right_w - 1.5, card_h - 0.8, 24, TXT_PRIMARY, True, 0)
        cy += card_h + card_gap

    _add_footer(slide, f"{clean_txt(org)} · {clean_txt(year)}", "HERO")


def build_catalog(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BASE_BG)

    apps = _safe_get(data, ["applications"], []) or []

    # Header
    add_txt(slide, "Catalog", M, M - 0.1, SLIDE_W - 2 * M, 0.6, 34, TXT_PRIMARY, True, 0)
    add_txt(slide, clean_txt("All applications (cards)."), M, M + 0.45, SLIDE_W - 2 * M, 0.4, 14, TXT_MUTED, False, 0)

    # Content area
    top_y = M + 1.05
    content_h = SLIDE_H - top_y - M - 0.55
    content_w = SLIDE_W - 2 * M

    # Asymmetric layout: left column wider than right, but grid of cards overall.
    # Use 3 columns with generous whitespace.
    cols = 3
    gap = 0.26
    card_w = (content_w - gap * (cols - 1)) / cols
    rows = max(1, int(content_h // 1.55))
    card_h = (content_h - gap * (rows - 1)) / rows
    capacity = rows * cols
    pages = max(1, ceil(len(apps) / capacity))

    # If more than capacity, we still show first page only (single catalog slide requirement)
    show_apps = apps[:capacity]

    # Cards
    i = 0
    for r in range(rows):
        for c in range(cols):
            if i >= len(show_apps):
                break
            app = show_apps[i] or {}
            aid = clean_txt(_safe_get(app, ["id"], ""))
            name = clean_txt(_safe_get(app, ["name"], "Untitled"))
            domain = clean_txt(_safe_get(app, ["domain"], ""))
            category = clean_txt(_safe_get(app, ["category"], ""))
            tagline = clean_txt(_safe_get(app, ["tagline"], ""))
            icon = clean_txt(_safe_get(app, ["icon"], ""))
            accent = _accent_normalize(_safe_get(app, ["accent"], ACCENT_10))

            x = M + c * (card_w + gap)
            y = top_y + r * (card_h + gap)

            add_rect(slide, x, y, card_w, card_h, SURFACE_30, LINE_SUBTLE, 1.0, r=0.14)

            # Top row: ID + icon + name
            add_txt(slide, f"{aid}".strip(), x + 0.35, y + 0.22, 0.7, 0.35, 12, TXT_MUTED, True, 0)
            add_txt(slide, icon, x + card_w - 0.75, y + 0.16, 0.4, 0.4, 16, TXT_PRIMARY, False, PP_ALIGN.RIGHT)

            add_txt(slide, name, x + 0.35, y + 0.55, card_w - 0.7, 0.55, 18, TXT_PRIMARY, True, 0)

            # Pills
            py = y + 1.1
            add_pill(slide, x + 0.35, py, domain if domain else "Domain", BASE_60, TXT_PRIMARY, size=9)
            add_pill(slide, x + 0.35 + 1.7, py, category if category else "Category", accent, "#FFFFFF", size=9)

            # Tagline (2 lines)
            add_txt(slide, tagline if tagline else "—", x + 0.35, y + 1.5, card_w - 0.7, card_h - 1.85, 12, TXT_MUTED, False, 0)

            i += 1

    # Page note if truncated
    if pages > 1:
        add_txt(
            slide,
            clean_txt(f"Showing {len(show_apps)} of {len(apps)} applications (single-slide catalog preview)."),
            M, SLIDE_H - M - 0.65, SLIDE_W - 2 * M, 0.35, 11, TXT_MUTED, False, 0
        )

    _add_footer(slide, clean_txt(_safe_get(data, ["portfolio", "title"], "")), "CATALOG")


def build_detail(prs, data, app):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BASE_BG)

    # Regions
    content_w = SLIDE_W - 2 * M
    left_w = content_w * LEFT_RATIO - G / 2
    right_w = content_w * RIGHT_RATIO - G / 2
    left_x = M
    right_x = M + left_w + G
    top_y = M
    body_h = SLIDE_H - 2 * M - 0.55

    add_rect(slide, left_x, top_y, left_w, body_h, BASE_60, LINE_SUBTLE, 1.0, r=0.12)
    add_rect(slide, right_x, top_y, right_w, body_h, SURFACE_30, LINE_SUBTLE, 1.0, r=0.12)

    # Data
    aid = clean_txt(_safe_get(app, ["id"], ""))
    name = clean_txt(_safe_get(app, ["name"], "Untitled Application"))
    domain = clean_txt(_safe_get(app, ["domain"], ""))
    category = clean_txt(_safe_get(app, ["category"], ""))
    tagline = clean_txt(_safe_get(app, ["tagline"], ""))
    overview = clean_txt(_safe_get(app, ["overview"], ""))
    impact = clean_txt(_safe_get(app, ["impact"], ""))
    icon = clean_txt(_safe_get(app, ["icon"], ""))
    accent = _accent_normalize(_safe_get(app, ["accent"], ACCENT_10))
    metrics = _safe_get(app, ["metrics"], []) or []
    features = _safe_get(app, ["features"], []) or []

    # Header (left)
    add_txt(slide, f"{aid} · {name}".strip(" ·"), left_x + 0.6, top_y + 0.55, left_w - 1.2, 0.55, 30, TXT_PRIMARY, True, 0)
    add_txt(slide, icon, left_x + left_w - 1.1, top_y + 0.52, 0.5, 0.5, 18, TXT_PRIMARY, False, PP_ALIGN.RIGHT)
    add_txt(slide, tagline, left_x + 0.6, top_y + 1.1, left_w - 1.2, 0.55, 14, TXT_MUTED, False, 0)

    # Pills row (left)
    py = top_y + 1.75
    add_pill(slide, left_x + 0.6, py, domain if domain else "Domain", SURFACE_30, TXT_PRIMARY, size=9)
    add_pill(slide, left_x + 0.6 + 1.8, py, category if category else "Category", accent, "#FFFFFF", size=9)

    # Overview (chunked)
    add_txt(slide, "Overview", left_x + 0.6, top_y + 2.25, left_w - 1.2, 0.35, 14, TXT_PRIMARY, True, 0)
    chunks = _chunk_text(overview, max_chars=230)
    oy = top_y + 2.65
    for ch in chunks:
        add_txt(slide, ch, left_x + 0.6, oy, left_w - 1.2, 0.62, 12, TXT_MUTED, False, 0)
        oy += 0.66

    # Features (left)
    fy = oy + 0.1
    add_txt(slide, "Key Features", left_x + 0.6, fy, left_w - 1.2, 0.35, 14, TXT_PRIMARY, True, 0)
    fy += 0.45

    # Feature list with bullet-like structure
    max_feat = 5
    feat_show = features[:max_feat]
    line_h = 0.55
    for f in feat_show:
        ft = clean_txt(_safe_get(f, ["title"], ""))
        fd = clean_txt(_safe_get(f, ["description"], ""))
        # Accent dot + title
        add_rect(slide, left_x + 0.6, fy + 0.12, 0.12, 0.12, accent, accent, 0.5, r=0.4)
        add_txt(slide, ft if ft else "—", left_x + 0.78, fy, left_w - 1.38, 0.3, 12, TXT_PRIMARY, True, 0)
        add_txt(slide, fd if fd else "", left_x + 0.78, fy + 0.24, left_w - 1.38, 0.35, 11, TXT_MUTED, False, 0)
        fy += line_h

    # Impact callout (left bottom)
    if impact:
        call_h = 1.0
        call_y = top_y + body_h - call_h - 0.55
        add_rect(slide, left_x + 0.6, call_y, left_w - 1.2, call_h, SURFACE_30, LINE_SUBTLE, 1.0, r=0.16)
        add_txt(slide, "Impact", left_x + 0.9, call_y + 0.18, left_w - 1.8, 0.3, 12, TXT_PRIMARY, True, 0)
        add_txt(slide, impact, left_x + 0.9, call_y + 0.44, left_w - 1.8, call_h - 0.52, 11, TXT_MUTED, False, 0)

    # Right panel: Technical Specs + Metrics
    add_txt(slide, "Technical Specs", right_x + 0.45, top_y + 0.55, right_w - 0.9, 0.35, 14, TXT_PRIMARY, True, 0)

    # Spec cards
    spec_gap = 0.22
    sy = top_y + 1.0
    spec_h = 0.9

    # Domain / Category / Status (best-effort)
    # Use metrics map for status/version/coverage if present.
    m_map = {}
    for m in metrics:
        k = clean_txt(_safe_get(m, ["label"], "")).lower()
        v = clean_txt(_safe_get(m, ["value"], ""))
        if k:
            m_map[k] = v

    specs = [
        ("Domain", domain if domain else "—"),
        ("Category", category if category else "—"),
        ("Status", m_map.get("status", "—")),
        ("Version", m_map.get("version", "—")),
    ]

    for (k, v) in specs:
        add_rect(slide, right_x + 0.45, sy, right_w - 0.9, spec_h, BASE_60, LINE_SUBTLE, 1.0, r=0.14)
        add_txt(slide, clean_txt(k).upper(), right_x + 0.75, sy + 0.18, right_w - 1.5, 0.3, 10, TXT_MUTED, True, 0)
        add_txt(slide, clean_txt(v), right_x + 0.75, sy + 0.42, right_w - 1.5, 0.4, 16, TXT_PRIMARY, True, 0)
        sy += spec_h + spec_gap

    # Metrics list card
    add_txt(slide, "Metrics", right_x + 0.45, sy + 0.1, right_w - 0.9, 0.35, 14, TXT_PRIMARY, True, 0)
    sy += 0.5

    met_h = top_y + body_h - sy - 0.45
    add_rect(slide, right_x + 0.45, sy, right_w - 0.9, max(1.4, met_h), BASE_60, LINE_SUBTLE, 1.0, r=0.14)

    # Render metrics in two columns inside the card
    inner_x = right_x + 0.75
    inner_y = sy + 0.25
    inner_w = right_w - 1.5
    col_gap = 0.3
    col_w = (inner_w - col_gap) / 2
    row_h = 0.42

    met_items = metrics if metrics else [{"label": "Coverage", "value": "—"}, {"label": "Status", "value": "—"}, {"label": "Version", "value": "—"}]
    for idx, m in enumerate(met_items[:10]):
        r = idx // 2
        c = idx % 2
        mx = inner_x + c * (col_w + col_gap)
        my = inner_y + r * row_h
        lab = clean_txt(_safe_get(m, ["label"], ""))
        val = clean_txt(_safe_get(m, ["value"], ""))
        add_txt(slide, lab, mx, my, col_w, 0.22, 10, TXT_MUTED, True, 0)
        add_txt(slide, val, mx, my + 0.18, col_w, 0.26, 12, TXT_PRIMARY, False, 0)

    _add_footer(slide, clean_txt(_safe_get(data, ["portfolio", "title"], "")), f"DETAIL · {aid}")


def main():
    parser = argparse.ArgumentParser(description="Generate Systech Application Portfolio Wireframes (PPTX).")
    parser.add_argument("--output", required=True, help="Output PPTX path")
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    data = DATA or {}
    apps = _safe_get(data, ["applications"], []) or []

    # Slides
    build_hero(prs, data)
    build_catalog(prs, data)
    for app in apps:
        build_detail(prs, data, app or {})

    prs.save(args.output)


if __name__ == "__main__":
    main()