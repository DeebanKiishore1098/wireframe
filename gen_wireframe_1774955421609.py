from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import argparse, re

# -----------------------------
# Required utilities (define all before using)
# -----------------------------
def clean_txt(t):
    return re.sub(r"<[^>]*>", "", str(t))[:200]

def _hex_to_rgb(hex_color):
    hx = str(hex_color).strip().lstrip("#")
    if len(hx) != 6:
        hx = "000000"
    return RGBColor(int(hx[0:2], 16), int(hx[2:4], 16), int(hx[4:6], 16))

def _fits(x, y, w, h, max_x=13.0, max_y=7.2):
    return (x + w) <= max_x and (y + h) <= max_y and x >= 0.0 and y >= 0.0 and w >= 0.0 and h >= 0.0

def _truncate(s, n):
    s = clean_txt(s)
    if len(s) <= n:
        return s
    return s[: max(0, n - 1)] + "…"

def _text_height_in(n_lines, size_pt):
    return n_lines * (size_pt / 72.0) * 1.4

def set_bg(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)

def add_rect(slide, x, y, w, h, fill_hex, line_hex="#334155", line_w=1.0, radius=0.08):
    if not _fits(x, y, w, h):
        return None
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _hex_to_rgb(fill_hex)
    shp.line.color.rgb = _hex_to_rgb(line_hex)
    shp.line.width = Pt(line_w)
    try:
        shp.adjustments[0] = float(radius)
    except Exception:
        pass
    return shp

def add_txt(slide, text, x, y, w, h, size_pt, color_hex, bold=False, align=1):
    if not _fits(x, y, w, h):
        return None
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = clean_txt(text)
    run.font.name = "Calibri"
    run.font.size = Pt(size_pt)
    run.font.bold = bool(bold)
    run.font.color.rgb = _hex_to_rgb(color_hex)
    if align == 2:
        p.alignment = PP_ALIGN.CENTER
    elif align == 3:
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT
    tf.word_wrap = True
    return tb

def add_button(slide, text, x, y, w, h, bg_hex="#3B82F6", fg_hex="#FFFFFF"):
    if not _fits(x, y, w, h):
        return None
    add_rect(slide, x, y, w, h, bg_hex, line_hex=bg_hex, line_w=1.0, radius=0.08)
    add_txt(slide, _truncate(text, 30), x, y + (h - _text_height_in(1, 12)) / 2.0, w, _text_height_in(1, 12), 12, fg_hex, True, 2)
    return True

def add_input_field(slide, label, x, y, w):
    label = _truncate(label, 30)
    label_h = _text_height_in(1, 11)
    gap = 0.06
    field_h = 0.42
    total_h = label_h + gap + field_h
    if not _fits(x, y, w, total_h):
        return 0.0
    add_txt(slide, label, x, y, w, label_h, 11, "#94A3B8", False, 1)
    add_rect(slide, x, y + label_h + gap, w, field_h, "#0F172A", line_hex="#475569", line_w=1.0, radius=0.08)
    add_txt(slide, "Enter " + label.lower(), x + 0.12, y + label_h + gap + 0.10, w - 0.24, _text_height_in(1, 11), 11, "#64748B", False, 1)
    return total_h

def add_table_placeholder(slide, headers, x, y, w, h):
    if not _fits(x, y, w, h):
        return None
    add_rect(slide, x, y, w, h, "#1E293B", line_hex="#334155", line_w=1.0, radius=0.08)

    header_h = 0.45
    row_h = (h - header_h - 0.35) / 4.0  # reserve pagination space
    if row_h < 0.35:
        row_h = 0.35

    # header row
    add_rect(slide, x + 0.02, y + 0.02, w - 0.04, header_h, "#0F172A", line_hex="#334155", line_w=1.0, radius=0.08)

    cols = len(headers)
    if cols <= 0:
        cols = 1
        headers = ["Column"]
    col_w = (w - 0.08) / cols
    for i, hdr in enumerate(headers):
        tx = x + 0.04 + i * col_w
        add_txt(slide, _truncate(hdr, 20), tx, y + 0.12, col_w - 0.04, _text_height_in(1, 11), 11, "#F1F5F9", True, 1)

    # rows (4)
    base_y = y + 0.02 + header_h + 0.06
    for r in range(4):
        ry = base_y + r * (row_h + 0.06)
        if not _fits(x + 0.02, ry, w - 0.04, row_h):
            continue
        fill = "#111B2E" if r % 2 == 0 else "#0F172A"
        add_rect(slide, x + 0.02, ry, w - 0.04, row_h, fill, line_hex="#334155", line_w=1.0, radius=0.08)
        for c in range(cols):
            tx = x + 0.04 + c * col_w
            sample = f"Sample {r+1}-{c+1}"
            add_txt(slide, _truncate(sample, 30), tx, ry + 0.10, col_w - 0.04, _text_height_in(1, 11), 11, "#94A3B8", False, 1)

    # pagination
    pag_y = y + h - 0.28
    if _fits(x + 0.02, pag_y, w - 0.04, 0.24):
        add_txt(slide, "Showing 1–4 of 128", x + 0.06, pag_y, 3.0, _text_height_in(1, 10), 10, "#94A3B8", False, 1)
        add_rect(slide, x + w - 2.0, pag_y - 0.02, 0.5, 0.28, "#0F172A", line_hex="#475569", line_w=1.0, radius=0.08)
        add_txt(slide, "Prev", x + w - 2.0, pag_y + 0.04, 0.5, _text_height_in(1, 10), 10, "#94A3B8", False, 2)
        add_rect(slide, x + w - 1.4, pag_y - 0.02, 0.5, 0.28, "#3B82F6", line_hex="#3B82F6", line_w=1.0, radius=0.08)
        add_txt(slide, "1", x + w - 1.4, pag_y + 0.04, 0.5, _text_height_in(1, 10), 10, "#FFFFFF", True, 2)
        add_rect(slide, x + w - 0.8, pag_y - 0.02, 0.5, 0.28, "#0F172A", line_hex="#475569", line_w=1.0, radius=0.08)
        add_txt(slide, "Next", x + w - 0.8, pag_y + 0.04, 0.5, _text_height_in(1, 10), 10, "#94A3B8", False, 2)
    return True

# -----------------------------
# Layout constants
# -----------------------------
NAV_H = 0.6
SIDEBAR_W = 2.0
CARD_GAP = 0.2
ROW_GAP = 0.15
SECTION_GAP = 0.3

BG = "#0F172A"
SURFACE = "#1E293B"
BORDER = "#334155"
ACCENT = "#3B82F6"
TXT = "#F1F5F9"
MUTED = "#94A3B8"
INPUT_BG = "#0F172A"
INPUT_BORDER = "#475569"

SAFE_X0, SAFE_Y0, SAFE_X1, SAFE_Y1 = 0.3, 0.3, 13.0, 7.2

# -----------------------------
# UI helpers
# -----------------------------
def add_top_nav(slide, app_name):
    x = SAFE_X0
    y = SAFE_Y0
    w = SAFE_X1 - SAFE_X0
    h = NAV_H
    add_rect(slide, x, y, w, h, SURFACE, line_hex=BORDER, line_w=1.0, radius=0.08)

    # left logo area
    add_txt(slide, _truncate(app_name, 20), x + 0.25, y + 0.18, 2.2, _text_height_in(1, 14), 14, TXT, True, 1)

    # center nav links
    links = ["Overview", "Reports", "Data", "Settings"]
    start_x = x + (w / 2.0) - 2.0
    for i, lnk in enumerate(links):
        lx = start_x + i * 1.05
        if i == 0:
            add_rect(slide, lx - 0.08, y + 0.14, 0.95, 0.32, "#0B2A5A", line_hex=ACCENT, line_w=1.0, radius=0.08)
            add_txt(slide, _truncate(lnk, 20), lx, y + 0.20, 0.8, _text_height_in(1, 11), 11, TXT, True, 2)
        else:
            add_txt(slide, _truncate(lnk, 20), lx, y + 0.20, 0.8, _text_height_in(1, 11), 11, MUTED, False, 2)

    # right icons + avatar
    bell_x = x + w - 1.85
    add_rect(slide, bell_x, y + 0.16, 0.34, 0.34, INPUT_BG, line_hex=INPUT_BORDER, line_w=1.0, radius=0.08)
    add_txt(slide, "🔔", bell_x, y + 0.19, 0.34, _text_height_in(1, 12), 12, MUTED, False, 2)

    avatar_x = x + w - 1.40
    add_rect(slide, avatar_x, y + 0.14, 0.42, 0.42, "#0B2A5A", line_hex=ACCENT, line_w=1.0, radius=0.08)
    add_txt(slide, "U", avatar_x, y + 0.20, 0.42, _text_height_in(1, 12), 12, "#FFFFFF", True, 2)

    user_x = x + w - 0.9
    add_txt(slide, "User", user_x, y + 0.20, 0.7, _text_height_in(1, 11), 11, MUTED, False, 1)

def add_sidebar(slide, active_label="Dashboard"):
    x = SAFE_X0
    y = SAFE_Y0 + NAV_H + ROW_GAP
    w = SIDEBAR_W
    h = SAFE_Y1 - y
    add_rect(slide, x, y, w, h, SURFACE, line_hex=BORDER, line_w=1.0, radius=0.08)

    items = [
        ("Dashboard", "🏠"),
        ("Data", "🗂"),
        ("Workflows", "⚙"),
        ("Alerts", "🔔"),
        ("Users", "👤"),
    ]
    y_cursor = y + 0.25
    for label, ico in items:
        item_h = 0.44
        if not _fits(x + 0.15, y_cursor, w - 0.3, item_h):
            break
        is_active = (label.lower() == str(active_label).lower())
        fill = "#0B2A5A" if is_active else INPUT_BG
        line = ACCENT if is_active else INPUT_BORDER
        add_rect(slide, x + 0.15, y_cursor, w - 0.3, item_h, fill, line_hex=line, line_w=1.0, radius=0.08)
        add_txt(slide, ico, x + 0.22, y_cursor + 0.10, 0.3, _text_height_in(1, 12), 12, TXT if is_active else MUTED, False, 2)
        add_txt(slide, _truncate(label, 20), x + 0.55, y_cursor + 0.12, w - 0.75, _text_height_in(1, 11), 11, TXT if is_active else MUTED, True if is_active else False, 1)
        y_cursor += item_h + ROW_GAP

def add_page_header(slide, title, breadcrumb="Home /"):
    content_x = SAFE_X0 + SIDEBAR_W + CARD_GAP
    content_y = SAFE_Y0 + NAV_H + ROW_GAP
    content_w = SAFE_X1 - content_x
    y_cursor = content_y

    # breadcrumb
    bc_h = _text_height_in(1, 10)
    if _fits(content_x, y_cursor, content_w * 0.7, bc_h):
        add_txt(slide, _truncate(breadcrumb, 60), content_x, y_cursor, content_w * 0.7, bc_h, 10, MUTED, False, 1)
    y_cursor += bc_h + 0.08

    # title + actions
    title_h = _text_height_in(1, 20)
    if _fits(content_x, y_cursor, content_w * 0.65, title_h):
        add_txt(slide, _truncate(title, 60), content_x, y_cursor, content_w * 0.65, title_h, 20, TXT, True, 1)

    btn_h = 0.34
    btn_y = y_cursor + 0.02
    bx = content_x + content_w - (0.95 + 0.85 + 0.9 + 0.2 + 0.2)
    if bx < content_x + 0.2:
        bx = content_x + 0.2

    add_button(slide, "+ New", bx, btn_y, 0.95, btn_h, ACCENT, "#FFFFFF")
    add_button(slide, "Export", bx + 0.95 + 0.2, btn_y, 0.85, btn_h, "#0F172A", MUTED)
    add_button(slide, "Filter", bx + 0.95 + 0.2 + 0.85 + 0.2, btn_y, 0.9, btn_h, "#0F172A", MUTED)

    y_cursor += title_h + SECTION_GAP
    return content_x, y_cursor, content_w

def add_kpi_row(slide, x, y, w, kpis):
    card_h = 0.9
    n = min(max(len(kpis), 3), 4)
    card_w = (w - (n - 1) * CARD_GAP) / n
    if card_w < 2.0:
        n = 3
        card_w = (w - (n - 1) * CARD_GAP) / n

    for i in range(n):
        cx = x + i * (card_w + CARD_GAP)
        if not _fits(cx, y, card_w, card_h):
            continue
        add_rect(slide, cx, y, card_w, card_h, SURFACE, line_hex=BORDER, line_w=1.0, radius=0.08)
        label, value, trend = kpis[i]
        add_txt(slide, _truncate(label, 30), cx + 0.2, y + 0.18, card_w - 0.4, _text_height_in(1, 11), 11, MUTED, False, 1)
        add_txt(slide, _truncate(value, 20), cx + 0.2, y + 0.42, card_w - 0.4, _text_height_in(1, 18), 18, TXT, True, 1)
        add_txt(slide, _truncate(trend, 20), cx + 0.2, y + 0.70, card_w - 0.4, _text_height_in(1, 10), 10, "#60A5FA", False, 1)
    return card_h

def add_chart_grid(slide, x, y, w):
    gap = CARD_GAP
    chart_h = 2.25
    col_w = (w - gap) / 2.0
    labels = ["Bar Chart — Revenue by Month", "Line Chart — Engagement Trend"]
    for i in range(2):
        cx = x + i * (col_w + gap)
        if not _fits(cx, y, col_w, chart_h):
            continue
        add_rect(slide, cx, y, col_w, chart_h, SURFACE, line_hex=BORDER, line_w=1.0, radius=0.08)
        add_txt(slide, labels[i], cx + 0.25, y + 0.20, col_w - 0.5, _text_height_in(2, 12), 12, TXT, True, 1)
        add_rect(slide, cx + 0.25, y + 0.70, col_w - 0.5, chart_h - 1.0, "#0F172A", line_hex="#475569", line_w=1.0, radius=0.08)
        add_txt(slide, "Chart placeholder", cx + 0.25, y + 1.25, col_w - 0.5, _text_height_in(1, 11), 11, MUTED, False, 2)
    return chart_h

def domain_items_for_app(app):
    name = clean_txt(app.get("name", "App"))
    domain = clean_txt(app.get("domain", "Items"))
    cat = clean_txt(app.get("category", ""))
    lname = name.lower()
    ldomain = domain.lower()
    lcat = cat.lower()

    if "casino" in ldomain or "maverick" in lname:
        return "Offers"
    if "bank" in ldomain or "trbank" in lname:
        return "Documents"
    if "esg" in ldomain or "sustain" in lname:
        return "Emission Records"
    if "ecovision" in lname or "vision" in lcat:
        return "Incidents"
    if "hospital" in ldomain or "concierge" in lname:
        return "Bookings"
    return "Records"

def table_headers_for_items(items):
    it = items.lower()
    if "offer" in it:
        return ["Offer ID", "Segment", "Reward", "Channel", "Expiry", "Status"]
    if "document" in it:
        return ["Doc ID", "Type", "Received", "Pages", "Confidence", "Status"]
    if "emission" in it:
        return ["Record ID", "Scope", "Source", "Period", "tCO₂e", "Status"]
    if "incident" in it:
        return ["Incident ID", "Site", "Severity", "Detected", "Owner", "Status"]
    if "booking" in it:
        return ["Booking ID", "Guest", "Room", "Check-in", "Nights", "Status"]
    return ["ID", "Name", "Category", "Owner", "Updated", "Status"]

def form_fields_for_items(items):
    it = items.lower()
    if "offer" in it:
        return ["Offer Name", "Player Segment", "Reward Type", "Reward Value", "Delivery Channel", "Expiry Date"]
    if "document" in it:
        return ["Document Type", "Source", "Batch Name", "Validation Ruleset", "Output Format", "Notify Email"]
    if "emission" in it:
        return ["Scope", "Activity Type", "Data Source", "Reporting Period", "Emission Factor Set", "Notes"]
    if "incident" in it:
        return ["Site", "Camera Feed", "Severity", "Assignee", "Corrective Action", "Due Date"]
    if "booking" in it:
        return ["Guest Name", "Room Type", "Check-in Date", "Check-out Date", "Amenity", "Special Requests"]
    return ["Name", "Type", "Owner", "Status", "Start Date", "Notes"]

# -----------------------------
# Slide builders
# -----------------------------
def build_title_slide(prs, app, org_name):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    name = clean_txt(app.get("name", "Application"))
    tagline = _truncate(app.get("tagline", "UI wireframes for the product."), 80)

    # centered card
    x = 1.2
    y = 1.6
    w = 11.0
    h = 3.2
    add_rect(slide, x, y, w, h, SURFACE, line_hex=BORDER, line_w=1.0, radius=0.08)

    add_txt(slide, name, x + 0.6, y + 0.55, w - 1.2, _text_height_in(1, 36), 36, TXT, True, 1)
    add_txt(slide, tagline, x + 0.6, y + 1.25, w - 1.2, _text_height_in(2, 18), 18, MUTED, False, 1)

    pill_w = 2.2
    pill_h = 0.42
    pill_x = x + 0.6
    pill_y = y + 2.15
    add_rect(slide, pill_x, pill_y, pill_w, pill_h, ACCENT, line_hex=ACCENT, line_w=1.0, radius=0.08)
    add_txt(slide, "UI Wireframes", pill_x, pill_y + 0.10, pill_w, _text_height_in(1, 12), 12, "#FFFFFF", True, 2)

    # bottom org
    add_txt(slide, _truncate(org_name, 60), SAFE_X0, SAFE_Y1 - 0.25, SAFE_X1 - SAFE_X0, _text_height_in(1, 11), 11, MUTED, False, 3)

def build_dashboard_slide(prs, app):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    app_name = clean_txt(app.get("name", "Application"))
    add_top_nav(slide, app_name)
    add_sidebar(slide, "Dashboard")

    content_x, y_cursor, content_w = add_page_header(slide, "Dashboard", "Home / Dashboard")

    # KPI row
    items = domain_items_for_app(app)
    kpis = [
        ("Active " + items, "128", "+6.2% vs last week"),
        ("Conversion Rate", "14.8%", "+1.1%"),
        ("Alerts (24h)", "9", "-2"),
        ("SLA Health", "98.3%", "+0.4%"),
    ]
    kpi_h = add_kpi_row(slide, content_x, y_cursor, content_w, kpis)
    y_cursor += kpi_h + SECTION_GAP

    # Charts
    ch_h = add_chart_grid(slide, content_x, y_cursor, content_w)
    y_cursor += ch_h + SECTION_GAP

    # small activity feed (optional; only if fits)
    feed_h = 1.05
    if _fits(content_x, y_cursor, content_w, feed_h):
        add_rect(slide, content_x, y_cursor, content_w, feed_h, SURFACE, line_hex=BORDER, line_w=1.0, radius=0.08)
        add_txt(slide, "Recent activity", content_x + 0.25, y_cursor + 0.18, content_w - 0.5, _text_height_in(1, 12), 12, TXT, True, 1)
        add_txt(slide, "• Scheduled job completed  • 3 new items created  • 1 policy updated",
                content_x + 0.25, y_cursor + 0.50, content_w - 0.5, _text_height_in(2, 11), 11, MUTED, False, 1)

def build_list_slide(prs, app):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    app_name = clean_txt(app.get("name", "Application"))
    add_top_nav(slide, app_name)
    add_sidebar(slide, "Data")

    items = domain_items_for_app(app)
    title = f"Manage {items}"
    content_x, y_cursor, content_w = add_page_header(slide, title, f"Home / Data / {items}")

    # filter bar
    bar_h = 0.92
    if _fits(content_x, y_cursor, content_w, bar_h):
        add_rect(slide, content_x, y_cursor, content_w, bar_h, SURFACE, line_hex=BORDER, line_w=1.0, radius=0.08)

        # two inputs + search button
        field_w = (content_w - 0.25 - 1.1 - 0.25 - 0.25) / 2.0
        if field_w < 2.6:
            field_w = 2.6
        fx1 = content_x + 0.25
        fy = y_cursor + 0.18
        used1 = add_input_field(slide, "Keyword", fx1, fy, field_w)
        fx2 = fx1 + field_w + 0.25
        used2 = add_input_field(slide, "Status", fx2, fy, field_w)

        btn_w = 1.1
        btn_x = content_x + content_w - 0.25 - btn_w
        btn_y = y_cursor + 0.40
        add_button(slide, "Search", btn_x, btn_y, btn_w, 0.36, ACCENT, "#FFFFFF")

    y_cursor += bar_h + SECTION_GAP

    # table
    table_h = 3.6
    if _fits(content_x, y_cursor, content_w, table_h):
        headers = table_headers_for_items(items)
        headers = [clean_txt(h) for h in headers]
        add_table_placeholder(slide, headers, content_x, y_cursor, content_w, table_h)
    y_cursor += table_h + SECTION_GAP

def build_form_slide(prs, app):
    # only for apps with 3+ features (per instruction)
    feats = app.get("features", [])
    if not isinstance(feats, list) or len(feats) < 3:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BG)

    app_name = clean_txt(app.get("name", "Application"))
    add_top_nav(slide, app_name)
    add_sidebar(slide, "Workflows")

    items = domain_items_for_app(app)
    content_x, y_cursor, content_w = add_page_header(slide, f"Create New {items[:-1] if items.endswith('s') else items}", f"Home / {items} / Create")

    # form container
    form_h = 4.35
    if not _fits(content_x, y_cursor, content_w, form_h):
        return
    add_rect(slide, content_x, y_cursor, content_w, form_h, SURFACE, line_hex=BORDER, line_w=1.0, radius=0.08)

    inner_x = content_x + 0.3
    inner_y = y_cursor + 0.25
    inner_w = content_w - 0.6

    col_gap = 0.35
    col_w = (inner_w - col_gap) / 2.0
    left_x = inner_x
    right_x = inner_x + col_w + col_gap

    fields = form_fields_for_items(items)
    fields = [clean_txt(f) for f in fields][:6]

    left_fields = fields[:3]
    right_fields = fields[3:6]

    yL = inner_y
    yR = inner_y
    for f in left_fields:
        used = add_input_field(slide, f, left_x, yL, col_w)
        if used <= 0:
            break
        yL += used + ROW_GAP

    for f in right_fields:
        used = add_input_field(slide, f, right_x, yR, col_w)
        if used <= 0:
            break
        yR += used + ROW_GAP

    # buttons bottom-right within form
    btn_h = 0.38
    btn_w1 = 0.95
    btn_w2 = 0.95
    btn_gap = 0.2
    btn_y = y_cursor + form_h - 0.55
    cancel_x = content_x + content_w - 0.3 - btn_w2
    save_x = cancel_x - btn_gap - btn_w1
    add_button(slide, "Save", save_x, btn_y, btn_w1, btn_h, ACCENT, "#FFFFFF")
    add_button(slide, "Cancel", cancel_x, btn_y, btn_w2, btn_h, "#0F172A", MUTED)

# -----------------------------
# Data (first 5 applications)
# -----------------------------
BRIEF = {
  "portfolio": {
    "title": "Systech Analytics · Application Portfolio",
    "organisation": "Systech Analytics",
    "year": "2026",
    "confidential": True,
    "stats": {
      "totalApps": "15",
      "domains": "AI, Data Engineering, Automation, Analytics",
      "aiPowered": "7+ AI / LLM-Powered",
      "buildTenure": "Not specified"
    }
  },
  "applications": [
    {
      "id": "01",
      "name": "Maverick",
      "domain": "Casino",
      "category": "AI Generation",
      "tagline": "AI-Powered Promotional Coupon Engine for Casino Players",
      "overview": "Maverick is a personalised incentive generation platform built for the gaming and casino industry. It analyses player behaviour, spending history, game preferences, and visit frequency to craft hyper-personalised promotional offers — from free spins and dining credits to VIP room upgrades. Powered by AI, Maverick replaces static promotions with dynamic, data-driven campaigns that adapt in real time to maximise engagement and lifetime value.",
      "features": [
        {"title": "Player segmentation","description": "Segments players by lifetime value, churn risk, and play patterns."},
        {"title": "Automated coupon generation","description": "Generates offers with configurable reward rules plus expiry and eligibility logic."},
        {"title": "Real-time offer delivery","description": "Delivers offers via SMS, app notifications, and kiosk displays."},
        {"title": "Campaign analytics & A/B testing","description": "Tracks performance and supports experimentation to optimise uplift."},
        {"title": "POS & loyalty integrations","description": "Integrates with casino POS and loyalty management systems."}
      ],
      "impact": "Increases promotional ROI and player retention by delivering the right offer to the right player at the right moment, replacing guesswork with precision targeting.",
      "metrics": [{"label": "Coverage","value": "Full"},{"label": "Status","value": "Live"},{"label": "Version","value": "v1.0"}],
      "accent": "#6B7280",
      "icon": "🧩"
    },
    {
      "id": "02",
      "name": "TRBank",
      "domain": "Banking",
      "category": "Data Processing",
      "tagline": "Unstructured-to-Structured Data Transformation for Banking",
      "overview": "TRBank is an intelligent document processing platform designed for financial institutions. It ingests unstructured banking documents — account statements, loan applications, KYC forms, remittance slips, and free-text correspondence — and transforms them into clean, validated, structured datasets ready for downstream analytics, compliance reporting, and core banking system ingestion.",
      "features": [
        {"title": "Multi-format ingestion","description": "Processes PDFs, scanned images, and handwritten forms."},
        {"title": "Entity extraction","description": "Extracts key fields such as account numbers, dates, amounts, and counterparties."},
        {"title": "Validation rules engine","description": "Flags anomalies, duplicates, and missing fields to improve data quality."},
        {"title": "Structured outputs","description": "Exports JSON/CSV and supports direct database writes for downstream use."},
        {"title": "Auditability & confidence","description": "Provides field-level audit trail and confidence scoring."}
      ],
      "impact": "Eliminates manual data entry for banking operations teams, reducing processing time from hours to minutes per batch while improving accuracy and compliance readiness.",
      "metrics": [{"label": "Coverage","value": "Full"},{"label": "Status","value": "Live"},{"label": "Version","value": "v1.0"}],
      "accent": "#6B7280",
      "icon": "🧩"
    },
    {
      "id": "03",
      "name": "Sustainability",
      "domain": "ESG",
      "category": "Emission Analytics",
      "tagline": "Full-Cycle ESG Intelligence and Emission Calculation Platform",
      "overview": "The Sustainability platform is an end-to-end ESG intelligence tool for organisations managing their environmental footprint. It ingests utility bills, fuel logs, travel records, and supply chain data to calculate carbon emissions across Scope 1, 2, and 3 categories. It generates audit-ready ESG reports aligned with GHG Protocol and GRI standards, and provides predictive analytics to model emission reduction scenarios for net-zero planning.",
      "features": [
        {"title": "Automated bill parsing","description": "Parses utility bills and maps activity data to appropriate emission factors."},
        {"title": "Scope 1/2/3 carbon accounting","description": "Calculates emissions aligned with the GHG Protocol across all scopes."},
        {"title": "AI-generated ESG reporting","description": "Produces narrative ESG reports suitable for board and regulatory submission."},
        {"title": "Benchmarking dashboards","description": "Provides peer benchmarking plus target-gap analysis visualisations."},
        {"title": "Scenario modelling","description": "Models reduction roadmaps and net-zero pathways with predictive analytics."}
      ],
      "impact": "Reduces ESG reporting effort from weeks of manual spreadsheet work to automated, audit-ready output generated in minutes — enabling organisations to meet sustainability reporting obligations with confidence.",
      "metrics": [{"label": "Coverage","value": "Full"},{"label": "Status","value": "Live"},{"label": "Version","value": "v1.0"}],
      "accent": "#6B7280",
      "icon": "🧩"
    },
    {
      "id": "04",
      "name": "Ecovision",
      "domain": "ESG",
      "category": "Vision AI",
      "tagline": "Video-Powered AI Auditing for Sustainability Compliance",
      "overview": "Ecovision brings computer vision into the sustainability space, transforming passive CCTV infrastructure into an active compliance monitoring system. It analyses live and recorded facility video feeds to detect energy waste, improper waste disposal, safety non-compliance, and resource misuse. Each detected event is logged with an extracted clip, timestamp, severity rating, and an AI-generated corrective action recommendation.",
      "features": [
        {"title": "Real-time video compliance analysis","description": "Detects sustainability and compliance events across facility feeds."},
        {"title": "Automated incident logging","description": "Logs incidents with clip extraction, timestamps, and severity classification."},
        {"title": "Corrective action recommendations","description": "Generates AI-backed guidance for remediation per detected incident."},
        {"title": "Compliance scoring & trends","description": "Dashboards compliance scores with historical trend analytics."},
        {"title": "ESG platform integration","description": "Integrates with the Sustainability platform for unified ESG reporting."}
      ],
      "impact": "Transforms sustainability auditing from periodic manual walkthroughs into a continuous, automated, evidence-driven compliance function — providing an always-on audit trail for regulators and internal teams.",
      "metrics": [{"label": "Coverage","value": "Full"},{"label": "Status","value": "Live"},{"label": "Version","value": "v1.0"}],
      "accent": "#6B7280",
      "icon": "🧩"
    },
    {
      "id": "05",
      "name": "Hotel Concierge",
      "domain": "Hospitality",
      "category": "AI Agent",
      "tagline": "AI Avatar Booking Agent for Hotels and Guest Amenities",
      "overview": "Hotel Concierge is a conversational AI agent that handles end-to-end hotel booking and guest services through a lifelike AI avatar interface. Guests can search rooms, check availability, book, and reserve amenities (dining, spa, leisure). The agent maintains context across multi-turn interactions and orchestrates reservations through hotel PMS integrations for a seamless 24/7 self-service experience.",
      "features": [
        {"title": "End-to-end booking flows","description": "Supports room search, booking, modification, and cancellation in natural language."},
        {"title": "Amenity reservations","description": "Books dining, spa, pool, gym, and activity slots via conversation."},
        {"title": "Conversational memory","description": "Maintains multi-turn context with guest preference tracking."},
        {"title": "AI avatar interface","description": "Provides a lifelike avatar experience for immersive guest self-service."},
        {"title": "PMS/channel manager integration","description": "Connects to PMS for real-time availability and pricing."}
      ],
      "impact": "Reduces front-desk workload and call centre volumes while increasing ancillary revenue by enabling 24/7 self-service booking through an engaging, human-like AI interface.",
      "metrics": [{"label": "Coverage","value": "Full"},{"label": "Status","value": "Live"},{"label": "Version","value": "v1.0"}],
      "accent": "#6B7280",
      "icon": "🧩"
    }
  ]
}

# -----------------------------
# Main
# -----------------------------
def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--output", required=True)
    args = parser.parse_args()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    portfolio = BRIEF.get("portfolio", {}) if isinstance(BRIEF.get("portfolio", {}), dict) else {}
    org_name = clean_txt(portfolio.get("organisation", "Organisation"))

    apps = BRIEF.get("applications", [])
    if not isinstance(apps, list):
        apps = []
    apps = apps[:5]

    for app in apps:
        if not isinstance(app, dict):
            continue
        build_title_slide(prs, app, org_name)
        build_dashboard_slide(prs, app)
        build_list_slide(prs, app)
        build_form_slide(prs, app)

    prs.save(